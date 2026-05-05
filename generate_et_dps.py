import os
import shutil
import copy
from pptx import Presentation
from lxml import etree

os.chdir('/home/user/oxxo')
TEMPLATE = "Formato DP Oficial.pptx"
OUTPUT_BASE = "Estrategia y Transformacion/Descripciones de Puestos"
NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# ── helpers ────────────────────────────────────────────────────────────────

def get_shape(slide, name, top=None, left=None, tol=200000):
    candidates = [s for s in slide.shapes if s.name == name]
    if top is not None:
        candidates = [s for s in candidates if abs(s.top - top) < tol]
    if left is not None:
        candidates = [s for s in candidates if abs(s.left - left) < tol]
    return candidates[0] if candidates else None


def replace_content(tf, items):
    """Keep paragraph 0 (section title), replace everything else with items."""
    txBody = tf._txBody
    direct_paras = [c for c in txBody if c.tag == f'{{{NS}}}p']
    template_para = copy.deepcopy(direct_paras[1] if len(direct_paras) > 1 else direct_paras[0])
    for p in direct_paras[1:]:
        txBody.remove(p)
    for item in items:
        new_p = copy.deepcopy(template_para)
        runs = new_p.findall(f'{{{NS}}}r')
        if runs:
            for r in runs[1:]:
                new_p.remove(r)
            t = runs[0].find(f'{{{NS}}}t')
            if t is not None:
                t.text = item
        else:
            r_el = etree.SubElement(new_p, f'{{{NS}}}r')
            t_el = etree.SubElement(r_el, f'{{{NS}}}t')
            t_el.text = item
        txBody.append(new_p)


def set_cell_text(cell, text):
    """Replace all text in a table cell while preserving cell formatting."""
    tf = cell.text_frame
    txBody = tf._txBody
    paras = [c for c in txBody if c.tag == f'{{{NS}}}p']
    template_para = copy.deepcopy(paras[0])
    for p in paras[1:]:
        txBody.remove(p)
    # Modify first paragraph
    first_para = [c for c in txBody if c.tag == f'{{{NS}}}p'][0]
    runs = first_para.findall(f'{{{NS}}}r')
    if runs:
        for r in runs[1:]:
            first_para.remove(r)
        t = runs[0].find(f'{{{NS}}}t')
        if t is not None:
            t.text = text
    else:
        r_el = etree.SubElement(first_para, f'{{{NS}}}r')
        t_el = etree.SubElement(r_el, f'{{{NS}}}t')
        t_el.text = text


def modify_header(slide, nombre_puesto, jefe_directo, area):
    shape = get_shape(slide, 'Title 2')
    if not shape:
        return
    tf = shape.text_frame
    if len(tf.paragraphs) >= 2:
        p = tf.paragraphs[1]
        if len(p.runs) >= 1:
            p.runs[0].text = f"Nombre Puesto: {nombre_puesto}"
        if len(p.runs) >= 2:
            p.runs[1].text = f"Nombre Puesto Jefe Directo: {jefe_directo}"
        if len(p.runs) >= 3:
            p.runs[2].text = f"Área Funcional: {area}"


def modify_proposito(slide, text):
    shape = get_shape(slide, 'Content Placeholder 4', top=969028)
    if shape:
        replace_content(shape.text_frame, [text])


def modify_responsabilidades(slide, items):
    shape = get_shape(slide, 'Content Placeholder 4', top=1823416)
    if shape:
        replace_content(shape.text_frame, items)


def modify_decisiones(slide, items):
    shape = get_shape(slide, 'Content Placeholder 4', top=5272810, left=74103)
    if shape:
        replace_content(shape.text_frame, items)


def modify_retos(slide, items):
    shape = get_shape(slide, 'Content Placeholder 4', top=5272810, left=4619672)
    if shape:
        replace_content(shape.text_frame, items)


def modify_dimensionamiento(slide, negocios, alcance, tramo):
    shape = get_shape(slide, 'Content Placeholder 4', top=3705118, left=139564)
    if not shape:
        return
    paras = shape.text_frame.paragraphs
    updates = {2: f"Negocios: {negocios}", 4: f"Alcance Geográfico: {alcance}", 6: f"Tramo de supervisión: {tramo}"}
    for idx, text in updates.items():
        if idx < len(paras) and paras[idx].runs:
            paras[idx].runs[0].text = text


def modify_estructura_org(slide, reportes, pares):
    shape = get_shape(slide, 'Content Placeholder 4', top=3704137, left=4620903)
    if not shape:
        return
    paras = shape.text_frame.paragraphs
    if len(paras) > 1 and paras[1].runs:
        paras[1].runs[0].text = f"Reportes Directos: {reportes}"
    if len(paras) > 2 and paras[2].runs:
        paras[2].runs[0].text = f"Pares: {pares}"


def modify_relations_table(slide, relations):
    """relations: list of (tipo, puesto, objetivo, frecuencia)"""
    for shape in slide.shapes:
        if hasattr(shape, 'table') and shape.name == '5 Tabla':
            tbl = shape.table
            row_idx = 1
            prev_tipo = None
            for rel in relations:
                if row_idx >= len(tbl.rows):
                    break
                tipo, puesto, objetivo, frecuencia = rel
                set_cell_text(tbl.rows[row_idx].cells[0], tipo if tipo != prev_tipo else '')
                set_cell_text(tbl.rows[row_idx].cells[1], puesto)
                set_cell_text(tbl.rows[row_idx].cells[2], objetivo)
                set_cell_text(tbl.rows[row_idx].cells[3], frecuencia)
                prev_tipo = tipo
                row_idx += 1
            break


def generate_dp(role, folder_name, filename):
    out_dir = os.path.join(OUTPUT_BASE, folder_name)
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, filename)
    shutil.copy(TEMPLATE, out_path)
    prs = Presentation(out_path)

    s1, s2 = prs.slides[0], prs.slides[1]
    modify_header(s1, role['nombre_puesto'], role['jefe_directo'], role['area'])
    modify_proposito(s1, role['proposito'])
    modify_responsabilidades(s1, role['responsabilidades'])
    modify_decisiones(s1, role['decisiones'])
    modify_retos(s1, role['retos'])

    modify_header(s2, role['nombre_puesto'], role['jefe_directo'], role['area'])
    modify_dimensionamiento(s2, role.get('negocios', 'OXXO México'), role.get('alcance', 'Nacional'), role.get('tramo', ''))
    modify_estructura_org(s2, role.get('reportes', ''), role.get('pares', ''))
    if role.get('relaciones'):
        modify_relations_table(s2, role['relaciones'])

    prs.save(out_path)
    print(f"  ✓ {out_path}")


# ── role data ──────────────────────────────────────────────────────────────

AREA = "Estrategia y Transformación"
DIRECTOR = "Director Estrategia y Transformación"

roles = [

# ══════════════════════════════════════════════════════════════════════════
# L2 – GERENTE SR FULLFILMENT
# ══════════════════════════════════════════════════════════════════════════
{
  '_folder': 'Gerente Sr Fullfilment',
  '_file':   'ET - L2 Gerente Sr Fullfilment.pptx',
  'nombre_puesto': 'Gerente Sr Fullfilment',
  'jefe_directo':  DIRECTOR,
  'area':          AREA,
  'proposito': (
    "Dirigir y evolucionar el modelo integral de Fulfillment, asegurando la correcta orquestación "
    "entre las capacidades comerciales, operativas, digitales y organizacionales, para habilitar el "
    "crecimiento rentable del negocio, mejorar la experiencia end-to-end del cliente y garantizar la "
    "ejecución eficiente de los servicios en tienda y plataformas, alineado a la estrategia corporativa."
  ),
  'responsabilidades': [
    "Definir y liderar la estrategia integral de Fulfillment, incluyendo los servicios digitales, e-commerce y pagos, asegurando su alineación con los objetivos del negocio, el modelo comercial y la experiencia del cliente.",
    "Gobernar el portafolio de servicios de Fulfillment, priorizando iniciativas, roadmaps, productos digitales y MVPs que maximicen valor, margen, escalabilidad y adopción.",
    "Orquestar la ejecución end-to-end del modelo de Fulfillment, asegurando la integración efectiva entre áreas comerciales, marketing, operaciones, TI, producto y diseño organizacional.",
    "Impulsar el desarrollo y gestión de socios estratégicos (marketplaces, fintechs, aliados tecnológicos y operativos), habilitando nuevos modelos de negocio y capacidades del ecosistema Fulfillment.",
    "Asegurar la correcta operación y despliegue de los servicios en tienda, manteniendo estándares de calidad, conectividad al punto de venta y mejores prácticas operativas.",
    "Liderar el análisis de información y desempeño del negocio, definiendo y monitoreando KPIs clave de crecimiento, margen, conversión, volumen, eficiencia operativa y experiencia del cliente.",
    "Diseñar y evolucionar los procesos de experiencia digital y operativa, garantizando interacciones integradas de extremo a extremo y habilitando servicios de valor agregado.",
    "Desarrollar y coordinar a los equipos a su cargo, fortaleciendo las capacidades organizacionales necesarias para la evolución y sostenibilidad del modelo Fulfillment.",
  ],
  'decisiones': [
    "Determinar prioridades del portafolio balanceando valor al cliente, rentabilidad, complejidad operativa y alineación estratégica.",
    "Definir el modelo de integración entre comercial, marketing, operaciones, TI y producto para asegurar ejecución consistente.",
    "Decidir alianzas con marketplaces, fintechs y proveedores clave, evaluando impacto en capacidades, costos y escalabilidad.",
    "Establecer y ajustar los KPIs clave que guían decisiones de crecimiento, margen, eficiencia y experiencia del cliente.",
    "Definir el ritmo de expansión de los servicios sin comprometer calidad, estabilidad en tienda ni experiencia del cliente.",
  ],
  'retos': [
    "Lograr una integración ágil y consistente entre áreas comerciales, marketing, operaciones, TI y producto, evitando silos y fricciones en la ejecución sin dependencia jerárquica directa.",
    "Crecer y evolucionar el portafolio de servicios de Fulfillment manteniendo estándares de calidad, simplicidad operativa y estabilidad en tienda.",
    "Asegurar que la estrategia, los roadmaps y las decisiones del portafolio se traduzcan en resultados concretos a nivel operativo, comercial y de experiencia del cliente.",
    "Balancear crecimiento, costos y eficiencia en un modelo con múltiples servicios, procesos y socios, cuidando márgenes y sostenibilidad del negocio.",
    "Impulsar la correcta implementación de los servicios de Fulfillment en tienda, minimizando fricciones operativas y garantizando una experiencia consistente para el cliente.",
  ],
  'negocios': 'OXXO México', 'alcance': 'Nacional', 'tramo': '5 reportes directos',
  'reportes': 'Coordinador Product Owner Servicios (Paulina Salinas), Responsable MKT Adquisición Cliente (Cesar Garza), Coordinador Op Fullfillment (Antonio Mora), Responsable Op Dis Fullfillment (Alberto Valdez), Gerente Área Comercial Fullfillment (Kenya Vazquez)',
  'pares': 'Gerente OXXO Cel, Gerente Sr Comercial Servicios Financieros, Gerente Sr Comercial Comercio Electrónico',
  'relaciones': [
    ('INTERNAS', 'Director Estrategia y Transformación', 'Alineación estratégica y reporte de resultados del área', 'Mensual'),
    ('INTERNAS', 'Gerente OXXO Cel', 'Coordinación de estrategia de servicios digitales y cliente', 'Semanal'),
    ('INTERNAS', 'Gerente Sr Comercial Servicios Financieros', 'Integración de capacidades financieras y de pago en el modelo', 'Quincenal'),
    ('INTERNAS', 'Áreas Comerciales y Operaciones', 'Ejecución y despliegue de servicios de Fulfillment en tienda', 'Semanal'),
    ('INTERNAS', 'TI', 'Integración tecnológica y entrega de soluciones digitales', 'Semanal'),
  ],
},

# ══════════════════════════════════════════════════════════════════════════
# L3s – FULLFILMENT
# ══════════════════════════════════════════════════════════════════════════
{
  '_folder': 'Gerente Sr Fullfilment',
  '_file':   'ET - L3 Coordinador Product Owner Servicios.pptx',
  'nombre_puesto': 'Coordinador Product Owner Servicios',
  'jefe_directo':  'Gerente Sr Fullfilment',
  'area':          AREA,
  'proposito': (
    "Definir y ejecutar el roadmap de productos digitales y MVPs de los servicios de Fulfillment, "
    "coordinando con TI y áreas clave la entrega de soluciones y evolutivos del producto, para asegurar "
    "que los servicios generen valor, escalabilidad y simplicidad operativa alineados a los objetivos del negocio."
  ),
  'responsabilidades': [
    "Definir y priorizar el roadmap de productos digitales y MVPs de Fulfillment, alineado a los objetivos estratégicos del negocio y a las necesidades del cliente.",
    "Gestionar el backlog de producto, documentando requerimientos, criterios de aceptación y prioridades en coordinación con los stakeholders del negocio.",
    "Coordinar con TI y áreas funcionales clave la entrega efectiva de soluciones y evolutivos del producto, asegurando cumplimiento de tiempos y especificaciones.",
    "Dar seguimiento a KPIs de adopción, uso y experiencia del cliente, proponiendo mejoras al producto con base en datos de desempeño.",
    "Apoyar la evolución de los procesos de experiencia digital y operativa, garantizando interacciones integradas y de valor para el cliente.",
    "Asegurar que los servicios desarrollados generen valor medible, escalabilidad y simplicidad operativa para el modelo Fulfillment.",
  ],
  'decisiones': [
    "Definir la priorización del backlog y los criterios de aceptación de producto en cada ciclo de desarrollo.",
    "Determinar los requerimientos funcionales de las soluciones digitales a desarrollar, balanceando valor al cliente y viabilidad técnica.",
  ],
  'retos': [
    "Asegurar alineación continua entre los equipos de negocio, TI y operaciones para una entrega ágil y de calidad.",
    "Gestionar expectativas de múltiples stakeholders con prioridades diversas en un entorno dinámico y de constante evolución.",
  ],
  'negocios': 'OXXO México', 'alcance': 'Nacional', 'tramo': 'Sin reportes directos',
  'reportes': '',
  'pares': 'Responsable MKT Adquisición Cliente, Coordinador Op Fullfillment, Responsable Op Dis Fullfillment, Gerente Área Comercial Fullfillment',
  'relaciones': [
    ('INTERNAS', 'Gerente Sr Fullfilment', 'Alineación de roadmap y prioridades de producto', 'Semanal'),
    ('INTERNAS', 'TI', 'Coordinación de desarrollos y entrega de soluciones digitales', 'Semanal'),
    ('INTERNAS', 'Operaciones', 'Validación de factibilidad operativa de los servicios', 'Quincenal'),
  ],
},
{
  '_folder': 'Gerente Sr Fullfilment',
  '_file':   'ET - L3 Responsable MKT Adquisicion Cliente.pptx',
  'nombre_puesto': 'Responsable MKT Adquisición Cliente',
  'jefe_directo':  'Gerente Sr Fullfilment',
  'area':          AREA,
  'proposito': (
    "Diseñar y ejecutar la estrategia de adquisición y comunicación de los servicios de Fulfillment, "
    "impulsando la adopción del cliente mediante mensajes claros, lanzamientos coordinados y seguimiento "
    "de KPIs de conversión y funnel de cliente."
  ),
  'responsabilidades': [
    "Diseñar y ejecutar la estrategia de adquisición de clientes para los servicios de Fulfillment, definiendo canales, mensajes e iniciativas de comunicación que incrementen la adopción.",
    "Coordinar lanzamientos y estrategias go-to-market con las áreas comerciales y de producto, asegurando alineación de mensajes y timings.",
    "Dar seguimiento a KPIs de conversión y funnel de cliente, analizando resultados y proponiendo ajustes a la estrategia de adquisición.",
    "Impulsar la adopción del cliente en tienda y canales digitales, asegurando mensajes claros, materiales de comunicación efectivos y experiencia consistente.",
    "Contribuir a la mejora continua de la experiencia del cliente, identificando oportunidades de optimización en el proceso de adquisición y onboarding.",
    "Coordinar con el equipo de Fulfillment y áreas de marketing la creación de contenidos y materiales de comunicación para los servicios.",
  ],
  'decisiones': [
    "Definir los canales de comunicación y adquisición más efectivos para cada servicio de Fulfillment.",
    "Establecer los mensajes clave y posicionamiento de los servicios en cada campaña de adquisición.",
  ],
  'retos': [
    "Incrementar la adopción de servicios de Fulfillment en un entorno donde los clientes tienen múltiples opciones y la diferenciación debe ser clara.",
    "Coordinar campañas de adquisición consistentes tanto en tienda como en canales digitales.",
  ],
  'negocios': 'OXXO México', 'alcance': 'Nacional', 'tramo': 'Sin reportes directos',
  'reportes': '',
  'pares': 'Coordinador Product Owner Servicios, Coordinador Op Fullfillment, Responsable Op Dis Fullfillment, Gerente Área Comercial Fullfillment',
  'relaciones': [
    ('INTERNAS', 'Gerente Sr Fullfilment', 'Alineación de estrategia de adquisición y KPIs', 'Semanal'),
    ('INTERNAS', 'Marketing Corporativo', 'Coordinación de campañas y materiales de comunicación', 'Quincenal'),
    ('INTERNAS', 'Producto', 'Sincronización de roadmap con lanzamientos y comunicación', 'Semanal'),
  ],
},
{
  '_folder': 'Gerente Sr Fullfilment',
  '_file':   'ET - L3 Coordinador Op Fullfillment.pptx',
  'nombre_puesto': 'Coordinador Op Fullfillment',
  'jefe_directo':  'Gerente Sr Fullfilment',
  'area':          AREA,
  'proposito': (
    "Coordinar la operación end-to-end del modelo Fulfillment en tienda, diseñando e implementando "
    "metodologías operativas, estándares y drills que aseguren la correcta ejecución de los servicios, "
    "la adopción de cambios y la sostenibilidad del modelo operativo."
  ),
  'responsabilidades': [
    "Habilitar la correcta ejecución del modelo Fulfillment operativamente en tienda, coordinando con operaciones zonales y despliegue para garantizar estándares de servicio.",
    "Diseñar e implementar drills, metodologías operativas y estándares de ejecución para los servicios de Fulfillment, asegurando consistencia en todas las tiendas.",
    "Apoyar el desarrollo de capacidades organizacionales del modelo Fulfillment, identificando necesidades de capacitación y habilitación del personal.",
    "Facilitar la adopción de cambios y evolutivos en la operación, coordinando la comunicación y el acompañamiento al personal de tienda.",
    "Monitorear KPIs operativos de los servicios de Fulfillment en tienda, identificando desviaciones y proponiendo acciones correctivas.",
    "Coordinar con las áreas de Operaciones, Despliegue y TI para garantizar la correcta implementación de los servicios a nivel de punto de venta.",
  ],
  'decisiones': [
    "Definir los estándares operativos y metodologías de ejecución de los servicios de Fulfillment en tienda.",
    "Establecer prioridades en el plan de drills y habilitación operativa según el impacto y urgencia de cada iniciativa.",
  ],
  'retos': [
    "Garantizar consistencia operativa en miles de tiendas con contextos y capacidades organizacionales diferentes.",
    "Gestionar el cambio operativo en tienda de manera ágil, minimizando fricciones y asegurando adopción efectiva.",
  ],
  'negocios': 'OXXO México', 'alcance': 'Nacional', 'tramo': 'Sin reportes directos',
  'reportes': '',
  'pares': 'Coordinador Product Owner Servicios, Responsable MKT Adquisición Cliente, Responsable Op Dis Fullfillment, Gerente Área Comercial Fullfillment',
  'relaciones': [
    ('INTERNAS', 'Gerente Sr Fullfilment', 'Reporte de ejecución operativa y KPIs de tienda', 'Semanal'),
    ('INTERNAS', 'Operaciones y Despliegue', 'Coordinación de implementación y drills en tienda', 'Semanal'),
    ('INTERNAS', 'TI', 'Soporte tecnológico para los servicios operativos de Fulfillment', 'Quincenal'),
  ],
},
{
  '_folder': 'Gerente Sr Fullfilment',
  '_file':   'ET - L3 Responsable Op Dis Fullfillment.pptx',
  'nombre_puesto': 'Responsable Op Dis Fullfillment',
  'jefe_directo':  'Gerente Sr Fullfilment',
  'area':          AREA,
  'proposito': (
    "Gestionar la operación y distribución de los servicios de Fulfillment, asegurando la disponibilidad, "
    "correcta distribución de recursos y materiales, y el cumplimiento de los estándares operativos de "
    "los servicios en tienda y canales digitales."
  ),
  'responsabilidades': [
    "Gestionar la operación de distribución de los servicios de Fulfillment, asegurando disponibilidad de recursos, materiales e insumos para la correcta ejecución en tienda.",
    "Coordinar los procesos logísticos y operativos que habilitan la prestación de los servicios de Fulfillment, garantizando continuidad y eficiencia.",
    "Monitorear indicadores clave de disponibilidad y desempeño operativo de la distribución, proponiendo mejoras para optimizar la cadena de servicio.",
    "Gestionar la relación con proveedores y áreas internas para asegurar el abasto oportuno de los recursos necesarios para la operación de Fulfillment.",
    "Identificar y escalar oportunidades de mejora en los procesos de distribución y operación, contribuyendo a la eficiencia y calidad del modelo.",
    "Apoyar la implementación de nuevos servicios y evolutivos en la operación de distribución, coordinando la habilitación y puesta en marcha.",
  ],
  'decisiones': [
    "Definir la estrategia de distribución de recursos y materiales para garantizar disponibilidad en tienda de los servicios de Fulfillment.",
    "Establecer prioridades en la atención de incidencias operativas de distribución con base en impacto al servicio y al cliente.",
  ],
  'retos': [
    "Asegurar disponibilidad y abasto continuo en una red de miles de tiendas con alta variabilidad en demanda y condiciones operativas.",
    "Optimizar costos de distribución sin comprometer la calidad y disponibilidad del servicio para el cliente.",
  ],
  'negocios': 'OXXO México', 'alcance': 'Nacional', 'tramo': 'Sin reportes directos',
  'reportes': '',
  'pares': 'Coordinador Product Owner Servicios, Responsable MKT Adquisición Cliente, Coordinador Op Fullfillment, Gerente Área Comercial Fullfillment',
  'relaciones': [
    ('INTERNAS', 'Gerente Sr Fullfilment', 'Reporte de disponibilidad y desempeño de distribución', 'Semanal'),
    ('INTERNAS', 'Cadena de Suministro', 'Coordinación de abasto y logística', 'Semanal'),
    ('INTERNAS', 'Operaciones', 'Sincronización de operación de tienda con distribución', 'Quincenal'),
  ],
},
{
  '_folder': 'Gerente Sr Fullfilment',
  '_file':   'ET - L3 Gerente Area Comercial Fullfillment.pptx',
  'nombre_puesto': 'Gerente Área Comercial Fullfillment',
  'jefe_directo':  'Gerente Sr Fullfilment',
  'area':          AREA,
  'proposito': (
    "Impulsar la estrategia comercial de los servicios de Fulfillment, asegurando el crecimiento y "
    "rentabilidad del portafolio mediante el desarrollo de socios comerciales estratégicos, la correcta "
    "ejecución en tienda y el seguimiento a KPIs de volumen, margen y conversión."
  ),
  'responsabilidades': [
    "Impulsar la estrategia comercial de los servicios de Fulfillment para asegurar crecimiento rentable y cumplimiento de objetivos de volumen y margen.",
    "Gestionar y desarrollar socios comerciales estratégicos, incluyendo marketplaces y comercios digitales, habilitando nuevas capacidades y modelos de negocio.",
    "Asegurar la adopción de los servicios desde el frente comercial, coordinando la correcta bajada a tienda e impulsando la ejecución efectiva.",
    "Contribuir a la priorización del portafolio de servicios desde la perspectiva comercial y de mercado, aportando visión de clientes y socios.",
    "Dar seguimiento a KPIs comerciales de volumen, margen y conversión, analizando resultados y proponiendo iniciativas de mejora.",
    "Coordinar con operaciones y marketing el lanzamiento y activación comercial de nuevos servicios y evolutivos del portafolio Fulfillment.",
  ],
  'decisiones': [
    "Definir las prioridades y acuerdos comerciales con socios estratégicos del ecosistema Fulfillment.",
    "Establecer la estrategia de activación y ejecución comercial en tienda para los servicios de Fulfillment.",
  ],
  'retos': [
    "Asegurar adopción y ejecución comercial consistente de los servicios en la red de tiendas, con alta variabilidad local.",
    "Desarrollar y escalar relaciones con socios estratégicos en un entorno competitivo donde la propuesta de valor debe diferenciarse continuamente.",
  ],
  'negocios': 'OXXO México', 'alcance': 'Nacional', 'tramo': 'Sin reportes directos',
  'reportes': '',
  'pares': 'Coordinador Product Owner Servicios, Responsable MKT Adquisición Cliente, Coordinador Op Fullfillment, Responsable Op Dis Fullfillment',
  'relaciones': [
    ('INTERNAS', 'Gerente Sr Fullfilment', 'Alineación de estrategia comercial y seguimiento de KPIs', 'Semanal'),
    ('INTERNAS', 'Áreas Comerciales y Tienda', 'Ejecución y activación comercial de servicios de Fulfillment', 'Semanal'),
    ('EXTERNAS', 'Socios Estratégicos (marketplaces, fintechs)', 'Gestión y desarrollo de relaciones comerciales', 'Mensual'),
  ],
},

# ══════════════════════════════════════════════════════════════════════════
# L2 – GERENTE OXXO CEL
# ══════════════════════════════════════════════════════════════════════════
{
  '_folder': 'Gerente OXXO Cel',
  '_file':   'ET - L2 Gerente OXXO Cel.pptx',
  'nombre_puesto': 'Gerente OXXO Cel',
  'jefe_directo':  DIRECTOR,
  'area':          AREA,
  'proposito': (
    "Liderar el crecimiento rentable y sostenible del negocio OXXO Cel, definiendo y ejecutando la "
    "estrategia de conectividad del ecosistema OXXO, desde la adquisición y fidelización de clientes "
    "hasta la operación y evolución de productos y servicios, asegurando una experiencia consistente, "
    "una operación escalable y una propuesta de valor competitiva que capitalice el ecosistema OXXO."
  ),
  'responsabilidades': [
    "Definir y liderar la estrategia integral del negocio OXXO Cel, asegurando crecimiento rentable, propuesta de valor competitiva, posicionamiento en el mercado de telefonía y alineación con la estrategia del ecosistema OXXO.",
    "Impulsar el crecimiento de la base de clientes y los ingresos, mediante el diseño y ejecución de estrategias comerciales end-to-end que integren adquisición, pricing, promociones, incentivos, campañas, activaciones en tienda y expansión de dispositivos y servicios.",
    "Liderar la evolución del portafolio de productos y servicios OXXO Cel, incluyendo telefonía, eSIM, bundles, servicios digitales y nuevas propuestas, alineadas a las necesidades del cliente y a las tendencias del mercado.",
    "Asegurar una experiencia de cliente consistente y diferenciada, integrando estrategias de CX, comunicación, fidelización y retención que incrementen la permanencia, reduzcan churn y maximicen el valor del cliente.",
    "Garantizar la correcta operación integral del negocio, asegurando la continuidad del servicio de líneas, chips, dispositivos y usuarios activos, así como la ejecución eficiente de procesos operativos en tienda y canales digitales.",
    "Dirigir la operación y ejecución comercial en tienda, asegurando capacitación, habilitación, onboarding, herramientas e incentivos que maximicen el desempeño del personal y la productividad del punto de venta.",
    "Asegurar la operación tecnológica de OXXO Cel, incluyendo plataformas, soporte, procesos digitales, call center, medios de pago e infraestructura, impulsando la mejora continua y la eficiencia operativa.",
    "Gestionar estratégicamente la relación con proveedores y socios clave, garantizando condiciones comerciales óptimas, continuidad de la cadena de suministro y soporte al crecimiento del negocio.",
    "Impulsar un modelo de gestión y toma de decisiones basado en datos, asegurando el uso efectivo de analítica, KPIs y dashboards para optimizar desempeño comercial, operativo, financiero y de experiencia del cliente.",
    "Coordinar y alinear equipos multidisciplinarios y stakeholders internos, asegurando claridad de objetivos, ejecución integrada, cumplimiento de indicadores clave y representación efectiva del negocio en foros estratégicos.",
  ],
  'decisiones': [
    "Definir la estrategia de crecimiento y rentabilidad del negocio, priorizando entre adquisición, retención, nuevos productos y expansión dentro del ecosistema OXXO.",
    "Establecer la estrategia comercial integral, incluyendo pricing, promociones, campañas e incentivos, balanceando volumen, margen y valor del cliente.",
    "Decidir la evolución del portafolio de productos y servicios, determinando lanzamientos, ajustes o descontinuaciones conforme a desempeño y mercado.",
    "Asignar recursos entre operación, experiencia, tecnología y crecimiento, asegurando estabilidad operativa y escalabilidad del negocio.",
    "Seleccionar y gestionar alianzas estratégicas y proveedores clave, definiendo modelos de colaboración que aseguren continuidad y ventaja competitiva.",
  ],
  'retos': [
    "Lograr crecimiento acelerado sin deteriorar rentabilidad ni experiencia del cliente, en un mercado altamente competitivo y sensible a precio.",
    "Diferenciar la propuesta de valor de OXXO Cel, capitalizando el ecosistema OXXO frente a operadores tradicionales y ofertas digitales.",
    "Mantener estabilidad operativa y tecnológica mientras se escala la base de usuarios, productos y transacciones.",
    "Reducir churn y aumentar el valor del cliente, asegurando una experiencia consistente en múltiples puntos de contacto.",
    "Alinear equipos y stakeholders con prioridades diversas, garantizando ejecución integrada y foco en resultados de negocio.",
  ],
  'negocios': 'OXXO México', 'alcance': 'Nacional', 'tramo': '5 reportes directos',
  'reportes': 'Analista Operación Teléfonos y Chips (Leonardo Hernández), Coordinador Product Owner CX (Verónica González), Coordinador Comercial Fidelización Clientes (Antonio Salazar), Analista Datos (Kleist Jimenez), Gerente Área Adquisición Cliente (Daniel Pardo)',
  'pares': 'Gerente Sr Fullfilment, Gerente Sr Comercial Servicios Financieros, Gerente Sr Comercial Comercio Electrónico',
  'relaciones': [
    ('INTERNAS', 'Director Estrategia y Transformación', 'Alineación estratégica y reporte de resultados del negocio OXXO Cel', 'Mensual'),
    ('INTERNAS', 'Gerente Sr Fullfilment', 'Coordinación de servicios y experiencia del cliente', 'Semanal'),
    ('INTERNAS', 'Áreas Comerciales y Operaciones', 'Ejecución comercial y operativa en tienda', 'Semanal'),
    ('INTERNAS', 'TI', 'Soporte tecnológico a plataformas y sistemas OXXO Cel', 'Quincenal'),
    ('EXTERNAS', 'Proveedores estratégicos (Telcel, fabricantes)', 'Negociación y gestión de relaciones comerciales', 'Mensual'),
  ],
},

# ══════════════════════════════════════════════════════════════════════════
# L3s – OXXO CEL
# ══════════════════════════════════════════════════════════════════════════
{
  '_folder': 'Gerente OXXO Cel',
  '_file':   'ET - L3 Analista Operacion Telefonos y Chips.pptx',
  'nombre_puesto': 'Analista Operación Teléfonos y Chips',
  'jefe_directo':  'Gerente OXXO Cel',
  'area':          AREA,
  'proposito': (
    "Asegurar la correcta operación del catálogo de teléfonos, chips y dispositivos de OXXO Cel, "
    "gestionando disponibilidad, calidad operativa y el seguimiento de indicadores clave para garantizar "
    "continuidad del servicio y una experiencia confiable al cliente."
  ),
  'responsabilidades': [
    "Gestionar y monitorear la disponibilidad del catálogo de teléfonos, chips y dispositivos en tienda, asegurando abasto oportuno y sin quiebres de stock.",
    "Coordinar los procesos operativos relacionados con la distribución, activación y soporte de teléfonos y chips, garantizando eficiencia y calidad.",
    "Dar seguimiento a indicadores clave de desempeño de la operación de dispositivos y chips, identificando desviaciones y proponiendo acciones correctivas.",
    "Gestionar incidencias operativas relacionadas con teléfonos, chips y dispositivos, escalando y coordinando su resolución oportuna.",
    "Coordinar con proveedores y áreas internas para asegurar la disponibilidad y correcta habilitación de nuevos dispositivos y SKUs.",
    "Apoyar la implementación de mejoras operativas en el proceso de gestión de teléfonos y chips, contribuyendo a la eficiencia del modelo OXXO Cel.",
  ],
  'decisiones': [
    "Definir prioridades de reabasto y distribución de teléfonos y chips según demanda y disponibilidad en tienda.",
    "Escalar incidencias críticas de operación de dispositivos conforme a su impacto en el servicio y la experiencia del cliente.",
  ],
  'retos': [
    "Mantener disponibilidad continua de un catálogo amplio de dispositivos en miles de tiendas con alta variabilidad de demanda.",
    "Coordinar con múltiples proveedores y áreas internas para garantizar abasto oportuno en un entorno con frecuentes cambios de portafolio.",
  ],
  'negocios': 'OXXO México', 'alcance': 'Nacional', 'tramo': 'Sin reportes directos',
  'reportes': '',
  'pares': 'Coordinador Product Owner CX, Coordinador Comercial Fidelización Clientes, Analista Datos, Gerente Área Adquisición Cliente',
  'relaciones': [
    ('INTERNAS', 'Gerente OXXO Cel', 'Reporte de operación de dispositivos y KPIs', 'Semanal'),
    ('INTERNAS', 'Cadena de Suministro', 'Coordinación logística y distribución de dispositivos', 'Semanal'),
    ('EXTERNAS', 'Proveedores de dispositivos', 'Gestión de abasto y soporte de dispositivos', 'Semanal'),
  ],
},
{
  '_folder': 'Gerente OXXO Cel',
  '_file':   'ET - L3 Coordinador Product Owner CX.pptx',
  'nombre_puesto': 'Coordinador Product Owner CX',
  'jefe_directo':  'Gerente OXXO Cel',
  'area':          AREA,
  'proposito': (
    "Gestionar la experiencia del cliente de OXXO Cel, coordinando el desarrollo del portafolio digital "
    "y la evolución de los puntos de contacto para asegurar una experiencia consistente, diferenciada y "
    "alineada a las necesidades del cliente."
  ),
  'responsabilidades': [
    "Gestionar el backlog de producto orientado a la experiencia del cliente, priorizando iniciativas de CX que generen mayor valor e impacto en la satisfacción.",
    "Coordinar con TI y áreas de diseño la evolución de los touchpoints digitales de OXXO Cel, asegurando usabilidad, consistencia y calidad.",
    "Dar seguimiento a métricas de experiencia del cliente (NPS, CSAT, churn) y proponer iniciativas de mejora basadas en datos.",
    "Coordinar lanzamientos de nuevas funcionalidades y evolutivos del producto digital, asegurando comunicación y habilitación adecuada al equipo comercial y operativo.",
    "Identificar oportunidades de mejora en el journey del cliente de OXXO Cel, desde la adquisición hasta la fidelización y retención.",
    "Gestionar la relación con equipos de UX, diseño y marketing para asegurar coherencia en la experiencia digital del cliente.",
  ],
  'decisiones': [
    "Priorizar el backlog de CX conforme al impacto en la experiencia del cliente y los objetivos de negocio de OXXO Cel.",
    "Definir los criterios de aceptación de funcionalidades de producto orientadas a la experiencia del cliente.",
  ],
  'retos': [
    "Asegurar una experiencia de cliente consistente y diferenciada en múltiples puntos de contacto (tienda, app, call center).",
    "Gestionar la evolución del producto digital en un entorno competitivo donde las expectativas del cliente cambian rápidamente.",
  ],
  'negocios': 'OXXO México', 'alcance': 'Nacional', 'tramo': 'Sin reportes directos',
  'reportes': '',
  'pares': 'Analista Operación Teléfonos y Chips, Coordinador Comercial Fidelización Clientes, Analista Datos, Gerente Área Adquisición Cliente',
  'relaciones': [
    ('INTERNAS', 'Gerente OXXO Cel', 'Alineación de roadmap de CX y métricas de experiencia', 'Semanal'),
    ('INTERNAS', 'TI', 'Coordinación de desarrollos y evolución de plataformas digitales', 'Semanal'),
    ('INTERNAS', 'Marketing', 'Coherencia de mensajes y experiencia del cliente', 'Quincenal'),
  ],
},
{
  '_folder': 'Gerente OXXO Cel',
  '_file':   'ET - L3 Coordinador Comercial Fidelizacion Clientes.pptx',
  'nombre_puesto': 'Coordinador Comercial Fidelización Clientes',
  'jefe_directo':  'Gerente OXXO Cel',
  'area':          AREA,
  'proposito': (
    "Coordinar las iniciativas de fidelización y retención de clientes de OXXO Cel, diseñando y ejecutando "
    "estrategias que incrementen la permanencia, reduzcan el churn y maximicen el valor del cliente "
    "a lo largo de su ciclo de vida."
  ),
  'responsabilidades': [
    "Diseñar y ejecutar estrategias de fidelización y retención de clientes de OXXO Cel, reduciendo el churn e incrementando la permanencia en la base de usuarios.",
    "Coordinar programas de beneficios, incentivos y comunicación que fortalezcan el vínculo del cliente con los servicios de OXXO Cel.",
    "Dar seguimiento a KPIs de retención, churn y valor del cliente (LTV), analizando tendencias y proponiendo mejoras a las estrategias de fidelización.",
    "Colaborar con las áreas comerciales y de producto para diseñar propuestas de valor diferenciadas que incentiven la permanencia del cliente.",
    "Gestionar campañas de comunicación y activación orientadas a clientes existentes, maximizando el engagement y la satisfacción.",
    "Identificar segmentos de clientes con mayor riesgo de churn y coordinar acciones proactivas de retención con base en analítica de datos.",
  ],
  'decisiones': [
    "Definir las estrategias y acciones de fidelización prioritarias para cada segmento de clientes de OXXO Cel.",
    "Establecer los incentivos y propuestas de valor que maximicen la retención en segmentos críticos.",
  ],
  'retos': [
    "Reducir el churn en un mercado competitivo donde los clientes tienen múltiples opciones de telefonía y servicios digitales.",
    "Diseñar propuestas de fidelización que sean rentables para el negocio y a la vez percibidas como valiosas por el cliente.",
  ],
  'negocios': 'OXXO México', 'alcance': 'Nacional', 'tramo': 'Sin reportes directos',
  'reportes': '',
  'pares': 'Analista Operación Teléfonos y Chips, Coordinador Product Owner CX, Analista Datos, Gerente Área Adquisición Cliente',
  'relaciones': [
    ('INTERNAS', 'Gerente OXXO Cel', 'Alineación de estrategias de fidelización y métricas de retención', 'Semanal'),
    ('INTERNAS', 'Marketing', 'Coordinación de campañas y comunicación a clientes', 'Semanal'),
    ('INTERNAS', 'Analítica de Datos', 'Segmentación de clientes y análisis de churn', 'Quincenal'),
  ],
},
{
  '_folder': 'Gerente OXXO Cel',
  '_file':   'ET - L3 Analista Datos.pptx',
  'nombre_puesto': 'Analista Datos',
  'jefe_directo':  'Gerente OXXO Cel',
  'area':          AREA,
  'proposito': (
    "Generar y gestionar información estratégica y operativa del negocio OXXO Cel mediante el análisis "
    "de datos, construyendo reportes, dashboards y modelos analíticos que habiliten la toma de decisiones "
    "basada en evidencia en todas las áreas del negocio."
  ),
  'responsabilidades': [
    "Desarrollar y mantener dashboards y reportes de KPIs del negocio OXXO Cel, asegurando disponibilidad oportuna de información estratégica y operativa.",
    "Realizar análisis de datos de desempeño comercial, operativo y de experiencia del cliente, identificando tendencias, oportunidades y áreas de mejora.",
    "Apoyar a los equipos de OXXO Cel con modelos analíticos y segmentaciones de clientes para orientar decisiones de adquisición, fidelización y operación.",
    "Gestionar y validar la calidad de los datos utilizados en los reportes y análisis, coordinando con TI la disponibilidad y confiabilidad de las fuentes.",
    "Elaborar análisis ad hoc y estudios de negocio que soporten la definición de estrategias y la evaluación de iniciativas del portafolio OXXO Cel.",
    "Automatizar procesos de generación de reportes y análisis recurrentes, incrementando la eficiencia del área y la frecuencia de actualización de la información.",
  ],
  'decisiones': [
    "Definir la metodología y fuentes de datos para cada análisis o modelo desarrollado, asegurando rigor y confiabilidad.",
    "Establecer la priorización de los reportes y análisis a desarrollar conforme a las necesidades del equipo y el impacto en la toma de decisiones.",
  ],
  'retos': [
    "Asegurar la calidad y confiabilidad de los datos en un entorno con múltiples fuentes y sistemas heterogéneos.",
    "Traducir grandes volúmenes de datos en insights accionables y comprensibles para equipos con distintos niveles de familiaridad analítica.",
  ],
  'negocios': 'OXXO México', 'alcance': 'Nacional', 'tramo': 'Sin reportes directos',
  'reportes': '',
  'pares': 'Analista Operación Teléfonos y Chips, Coordinador Product Owner CX, Coordinador Comercial Fidelización Clientes, Gerente Área Adquisición Cliente',
  'relaciones': [
    ('INTERNAS', 'Gerente OXXO Cel', 'Entrega de reportes y análisis estratégicos del negocio', 'Semanal'),
    ('INTERNAS', 'TI', 'Acceso a fuentes de datos e integración de sistemas de información', 'Semanal'),
    ('INTERNAS', 'Todas las áreas de OXXO Cel', 'Soporte analítico para toma de decisiones', 'Semanal'),
  ],
},
{
  '_folder': 'Gerente OXXO Cel',
  '_file':   'ET - L3 Gerente Area Adquisicion Cliente.pptx',
  'nombre_puesto': 'Gerente Área Adquisición Cliente',
  'jefe_directo':  'Gerente OXXO Cel',
  'area':          AREA,
  'proposito': (
    "Liderar la estrategia y ejecución de adquisición de nuevos clientes para OXXO Cel, diseñando y "
    "coordinando iniciativas comerciales, campañas y activaciones en tienda que maximicen el crecimiento "
    "de la base de usuarios y el volumen de nuevas activaciones."
  ),
  'responsabilidades': [
    "Definir y ejecutar la estrategia de adquisición de nuevos clientes de OXXO Cel, identificando segmentos objetivo y canales más efectivos para incrementar activaciones.",
    "Coordinar campañas comerciales de adquisición en tienda y canales digitales, asegurando coherencia en mensajes y ejecución consistente en el punto de venta.",
    "Gestionar la relación con proveedores y categorías para asegurar la correcta ejecución en tienda de actividades de adquisición.",
    "Dar seguimiento a KPIs de adquisición de clientes, incluyendo nuevas activaciones, market share y costo de adquisición, proponiendo ajustes a la estrategia.",
    "Coordinar con los equipos de marketing, operaciones y socios estratégicos la implementación de iniciativas de activación y lanzamientos de nuevos servicios.",
    "Analizar la competencia y tendencias del mercado de telefonía para identificar oportunidades de diferenciación en la estrategia de adquisición.",
  ],
  'decisiones': [
    "Definir la estrategia y prioridades de adquisición de clientes por segmento, canal y periodo, balanceando volumen y costo.",
    "Determinar las actividades de activación en tienda y los apoyos comerciales necesarios para cada campaña de adquisición.",
  ],
  'retos': [
    "Crecer la base de clientes de OXXO Cel en un mercado saturado y competitivo, donde la diferenciación de la propuesta de valor es clave.",
    "Asegurar consistencia en la ejecución de estrategias de adquisición en miles de puntos de venta con realidades comerciales muy diversas.",
  ],
  'negocios': 'OXXO México', 'alcance': 'Nacional', 'tramo': 'Sin reportes directos',
  'reportes': '',
  'pares': 'Analista Operación Teléfonos y Chips, Coordinador Product Owner CX, Coordinador Comercial Fidelización Clientes, Analista Datos',
  'relaciones': [
    ('INTERNAS', 'Gerente OXXO Cel', 'Alineación de estrategia de adquisición y resultados', 'Semanal'),
    ('INTERNAS', 'Marketing', 'Coordinación de campañas y comunicación de adquisición', 'Semanal'),
    ('INTERNAS', 'Socios y Categorías', 'Ejecución en tienda y abasto de dispositivos para activación', 'Semanal'),
  ],
},

# ══════════════════════════════════════════════════════════════════════════
# L2 – GERENTE SR COMERCIAL SERVICIOS FINANCIEROS
# ══════════════════════════════════════════════════════════════════════════
{
  '_folder': 'Gerente Sr Comercial Servicios Financieros',
  '_file':   'ET - L2 Gerente Sr Comercial Servicios Financieros.pptx',
  'nombre_puesto': 'Gerente Sr Comercial Servicios Financieros',
  'jefe_directo':  DIRECTOR,
  'area':          AREA,
  'proposito': (
    "Diseñar, desarrollar y maximizar el valor del portafolio de servicios financieros del retail, "
    "asegurando crecimiento rentable y sostenible del margen anual a través de la innovación en servicios, "
    "la excelencia operativa en tienda y la orquestación efectiva del ecosistema financiero, fortaleciendo "
    "la propuesta de valor para clientes, tiendas y socios estratégicos en México."
  ),
  'responsabilidades': [
    "Definir la estrategia integral de servicios financieros para maximizar el crecimiento rentable y la contribución al margen del negocio, mediante la priorización de iniciativas, alineación con objetivos corporativos y seguimiento a indicadores estratégicos.",
    "Gestionar la rentabilidad del portafolio de servicios financieros actuales y nuevos para incrementar el margen anual, a través de modelos de negocio, esquemas tarifarios y optimización de costos operativos.",
    "Impulsar el diseño, lanzamiento y escalamiento de nuevos productos financieros para diversificar fuentes de ingreso y fortalecer la propuesta de valor al cliente, mediante la gestión end-to-end del ciclo de producto y la evaluación de business cases.",
    "Liderar la gestión comercial y negociación con socios financieros estratégicos para asegurar condiciones competitivas y modelos de colaboración sostenibles, mediante acuerdos comerciales, administración de contratos y evaluación de desempeño de aliados.",
    "Asegurar la correcta ejecución operativa de los servicios financieros en tiendas para garantizar continuidad operativa y una experiencia confiable al cliente, a través de modelos de habilitación, seguimiento a KPIs operativos y estandarización de procesos.",
    "Implementar iniciativas de mejora continua en la operación de servicios financieros para reducir fricción, errores y costos, mediante análisis de procesos, automatización y coordinación con las áreas de Operaciones y Tecnología.",
    "Garantizar el cumplimiento regulatorio y la adecuada gestión de riesgos asociados a los servicios financieros, mediante la coordinación con áreas de Compliance, Legal, Asuntos Corporativos y atención a auditorías.",
    "Coordinar a las áreas comerciales, de producto, operaciones, tecnología y corporativas para asegurar la ejecución alineada de la estrategia de servicios financieros, a través de gobiernos de proyecto, comités y gestión transversal de stakeholders.",
  ],
  'decisiones': [
    "Definir el portafolio de servicios financieros a desarrollar, escalar o descontinuar con base en rentabilidad, viabilidad operativa y alineación estratégica.",
    "Seleccionar socios financieros y los modelos de asociación más convenientes para cada servicio.",
    "Definir la estrategia comercial y de monetización de los servicios financieros (tarifas, precios, incentivos y condiciones comerciales).",
    "Priorizar el roadmap de iniciativas y productos, balanceando impacto financiero, time to market, riesgo regulatorio y capacidades operativas.",
  ],
  'retos': [
    "Mantener el crecimiento sostenido del margen y la contribución de los servicios financieros en un entorno altamente competitivo y regulado.",
    "Desarrollar y escalar nuevas avenidas de crecimiento (nuevos servicios) con impacto real en el negocio.",
    "Reducir el time to market de nuevos servicios financieros sin comprometer cumplimiento ni estabilidad operativa.",
    "Alinear de manera efectiva a servicios, operaciones, tecnología y áreas corporativas para una ejecución integrada.",
    "Lograr la alineación del equipo SPIN para el desarrollo de nuevas iniciativas, asegurando que las propuestas se construyan a partir del producto de casa y fortalezcan la marca interna.",
  ],
  'negocios': 'OXXO México', 'alcance': 'Nacional', 'tramo': '5 reportes directos',
  'reportes': 'Responsable Comercial Medios Pago (Vacante), Gerente Área Operaciones iCash (Juan Manuel Vazquez), Gerente Área Comercial Dispersión Efectivo (Jorge Castaños), Gerente Área Producto (José Manuel Vallejo), Gerente Área Comercial Corresponsalías (Fernando Bernal)',
  'pares': 'Gerente Sr Fullfilment, Gerente OXXO Cel, Gerente Sr Comercial Comercio Electrónico',
  'relaciones': [
    ('INTERNAS', 'Director Estrategia y Transformación', 'Alineación estratégica y rendición de cuentas del área', 'Mensual'),
    ('INTERNAS', 'TI', 'Integración tecnológica y habilitación de servicios financieros', 'Semanal'),
    ('INTERNAS', 'Finanzas y Legal', 'Cumplimiento regulatorio y gestión de riesgos financieros', 'Mensual'),
    ('INTERNAS', 'Operaciones', 'Ejecución y habilitación de servicios financieros en tienda', 'Quincenal'),
    ('EXTERNAS', 'Socios Financieros (bancos, fintechs)', 'Negociación y desarrollo de nuevos productos y servicios', 'Mensual'),
  ],
},

# ══════════════════════════════════════════════════════════════════════════
# L3s – SERVICIOS FINANCIEROS
# ══════════════════════════════════════════════════════════════════════════
{
  '_folder': 'Gerente Sr Comercial Servicios Financieros',
  '_file':   'ET - L3 Responsable Comercial Medios Pago.pptx',
  'nombre_puesto': 'Responsable Comercial Medios Pago',
  'jefe_directo':  'Gerente Sr Comercial Servicios Financieros',
  'area':          AREA,
  'proposito': (
    "Gestionar y desarrollar la estrategia comercial de los medios de pago de OXXO, asegurando el "
    "crecimiento del portafolio de servicios de pago, la optimización de las relaciones con socios "
    "estratégicos y el cumplimiento de los objetivos de volumen y rentabilidad."
  ),
  'responsabilidades': [
    "Gestionar la estrategia comercial del portafolio de medios de pago, asegurando crecimiento de volumen, rentabilidad y adopción por parte de los clientes.",
    "Coordinar con socios financieros estratégicos (redes de pago, procesadoras, bancos) la negociación y mantenimiento de condiciones comerciales competitivas.",
    "Dar seguimiento a KPIs comerciales de medios de pago, analizando desempeño y proponiendo iniciativas de mejora para cumplir los objetivos del área.",
    "Identificar y desarrollar nuevas oportunidades de negocio en medios de pago, evaluando viabilidad comercial y propuesta de valor para clientes y tiendas.",
    "Coordinar con TI, Operaciones y Compliance la implementación de nuevos medios de pago y la resolución de incidencias comerciales.",
    "Gestionar la relación con socios de medios de pago, asegurando cumplimiento de acuerdos y evolución de la colaboración estratégica.",
  ],
  'decisiones': [
    "Definir prioridades en el desarrollo y negociación de nuevos acuerdos comerciales de medios de pago.",
    "Establecer estrategias comerciales para maximizar adopción y volumen de transacciones en el portafolio de medios de pago.",
  ],
  'retos': [
    "Mantener competitividad del portafolio de medios de pago frente a la evolución acelerada del mercado fintech y digital.",
    "Alinear las necesidades comerciales con los requerimientos de cumplimiento regulatorio en un entorno de alta regulación.",
  ],
  'negocios': 'OXXO México', 'alcance': 'Nacional', 'tramo': 'Sin reportes directos',
  'reportes': '',
  'pares': 'Gerente Área Operaciones iCash, Gerente Área Comercial Dispersión Efectivo, Gerente Área Producto, Gerente Área Comercial Corresponsalías',
  'relaciones': [
    ('INTERNAS', 'Gerente Sr Comercial Servicios Financieros', 'Alineación de estrategia y KPIs de medios de pago', 'Semanal'),
    ('INTERNAS', 'TI y Operaciones', 'Implementación y soporte de medios de pago', 'Quincenal'),
    ('EXTERNAS', 'Socios Financieros (redes de pago)', 'Negociación y gestión de relaciones comerciales', 'Mensual'),
  ],
},
{
  '_folder': 'Gerente Sr Comercial Servicios Financieros',
  '_file':   'ET - L3 Gerente Area Operaciones iCash.pptx',
  'nombre_puesto': 'Gerente Área Operaciones iCash',
  'jefe_directo':  'Gerente Sr Comercial Servicios Financieros',
  'area':          AREA,
  'proposito': (
    "Gestionar y optimizar la operación del servicio iCash, asegurando la correcta implementación, "
    "funcionamiento y mantenimiento de las cajas recicladoras y procesos de cash management en tienda, "
    "maximizando eficiencia operativa y rentabilidad del servicio."
  ),
  'responsabilidades': [
    "Gestionar la operación integral del servicio iCash, coordinando la instalación, funcionamiento y mantenimiento de cajas recicladoras en tienda.",
    "Dar seguimiento a KPIs operativos de iCash, incluyendo número de cajas instaladas y encendidas, disponibilidad y rentabilidad del servicio.",
    "Coordinar con TI, Finanzas y Operaciones los procesos de conciliación, reporte y escalamiento de incidencias del servicio iCash.",
    "Gestionar la relación con proveedores de equipos y servicios de cash management, asegurando soporte oportuno y condiciones contractuales adecuadas.",
    "Identificar e implementar mejoras operativas en el proceso de iCash que reduzcan costos, tiempos de inactividad y riesgo operativo.",
    "Asegurar el cumplimiento regulatorio y los protocolos de seguridad relacionados con el manejo de efectivo en las operaciones de iCash.",
  ],
  'decisiones': [
    "Definir la priorización de instalación y mantenimiento de cajas iCash según rentabilidad y necesidades operativas de cada tienda.",
    "Establecer protocolos de respuesta ante incidencias críticas de operación iCash que afecten la continuidad del servicio.",
  ],
  'retos': [
    "Escalar la implementación de iCash cumpliendo los planes de instalación sin comprometer la calidad y disponibilidad del servicio.",
    "Asegurar la rentabilidad del servicio iCash balanceando costos de operación, mantenimiento e ingresos generados por las cajas.",
  ],
  'negocios': 'OXXO México', 'alcance': 'Nacional', 'tramo': 'Sin reportes directos',
  'reportes': '',
  'pares': 'Responsable Comercial Medios Pago, Gerente Área Comercial Dispersión Efectivo, Gerente Área Producto, Gerente Área Comercial Corresponsalías',
  'relaciones': [
    ('INTERNAS', 'Gerente Sr Comercial Servicios Financieros', 'Reporte de KPIs operativos y expansión de iCash', 'Semanal'),
    ('INTERNAS', 'TI', 'Integración tecnológica y soporte de plataformas iCash', 'Semanal'),
    ('INTERNAS', 'Operaciones y Despliegue', 'Coordinación de instalación y mantenimiento en tienda', 'Semanal'),
  ],
},
{
  '_folder': 'Gerente Sr Comercial Servicios Financieros',
  '_file':   'ET - L3 Gerente Area Comercial Dispersion Efectivo.pptx',
  'nombre_puesto': 'Gerente Área Comercial Dispersión Efectivo',
  'jefe_directo':  'Gerente Sr Comercial Servicios Financieros',
  'area':          AREA,
  'proposito': (
    "Gestionar y desarrollar el negocio de dispersión de efectivo, asegurando el crecimiento comercial "
    "del portafolio, la rentabilidad del servicio y la correcta ejecución operativa, fortaleciendo las "
    "relaciones con socios estratégicos y maximizando el valor para clientes, tiendas y el negocio."
  ),
  'responsabilidades': [
    "Gestionar la estrategia comercial del servicio de dispersión de efectivo, asegurando crecimiento de volumen, margen y adopción del servicio.",
    "Desarrollar y mantener relaciones con socios estratégicos de dispersión de efectivo, negociando condiciones comerciales competitivas y modelos de colaboración sostenibles.",
    "Dar seguimiento a KPIs comerciales y financieros del servicio de dispersión, proponiendo iniciativas de mejora para cumplir los objetivos del área.",
    "Coordinar con TI y Operaciones la correcta ejecución operativa del servicio de dispersión en tienda, resolviendo incidencias y mejorando procesos.",
    "Identificar oportunidades de crecimiento en el mercado de dispersión de efectivo, evaluando nuevos segmentos, socios y modalidades de servicio.",
    "Asegurar el cumplimiento regulatorio y de riesgos en todas las operaciones de dispersión de efectivo, coordinando con Compliance y Legal.",
  ],
  'decisiones': [
    "Definir estrategias comerciales y de negociación con socios del servicio de dispersión de efectivo.",
    "Priorizar el desarrollo de nuevas oportunidades de negocio en dispersión conforme a rentabilidad e impacto estratégico.",
  ],
  'retos': [
    "Crecer el volumen de dispersión de efectivo en un entorno de creciente digitalización de pagos y alta competencia.",
    "Gestionar los riesgos regulatorios y operativos del servicio cumpliendo estrictamente con la normativa aplicable.",
  ],
  'negocios': 'OXXO México', 'alcance': 'Nacional', 'tramo': 'Sin reportes directos',
  'reportes': '',
  'pares': 'Responsable Comercial Medios Pago, Gerente Área Operaciones iCash, Gerente Área Producto, Gerente Área Comercial Corresponsalías',
  'relaciones': [
    ('INTERNAS', 'Gerente Sr Comercial Servicios Financieros', 'Alineación de estrategia y KPIs de dispersión de efectivo', 'Semanal'),
    ('INTERNAS', 'Compliance y Legal', 'Cumplimiento regulatorio del servicio de dispersión', 'Mensual'),
    ('EXTERNAS', 'Socios Estratégicos (empresas dispersoras)', 'Negociación y desarrollo de relación comercial', 'Mensual'),
  ],
},
{
  '_folder': 'Gerente Sr Comercial Servicios Financieros',
  '_file':   'ET - L3 Gerente Area Producto.pptx',
  'nombre_puesto': 'Gerente Área Producto',
  'jefe_directo':  'Gerente Sr Comercial Servicios Financieros',
  'area':          AREA,
  'proposito': (
    "Liderar el desarrollo y gestión del portafolio de nuevos productos y servicios financieros, "
    "definiendo el ciclo end-to-end desde la ideación hasta el lanzamiento, evaluando business cases y "
    "asegurando que las nuevas soluciones generen valor diferenciado para clientes, tiendas y el negocio."
  ),
  'responsabilidades': [
    "Gestionar el portafolio de nuevos productos financieros, coordinando el ciclo end-to-end desde la conceptualización, desarrollo, evaluación hasta el lanzamiento de nuevas soluciones.",
    "Elaborar y validar business cases de nuevos productos financieros, asegurando viabilidad comercial, financiera y operativa antes de su lanzamiento.",
    "Coordinar con las áreas de TI, Operaciones, Legal y Compliance el desarrollo e implementación de nuevos productos financieros.",
    "Dar seguimiento al desempeño de los productos lanzados, monitoreando KPIs de adopción, rentabilidad y experiencia del cliente, proponiendo mejoras.",
    "Identificar tendencias del mercado financiero y fintech que representen oportunidades de innovación en el portafolio de OXXO.",
    "Gestionar la documentación de productos financieros, incluyendo especificaciones, procesos, materiales de capacitación y comunicación para socios y operaciones.",
  ],
  'decisiones': [
    "Definir la viabilidad y prioridad de desarrollo de nuevos productos financieros con base en business cases y alineación estratégica.",
    "Establecer los criterios de lanzamiento y los KPIs de éxito para cada nuevo producto del portafolio financiero.",
  ],
  'retos': [
    "Acelerar el time to market de nuevos productos financieros sin comprometer el rigor regulatorio y la calidad operativa.",
    "Desarrollar productos innovadores que se diferencien en un mercado financiero altamente competitivo y regulado.",
  ],
  'negocios': 'OXXO México', 'alcance': 'Nacional', 'tramo': 'Sin reportes directos',
  'reportes': '',
  'pares': 'Responsable Comercial Medios Pago, Gerente Área Operaciones iCash, Gerente Área Comercial Dispersión Efectivo, Gerente Área Comercial Corresponsalías',
  'relaciones': [
    ('INTERNAS', 'Gerente Sr Comercial Servicios Financieros', 'Alineación de roadmap de producto y evaluación de business cases', 'Semanal'),
    ('INTERNAS', 'TI', 'Desarrollo tecnológico de nuevos productos financieros', 'Semanal'),
    ('INTERNAS', 'Legal y Compliance', 'Validación regulatoria de nuevos productos financieros', 'Mensual'),
  ],
},
{
  '_folder': 'Gerente Sr Comercial Servicios Financieros',
  '_file':   'ET - L3 Gerente Area Comercial Corresponsalias.pptx',
  'nombre_puesto': 'Gerente Área Comercial Corresponsalías',
  'jefe_directo':  'Gerente Sr Comercial Servicios Financieros',
  'area':          AREA,
  'proposito': (
    "Gestionar y desarrollar el portafolio de servicios de corresponsalías bancarias de OXXO, asegurando "
    "el crecimiento comercial, la operación eficiente del servicio y el fortalecimiento de las relaciones "
    "con socios bancarios, para maximizar el valor del servicio para clientes, tiendas y el negocio."
  ),
  'responsabilidades': [
    "Gestionar la estrategia comercial del portafolio de corresponsalías bancarias, asegurando crecimiento de volumen, margen y disponibilidad del servicio.",
    "Desarrollar y mantener relaciones con socios bancarios, negociando condiciones comerciales, comisiones y modelos de colaboración que maximicen el valor del servicio.",
    "Dar seguimiento a KPIs comerciales de corresponsalías, incluyendo transacciones, comisiones y satisfacción del cliente, proponiendo acciones de mejora.",
    "Coordinar con TI y Operaciones la correcta ejecución operativa del servicio de corresponsalías en tienda, asegurando disponibilidad y calidad del servicio.",
    "Asegurar el cumplimiento regulatorio y los requerimientos de auditoría bancaria en todas las operaciones de corresponsalías, coordinando con Compliance y Legal.",
    "Identificar oportunidades de expansión del portafolio de corresponsalías, incorporando nuevos bancos o servicios que amplíen la propuesta de valor para el cliente.",
  ],
  'decisiones': [
    "Definir la estrategia de negociación y desarrollo de nuevos acuerdos con socios bancarios de corresponsalías.",
    "Priorizar la incorporación de nuevos bancos o servicios al portafolio de corresponsalías conforme a rentabilidad e impacto al cliente.",
  ],
  'retos': [
    "Mantener y crecer el volumen de transacciones de corresponsalías en un contexto de mayor digitalización bancaria y alternativas de pago.",
    "Gestionar la complejidad regulatoria y los requerimientos de auditoría de múltiples socios bancarios simultáneamente.",
  ],
  'negocios': 'OXXO México', 'alcance': 'Nacional', 'tramo': 'Sin reportes directos',
  'reportes': '',
  'pares': 'Responsable Comercial Medios Pago, Gerente Área Operaciones iCash, Gerente Área Comercial Dispersión Efectivo, Gerente Área Producto',
  'relaciones': [
    ('INTERNAS', 'Gerente Sr Comercial Servicios Financieros', 'Alineación de estrategia y KPIs de corresponsalías', 'Semanal'),
    ('INTERNAS', 'Compliance y Legal', 'Cumplimiento de requerimientos regulatorios bancarios', 'Mensual'),
    ('EXTERNAS', 'Socios Bancarios', 'Negociación y gestión de relaciones de corresponsalía', 'Mensual'),
  ],
},

]  # end roles list

# ── generate all files ──────────────────────────────────────────────────

print("Generating ET Descripciones de Puestos...\n")
for role in roles:
    generate_dp(role, role['_folder'], role['_file'])

print(f"\nDone! {len(roles)} files generated.")
