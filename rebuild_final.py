#!/usr/bin/env python3
"""
Final Finanzas Playbook restructure:
 - Slide 6:   NEW L1→L2 Finanzas overview (repurposed from old idx 59)
 - Slides 7-52: RASCI (unchanged)
 - Slides 53-59: 7 L2→L3 org charts (moved from idx 5-11, fixed visuals)
 - Slides 60-65: old bad slides (cleared + hidden)
 - Slides 66-67: Anexos, Gracias

Visual fixes on L2→L3 slides:
 - Title and subtitle no longer overlap
 - Standardized compact description box height
 - Better font sizes (no 7pt)
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import zipfile, re

PPTX_PATH = '/home/user/oxxo/Finanzas/Playbook/Finanzas Playbook.pptx'
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'

SLIDE_W = 12192000
SLIDE_H = 6858000
SIDE_M  = 304800

# ── L2 DATA (for the 7 org chart slides) ─────────────────────
L2_DATA = [
    {
        'subtitle': 'PF&A y Commercial Finance',
        'l2_title': 'Gte Sr Planeación Financiera\ny Análisis & Commercial Finance',
        'l3s': [
            {'title': 'Gerente BP Finanzas\nPlataformas RH, TI',
             'bullets': ['Elaborar y dar seguimiento al presupuesto anual y las previsiones de las plataformas RH y TI.',
                         'Coordinar el análisis de variaciones del desempeño financiero de RH y TI.',
                         'Generar información y asesoramiento sobre el impacto financiero de las decisiones de RH y TI.',
                         'Elaborar reportes financieros periódicos de las plataformas RH y TI.']},
            {'title': 'Gerente Commercial\nFinance',
             'bullets': ['Elaborar el modelamiento de trade-offs económicos para respaldar decisiones comerciales.',
                         'Monitorear y analizar el rendimiento del margen bruto, la rentabilidad y sus drivers.',
                         'Identificar y reportar riesgos financieros con impacto potencial en la rentabilidad.',
                         'Elaborar reportes de visibilidad unificada de ingresos y márgenes para los socios comerciales.']},
            {'title': 'Gerente Planeación\nFinanciera',
             'bullets': ['Coordinar el proceso de presupuestación anual y previsiones continuas de las funciones asignadas.',
                         'Elaborar el análisis de variaciones del desempeño financiero real versus lo planeado.',
                         'Desarrollar modelos financieros para la planeación, análisis de escenarios y pruebas de sensibilidad.',
                         'Generar los reportes financieros periódicos para los responsables de función y la dirección.']},
        ]
    },
    {
        'subtitle': 'Fiscal y Contraloría',
        'l2_title': 'Gte Sr Fiscal y Contraloría',
        'l3s': [
            {'title': 'Gerente Área\nNormatividad y Control',
             'bullets': ['Diseñar e implementar el plan anual de normatividad y control interno.',
                         'Coordinar la ejecución de revisiones de control interno en las unidades de negocio.',
                         'Elaborar y reportar los resultados del sistema de control interno.',
                         'Homologar criterios y políticas contables alineados con FEMSA e IFRS.']},
            {'title': 'Gerente Control\nFiscal',
             'bullets': ['Ejecutar el cumplimiento de las obligaciones fiscales en materia de impuestos federales, estatales y locales.',
                         'Implementar y mantener las herramientas de documentación y respaldo de los procesos fiscales.',
                         'Coordinar con Auditorías Fiscales la preparación de información para revisiones de autoridades.',
                         'Elaborar análisis de los cambios en la legislación fiscal y su aplicación en la operación.']},
            {'title': 'Gerente Gestión\nRiesgos SI',
             'bullets': ['Diseñar e implementar la metodología de prevención de lavado de dinero y financiamiento al terrorismo.',
                         'Desarrollar y mantener los modelos de detección y reducción de fraude monetario.',
                         'Coordinar la implementación del sistema integrado de gestión de riesgos.',
                         'Elaborar reportes de exposición al riesgo y efectividad de los controles de prevención.']},
            {'title': 'Gerente Sr\nContraloría',
             'bullets': ['Supervisar la emisión de información financiera confiable y oportuna bajo normas IFRS.',
                         'Coordinar la representación del negocio ante auditoría interna y externa.',
                         'Implementar y monitorear el sistema de control interno orientado a la integridad de reportes financieros.',
                         'Coordinar la consolidación financiera de OXXO México.']},
            {'title': 'Responsable\nAuditorías Fiscales',
             'bullets': ['Coordinar la atención y seguimiento a requerimientos de autoridades fiscales.',
                         'Ejecutar el seguimiento a las auditorías externas e internas en materia fiscal.',
                         'Elaborar análisis de cumplimiento de las obligaciones fiscales de OXXO México.',
                         'Coordinar la recuperación de saldos de impuestos a favor.']},
        ]
    },
    {
        'subtitle': 'Servicios Compartidos',
        'l2_title': 'Gte Sr Servicios Compartidos\nAdministrativo',
        'l3s': [
            {'title': 'Coordinador\nInformación y Gestión CSF',
             'bullets': ['Coordinar la consolidación y reporte de los indicadores de desempeño del Centro de Servicios.',
                         'Dar seguimiento al cumplimiento de los acuerdos de nivel de servicio con los clientes internos.',
                         'Coordinar los rituales de gestión del Centro de Servicios.',
                         'Apoyar la gestión de la relación con los clientes internos del Centro de Servicios.']},
            {'title': 'Gerente Order\nto Cash',
             'bullets': ['Ejecutar y supervisar el ciclo Order to Cash end-to-end.',
                         'Coordinar la gestión de la cartera de cuentas por cobrar.',
                         'Implementar y mantener los controles del proceso O2C.',
                         'Elaborar y reportar los indicadores de desempeño del proceso O2C.']},
            {'title': 'Gerente Purchase\nto Pay',
             'bullets': ['Ejecutar y supervisar el ciclo Purchase to Pay end-to-end.',
                         'Coordinar la optimización de los flujos de aprobación, condiciones de pago y gestión de caja.',
                         'Implementar y mantener los controles del proceso P2P.',
                         'Gestionar la relación operativa con proveedores en el ámbito del proceso P2P.']},
            {'title': 'Gerente Record\nto Report',
             'bullets': ['Ejecutar y supervisar el ciclo Record to Report end-to-end.',
                         'Coordinar el proceso de cierre contable mensual y anual.',
                         'Implementar y mantener los controles del proceso R2R.',
                         'Coordinar con Contraloría la alineación de criterios contables y el proceso de auditoría.']},
            {'title': 'Responsable\nGobierno Datos y Cond',
             'bullets': ['Implementar y mantener los estándares de gobierno de datos del Centro de Servicios.',
                         'Coordinar el mantenimiento y actualización de los catálogos y condiciones comerciales.',
                         'Identificar y reportar desviaciones en la calidad de datos que impacten los procesos transaccionales.',
                         'Coordinar con las áreas usuarias la resolución de inconsistencias en datos y condiciones.']},
        ]
    },
    {
        'subtitle': 'Legal y Regulatorio',
        'l2_title': 'Gte Sr Legal y Regulatorio',
        'l3s': [
            {'title': 'Gerente Área\nCumplimiento',
             'bullets': ['Diseñar e implementar el programa de cumplimiento legal y regulatorio de OXXO México.',
                         'Coordinar el monitoreo del cumplimiento de las obligaciones legales y regulatorias en las unidades de negocio.',
                         'Elaborar y reportar los indicadores de desempeño del programa de cumplimiento.',
                         'Coordinar con las áreas del negocio la implementación de acciones correctivas derivadas de hallazgos.']},
            {'title': 'Gerente Área\nLegal Comercial',
             'bullets': ['Ejecutar la revisión y análisis legal de contratos y acuerdos comerciales de OXXO.',
                         'Gestionar las contingencias y litigios legales del área comercial.',
                         'Elaborar y actualizar los formatos de contratos comerciales de OXXO.',
                         'Asesorar al área comercial en procesos legales, negociaciones y contingencias.']},
            {'title': 'Gerente Legal CdS',
             'bullets': ['Ejecutar la revisión y análisis legal de contratos y acuerdos de la Cadena de Suministro.',
                         'Gestionar las contingencias y litigios legales de la Cadena de Suministro.',
                         'Coordinar con el área de Gestoría la atención de trámites legales y regulatorios de CEDIS y Distribución.',
                         'Asesorar al área de Cadena de Suministro en procesos legales, negociaciones y contingencias.']},
            {'title': 'Gerente Sr Legal\nOperaciones',
             'bullets': ['Ejecutar el análisis y atención de los requerimientos legales de la operación de tiendas y plazas.',
                         'Gestionar los litigios y procesos legales operativos de OXXO a nivel nacional.',
                         'Elaborar y revisar los documentos legales requeridos por la operación de tiendas.',
                         'Coordinar con los Gerentes de Plaza y Región la atención de contingencias legales operativas.']},
            {'title': 'Responsable\nLegal (i78)',
             'bullets': ['Ejecutar el análisis y atención de los requerimientos legales de la plataforma i78.',
                         'Elaborar y revisar los documentos legales requeridos por la operación de i78.',
                         'Coordinar con las áreas operativas de i78 la atención de contingencias legales.',
                         'Dar seguimiento a los litigios y procesos legales activos de la plataforma i78.']},
        ]
    },
    {
        'subtitle': 'Asuntos Corporativos',
        'l2_title': 'Gte Sr Asuntos Corporativos',
        'l3s': [
            {'title': 'Analista\nComunicación Externa',
             'bullets': ['Elaborar contenidos y materiales de comunicación externa para medios y audiencias clave.',
                         'Coordinar el monitoreo de medios de comunicación y redes sociales.',
                         'Elaborar reportes de monitoreo de medios y análisis de cobertura.',
                         'Apoyar la implementación de la estrategia de comunicación proactiva y reactiva del área.']},
            {'title': 'Coordinador\nAsuntos Corporativos',
             'bullets': ['Elaborar análisis de riesgo e informes políticos basados en tendencias regulatorias del entorno.',
                         'Coordinar el monitoreo sistemático del entorno político, económico y social.',
                         'Elaborar reportes de la agenda de riesgos del entorno para la toma de decisiones.',
                         'Coordinar la gestión de proyectos con consultores especializados en políticas públicas.']},
            {'title': 'Gerente Área\nAtención Autoridades',
             'bullets': ['Gestionar y dar seguimiento a los requerimientos de autoridades gubernamentales hacia OXXO.',
                         'Elaborar las respuestas y documentación requerida ante requerimientos de autoridades.',
                         'Coordinar con las áreas operativas y de soporte la atención de requerimientos de autoridades.',
                         'Monitorear y reportar el estatus de los requerimientos de autoridades en curso.']},
            {'title': 'Gerente Área\nRelaciones Gobierno',
             'bullets': ['Ejecutar los planes de relacionamiento con autoridades gubernamentales en la región o ámbito asignado.',
                         'Monitorear el entorno político-regulatorio local y reportar tendencias relevantes para el negocio.',
                         'Participar en organismos empresariales y foros de influencia en representación de OXXO.',
                         'Ejecutar proyectos de cabildeo e incidencia en política pública en el ámbito asignado.']},
            {'title': 'Gerente\nRelaciones Gobierno',
             'bullets': ['Coordinar la implementación de la estrategia de relacionamiento con autoridades gubernamentales.',
                         'Elaborar y dar seguimiento a los planes de relacionamiento con autoridades y actores clave.',
                         'Coordinar el monitoreo del entorno regulatorio y político relevante para el negocio.',
                         'Gestionar proyectos de cabildeo y relacionamiento institucional con consultores especializados.']},
        ]
    },
    {
        'subtitle': 'Administración Nacional OXXO',
        'l2_title': 'Director Nacional\nAdmvo OXXO',
        'l3s': [
            {'title': 'Gerente Área\nAdmvo Neg Dllo',
             'bullets': ['Implementar los procesos administrativos y financieros de los negocios en desarrollo.',
                         'Coordinar la generación y entrega de información financiera de los negocios en desarrollo.',
                         'Elaborar el seguimiento presupuestal de los negocios en desarrollo.',
                         'Ejecutar los controles de normatividad y control interno aplicables a los negocios en desarrollo.']},
            {'title': 'Gerente Área BP\nFinanzas Plataformas Expansión',
             'bullets': ['Elaborar los modelos de evaluación financiera de proyectos de expansión.',
                         'Coordinar el seguimiento presupuestal de los proyectos de expansión en curso.',
                         'Generar reportes financieros periódicos de la cartera de expansión.',
                         'Implementar los controles administrativos y financieros para los procesos de expansión.']},
            {'title': 'Gerente Gestión\nControl Zona',
             'bullets': ['Ejecutar el plan de normatividad y control interno en las plazas de la zona asignada.',
                         'Implementar y dar seguimiento a los procedimientos contables y administrativos en tiendas de la zona.',
                         'Generar y reportar los indicadores de gestión y control de la zona.',
                         'Coordinar con los Gerentes de Plaza y Región la implementación de acciones correctivas.']},
            {'title': 'Responsable\nGestión y Control',
             'bullets': ['Ejecutar los procesos de control administrativo y financiero en las plazas asignadas.',
                         'Generar y consolidar los reportes de gestión e indicadores operativos de las plazas.',
                         'Coordinar la ejecución de revisiones de control interno en plazas y tiendas.',
                         'Implementar y dar seguimiento a las herramientas de gestión y KPIs en las plazas.']},
            {'title': 'Responsable\nTraslado de Valores',
             'bullets': ['Coordinar la operación del modelo de traslado y resguardo de valores en las tiendas.',
                         'Implementar y dar seguimiento al modelo I-Cash y su caso de negocio.',
                         'Coordinar la relación con los proveedores de servicios de traslado de valores.',
                         'Elaborar reportes de desempeño y costos del modelo de traslado de valores.']},
        ]
    },
    {
        'subtitle': 'Administración Nacional CdS',
        'l2_title': 'Gte Sr Nac. Administrativo\nCentros de Suministro',
        'l3s': [
            {'title': 'Gerente Área\nAdmin Control CEDIS',
             'bullets': ['Implementar y mantener los procedimientos contables y administrativos en CEDIS.',
                         'Ejecutar el plan anual de normatividad y control interno en CEDIS.',
                         'Coordinar la generación de información financiera y operativa de CEDIS.',
                         'Elaborar y dar seguimiento a los indicadores de productividad de CEDIS.']},
            {'title': 'Gerente Área Admin\ny Finanzas Dist',
             'bullets': ['Implementar y mantener los procedimientos contables y administrativos en Distribución.',
                         'Coordinar la generación de información financiera de rentabilidad en Distribución.',
                         'Elaborar y dar seguimiento a los indicadores de productividad de Distribución.',
                         'Ejecutar el plan de normatividad y control interno en Distribución.']},
            {'title': 'Responsable Gestión\ny Control CEDIS',
             'bullets': ['Ejecutar los procesos de control administrativo en las unidades de CEDIS asignadas.',
                         'Generar y consolidar los reportes de gestión e indicadores de productividad de CEDIS.',
                         'Coordinar la ejecución de revisiones de control interno en CEDIS.',
                         'Implementar los mecanismos de seguimiento al capital de trabajo en CEDIS.']},
            {'title': 'Responsable\nGestoría',
             'bullets': ['Ejecutar y dar seguimiento a los trámites y gestiones administrativas de CEDIS y Distribución.',
                         'Coordinar la gestión documental y archivo de expedientes normativos y regulatorios.',
                         'Identificar y reportar vencimientos y renovaciones de permisos operativos de CEDIS y Distribución.',
                         'Coordinar con el área Legal los requerimientos de gestoría que impliquen procesos jurídicos.']},
        ]
    },
]

# L2 titles and descriptions for the L1→L2 overview slide
L1_L2_OVERVIEW = {
    'l1_title': 'Director de Administración\ny Finanzas (DAF)',
    'l2s': [
        ('Gte Sr Planeación\nFinanciera & CF',
         'Copiloto financiero del negocio. Planeación, presupuesto y análisis de rendimiento.'),
        ('Gte Sr Fiscal\ny Contraloría',
         'Cumplimiento fiscal, control interno, contraloría y gestión de riesgos de SI.'),
        ('Gte Sr Servicios\nCompartidos Admvo',
         'Operación centralizada de P2P, O2C y R2R. Cierres contables y gobierno de datos.'),
        ('Gte Sr Legal\ny Regulatorio',
         'Gestión legal, cumplimiento regulatorio, litigios y contratos a nivel nacional.'),
        ('Gte Sr Asuntos\nCorporativos',
         'Relacionamiento con gobierno, comunicación externa y gestión de riesgos políticos.'),
        ('Director Nacional\nAdmvo OXXO',
         'Control administrativo en tiendas, plazas y negocios en desarrollo.'),
        ('Gte Sr Nac. Admvo\nCentros de Suministro',
         'Control administrativo en CEDIS y Distribución. Gestoría y normatividad.'),
    ]
}


# ── XML helpers ───────────────────────────────────────────────
def clear_slide(slide):
    spTree = slide.shapes._spTree
    keep = {f'{{{NS_P}}}nvGrpSpPr', f'{{{NS_P}}}grpSpPr'}
    for child in list(spTree):
        if child.tag not in keep:
            spTree.remove(child)


def set_fill_xml(spPr, hex_color=None, no_fill=False):
    for tag in [f'{{{NS_A}}}noFill', f'{{{NS_A}}}solidFill']:
        for el in spPr.findall(tag):
            spPr.remove(el)
    if no_fill:
        etree.SubElement(spPr, f'{{{NS_A}}}noFill')
    elif hex_color:
        sf = etree.SubElement(spPr, f'{{{NS_A}}}solidFill')
        etree.SubElement(sf, f'{{{NS_A}}}srgbClr').set('val', hex_color)


def set_border_xml(spPr, hex_color='000000', width_pt=0.75):
    ln = spPr.find(f'{{{NS_A}}}ln')
    if ln is None:
        ln = etree.SubElement(spPr, f'{{{NS_A}}}ln')
    ln.set('w', str(int(width_pt * 12700)))
    for child in list(ln):
        ln.remove(child)
    sf = etree.SubElement(ln, f'{{{NS_A}}}solidFill')
    etree.SubElement(sf, f'{{{NS_A}}}srgbClr').set('val', hex_color)


def set_no_border_xml(spPr):
    ln = spPr.find(f'{{{NS_A}}}ln')
    if ln is None:
        ln = etree.SubElement(spPr, f'{{{NS_A}}}ln')
    for child in list(ln):
        ln.remove(child)
    etree.SubElement(ln, f'{{{NS_A}}}noFill')


def make_sp(spTree, shape_id, name, left, top, width, height,
            fill_hex=None, no_fill=False, border_hex=None, no_border=False,
            border_pt=0.75, wrap='square', lIns=45720, rIns=45720, tIns=30000, bIns=30000):
    sp = etree.SubElement(spTree, f'{{{NS_P}}}sp')
    nvSpPr = etree.SubElement(sp, f'{{{NS_P}}}nvSpPr')
    cNvPr = etree.SubElement(nvSpPr, f'{{{NS_P}}}cNvPr')
    cNvPr.set('id', str(shape_id)); cNvPr.set('name', name)
    etree.SubElement(nvSpPr, f'{{{NS_P}}}cNvSpPr').set('txBox', '1')
    etree.SubElement(nvSpPr, f'{{{NS_P}}}nvPr')
    spPr = etree.SubElement(sp, f'{{{NS_P}}}spPr')
    xfrm = etree.SubElement(spPr, f'{{{NS_A}}}xfrm')
    off = etree.SubElement(xfrm, f'{{{NS_A}}}off')
    off.set('x', str(int(left))); off.set('y', str(int(top)))
    ext = etree.SubElement(xfrm, f'{{{NS_A}}}ext')
    ext.set('cx', str(int(width))); ext.set('cy', str(int(height)))
    pg = etree.SubElement(spPr, f'{{{NS_A}}}prstGeom'); pg.set('prst', 'rect')
    etree.SubElement(pg, f'{{{NS_A}}}avLst')
    set_fill_xml(spPr, hex_color=fill_hex, no_fill=no_fill)
    if no_border:
        set_no_border_xml(spPr)
    elif border_hex:
        set_border_xml(spPr, hex_color=border_hex, width_pt=border_pt)
    txBody = etree.SubElement(sp, f'{{{NS_P}}}txBody')
    bodyPr = etree.SubElement(txBody, f'{{{NS_A}}}bodyPr')
    bodyPr.set('wrap', wrap)
    bodyPr.set('lIns', str(lIns)); bodyPr.set('rIns', str(rIns))
    bodyPr.set('tIns', str(tIns)); bodyPr.set('bIns', str(bIns))
    etree.SubElement(txBody, f'{{{NS_A}}}lstStyle')
    return txBody


def add_para(txBody, text, size_pt, bold=False, italic=False, color='000000',
             align='l', spc_before_pt=0):
    para = etree.SubElement(txBody, f'{{{NS_A}}}p')
    pPr = etree.SubElement(para, f'{{{NS_A}}}pPr')
    pPr.set('algn', align)
    if spc_before_pt:
        spcBef = etree.SubElement(pPr, f'{{{NS_A}}}spcBef')
        etree.SubElement(spcBef, f'{{{NS_A}}}spcPts').set('val', str(int(spc_before_pt * 100)))
    r = etree.SubElement(para, f'{{{NS_A}}}r')
    rPr = etree.SubElement(r, f'{{{NS_A}}}rPr')
    rPr.set('sz', str(int(size_pt * 100)))
    rPr.set('b', '1' if bold else '0')
    rPr.set('i', '1' if italic else '0')
    sf = etree.SubElement(rPr, f'{{{NS_A}}}solidFill')
    etree.SubElement(sf, f'{{{NS_A}}}srgbClr').set('val', color)
    etree.SubElement(r, f'{{{NS_A}}}t').text = text
    return para


def add_title_ph(spTree, shape_id, text, font_pt=22):
    sp = etree.SubElement(spTree, f'{{{NS_P}}}sp')
    nvSpPr = etree.SubElement(sp, f'{{{NS_P}}}nvSpPr')
    cNvPr = etree.SubElement(nvSpPr, f'{{{NS_P}}}cNvPr')
    cNvPr.set('id', str(shape_id)); cNvPr.set('name', 'Title 1')
    cNvSpPr = etree.SubElement(nvSpPr, f'{{{NS_P}}}cNvSpPr')
    etree.SubElement(cNvSpPr, f'{{{NS_A}}}spLocks').set('noGrp', '1')
    nvPr = etree.SubElement(nvSpPr, f'{{{NS_P}}}nvPr')
    etree.SubElement(nvPr, f'{{{NS_P}}}ph').set('type', 'title')
    etree.SubElement(sp, f'{{{NS_P}}}spPr')
    txBody = etree.SubElement(sp, f'{{{NS_P}}}txBody')
    etree.SubElement(txBody, f'{{{NS_A}}}bodyPr')
    etree.SubElement(txBody, f'{{{NS_A}}}lstStyle')
    para = etree.SubElement(txBody, f'{{{NS_A}}}p')
    r = etree.SubElement(para, f'{{{NS_A}}}r')
    rPr = etree.SubElement(r, f'{{{NS_A}}}rPr')
    rPr.set('b', '1'); rPr.set('sz', str(int(font_pt * 100)))
    sf = etree.SubElement(rPr, f'{{{NS_A}}}solidFill')
    etree.SubElement(sf, f'{{{NS_A}}}srgbClr').set('val', '000000')
    etree.SubElement(r, f'{{{NS_A}}}t').text = text


def add_line(spTree, shape_id, x, y, length, vertical=False):
    if vertical:
        make_sp(spTree, shape_id, f'Line{shape_id}', x - 6350, y, 12700, length,
                fill_hex='000000', no_border=True, wrap='none')
    else:
        make_sp(spTree, shape_id, f'Line{shape_id}', x, y - 6350, length, 12700,
                fill_hex='000000', no_border=True, wrap='none')


# ── Layout constants for org chart slides ─────────────────────
# Title placeholder is at top=250442 per layout (about 380000 tall at 22pt)
TITLE_BOTTOM  = 640000   # safe estimate: title ends here
SUBTITLE_TOP  = 660000   # subtitle starts just below
SUBTITLE_H    = 220000
L2_TOP        = 930000
L2_H          = 500000
L2_W          = 3200000

# Standard description box height (compact, same for all slides)
# Computed to fit 4 bullets at ~2 lines each at 9-10pt
DESC_H_STD    = 2200000  # ~2.4 inches — fits all content, no excess white space

L3_TOP        = 1730000
L3_H          = 560000
DESC_TOP      = L3_TOP + L3_H + 50000  # = 2340000

CONN_H_TOP   = L2_TOP + L2_H           # where connector starts
HBAR_Y       = L3_TOP - 100000


def build_org_chart_slide(slide, l2_info, start_id=200):
    spTree = slide.shapes._spTree
    clear_slide(slide)
    sid = start_id
    N = len(l2_info['l3s'])

    # Title placeholder
    add_title_ph(spTree, sid, 'Estructura Organizacional', font_pt=22); sid += 1

    # Subtitle (red italic) — positioned BELOW the title, no overlap
    txb = make_sp(spTree, sid, 'Subtitle', 330200, SUBTITLE_TOP, 10000000, SUBTITLE_H,
                  no_fill=True, no_border=True, tIns=0, bIns=0); sid += 1
    add_para(txb, l2_info['subtitle'], 18, italic=True, color='CC0000', align='l')

    # L2 box (centered)
    L2_L = (SLIDE_W - L2_W) // 2
    txb = make_sp(spTree, sid, 'L2Box', L2_L, L2_TOP, L2_W, L2_H,
                  fill_hex='FFFFFF', border_hex='000000', border_pt=1,
                  tIns=35000, bIns=35000); sid += 1
    first = True
    for line in l2_info['l2_title'].split('\n'):
        add_para(txb, line, 10, bold=True, color='000000', align='ctr',
                 spc_before_pt=0 if first else 2)
        first = False

    # Connectors
    L2_CX = L2_L + L2_W // 2
    VERT_LEN = HBAR_Y - CONN_H_TOP
    add_line(spTree, sid, L2_CX, CONN_H_TOP, VERT_LEN, vertical=True); sid += 1

    GAP  = 50000 if N <= 4 else 35000
    L3_W = (SLIDE_W - 2 * SIDE_M - (N - 1) * GAP) // N
    L3_LEFTS = [SIDE_M + i * (L3_W + GAP) for i in range(N)]
    L3_CXS   = [L + L3_W // 2 for L in L3_LEFTS]

    if N > 1:
        add_line(spTree, sid, L3_CXS[0], HBAR_Y, L3_CXS[-1] - L3_CXS[0], vertical=False); sid += 1

    for cx in L3_CXS:
        add_line(spTree, sid, cx, HBAR_Y, L3_TOP - HBAR_Y, vertical=True); sid += 1

    # L3 title boxes and description boxes
    title_pt  = 10 if N <= 4 else 9
    bullet_pt = 9  if N <= 4 else 8

    for i, l3 in enumerate(l2_info['l3s']):
        L3_L = L3_LEFTS[i]

        # L3 header box
        txb = make_sp(spTree, sid, f'L3_{i}', L3_L, L3_TOP, L3_W, L3_H,
                      fill_hex='FFFFFF', border_hex='1F3864', border_pt=1,
                      tIns=25000, bIns=25000); sid += 1
        first = True
        for line in l3['title'].split('\n'):
            add_para(txb, line, title_pt, bold=True, color='1F3864', align='ctr',
                     spc_before_pt=0 if first else 1)
            first = False

        # Description box
        txb = make_sp(spTree, sid, f'Desc_{i}', L3_L, DESC_TOP, L3_W, DESC_H_STD,
                      fill_hex='F2F6FC', border_hex='BDD0E9', border_pt=0.75); sid += 1
        first = True
        for bullet in l3['bullets']:
            add_para(txb, '• ' + bullet, bullet_pt, color='1F3864',
                     align='l', spc_before_pt=0 if first else 5)
            first = False

    print(f'  Built L2→L3: {l2_info["subtitle"]} ({N} L3s)')


# ── Build L1→L2 overview slide ────────────────────────────────
def build_l1_l2_overview(slide):
    spTree = slide.shapes._spTree
    clear_slide(slide)
    sid = 100
    ov = L1_L2_OVERVIEW
    N = len(ov['l2s'])

    add_title_ph(spTree, sid, 'Estructura Organizacional', font_pt=22); sid += 1

    # Subtitle
    txb = make_sp(spTree, sid, 'Subtitle', 330200, SUBTITLE_TOP, 10000000, SUBTITLE_H,
                  no_fill=True, no_border=True, tIns=0, bIns=0); sid += 1
    add_para(txb, 'Finanzas', 18, italic=True, color='CC0000', align='l')

    # L1 box (DAF)
    L1_W = 3400000; L1_H = 500000
    L1_L = (SLIDE_W - L1_W) // 2
    txb = make_sp(spTree, sid, 'L1Box', L1_L, L2_TOP, L1_W, L1_H,
                  fill_hex='1F3864', border_hex='1F3864', border_pt=1,
                  tIns=35000, bIns=35000); sid += 1
    first = True
    for line in ov['l1_title'].split('\n'):
        add_para(txb, line, 11, bold=True, color='FFFFFF', align='ctr',
                 spc_before_pt=0 if first else 2)
        first = False

    # Connectors
    L1_CX    = L1_L + L1_W // 2
    HBAR_Y_  = L3_TOP - 100000
    add_line(spTree, sid, L1_CX, L2_TOP + L1_H, HBAR_Y_ - (L2_TOP + L1_H), vertical=True); sid += 1

    GAP  = 35000
    L2_W_ = (SLIDE_W - 2 * SIDE_M - (N - 1) * GAP) // N
    L2_LEFTS = [SIDE_M + i * (L2_W_ + GAP) for i in range(N)]
    L2_CXS   = [L + L2_W_ // 2 for L in L2_LEFTS]

    if N > 1:
        add_line(spTree, sid, L2_CXS[0], HBAR_Y_, L2_CXS[-1] - L2_CXS[0], vertical=False); sid += 1

    for cx in L2_CXS:
        add_line(spTree, sid, cx, HBAR_Y_, L3_TOP - HBAR_Y_, vertical=True); sid += 1

    # L2 boxes (with short description below title)
    for i, (l2_title, l2_desc) in enumerate(ov['l2s']):
        L2_L_ = L2_LEFTS[i]
        txb = make_sp(spTree, sid, f'L2_{i}', L2_L_, L3_TOP, L2_W_, L3_H,
                      fill_hex='1F3864', border_hex='1F3864', border_pt=1,
                      tIns=20000, bIns=20000); sid += 1
        first = True
        for line in l2_title.split('\n'):
            add_para(txb, line, 9, bold=True, color='FFFFFF', align='ctr',
                     spc_before_pt=0 if first else 1)
            first = False

        # Description box below each L2 box
        txb = make_sp(spTree, sid, f'L2Desc_{i}', L2_L_, DESC_TOP, L2_W_, DESC_H_STD,
                      fill_hex='EEF3FB', border_hex='BDD0E9', border_pt=0.75); sid += 1
        add_para(txb, l2_desc, 8, color='1F3864', align='l')

    print('  Built L1→L2 overview')


# ── Hide a slide (show=0 so it doesn't appear in presentation) ──
def hide_slide(slide):
    slide._element.set('show', '0')


# ── Reorder sldIdLst ──────────────────────────────────────────
def reorder_slides(prs):
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    # Current: 0-4 intro, 5-11 good L2→L3, 12 RASCI header, 13-58 RASCI detail,
    #          59 repurposed as L1→L2 overview, 60-65 old bad slides, 66-67 Anexos/Gracias
    desired = (
        ids[0:5]    +  # intro slides
        [ids[59]]   +  # new L1→L2 overview (will be slide 6)
        [ids[12]]   +  # RASCI methodology (slide 7)
        ids[13:59]  +  # 46 RASCI detail slides (slides 8-53)
        ids[5:12]   +  # 7 L2→L3 org charts (slides 54-60)
        ids[60:66]  +  # old bad slides cleared (61-66)
        ids[66:68]     # Anexos, Gracias (67-68)
    )
    assert len(desired) == 68, f"Expected 68, got {len(desired)}"
    for sldId in list(sldIdLst):
        sldIdLst.remove(sldId)
    for sldId in desired:
        sldIdLst.append(sldId)
    print('  Reordered sldIdLst (68 slides)')


# ── Main ──────────────────────────────────────────────────────
def main():
    print(f'Loading: {PPTX_PATH}')
    prs = Presentation(PPTX_PATH)
    print(f'  Slides: {len(prs.slides)}')

    # 1. Build L1→L2 overview into slide at idx 59 (old PF&A)
    print('Building L1→L2 overview (repurposing idx 59)...')
    build_l1_l2_overview(prs.slides[59])

    # 2. Rebuild the 7 L2→L3 org chart slides (idx 5-11) with fixed visuals
    print('Rebuilding 7 L2→L3 org chart slides with fixed visuals...')
    for i, l2_info in enumerate(L2_DATA):
        build_org_chart_slide(prs.slides[5 + i], l2_info, start_id=200 + i * 50)

    # 3. Clear + hide old bad slides (idx 60-65)
    print('Clearing old bad slides (idx 60-65)...')
    for i in range(60, 66):
        clear_slide(prs.slides[i])
        hide_slide(prs.slides[i])
    print('  Done')

    # 4. Reorder
    print('Reordering slides...')
    reorder_slides(prs)

    # 5. Save
    prs.save(PPTX_PATH)
    print(f'Saved: {PPTX_PATH}')

    # 6. Verify ZIP
    with zipfile.ZipFile(PPTX_PATH, 'r') as z:
        names = z.namelist()
        dupes = [n for n in names if names.count(n) > 1]
        ct_slides = set(re.findall(r'PartName="/ppt/slides/(slide\d+\.xml)"',
                                    z.read('[Content_Types].xml').decode()))
        zip_slides = set(n.replace('ppt/slides/', '') for n in names
                         if n.startswith('ppt/slides/slide') and '.rels' not in n)
        print(f'ZIP dupes: {set(dupes)}  |  Orphans: {zip_slides - ct_slides}')

    # 7. Verify structure
    prs2 = Presentation(PPTX_PATH)
    print(f'Final slides: {len(prs2.slides)}')
    for i in [0, 4, 5, 6, 7, 52, 53, 59, 60, 65, 66, 67]:
        if i >= len(prs2.slides): continue
        s = prs2.slides[i]
        title = next((sh.text_frame.text[:55] for sh in s.shapes
                      if 'Title' in sh.name and hasattr(sh, 'text_frame')), '?')
        print(f'  Slide {i+1}: {title!r}')


if __name__ == '__main__':
    main()
