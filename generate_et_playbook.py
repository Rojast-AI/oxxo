#!/usr/bin/env python3
"""Generate Estrategia y Transformación Playbook PPTX from Finanzas Playbook base."""

import shutil, os, zipfile, re
from lxml import etree
from pptx import Presentation

SRC = '/home/user/oxxo/Finanzas/Playbook/Finanzas Playbook.pptx'
OUT = '/home/user/oxxo/Estrategia y Transformacion/Playbook/ET Playbook.pptx'

NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
SLIDE_W = 12192000; SLIDE_H = 6858000
MARGIN = 365760; CONTENT_W = SLIDE_W - 2 * MARGIN
SIDE_M = 304800

BADGE_COLORS = {
    'R':'70AD47','A':'FFC000','A/R':'FFC000','S':'FF7070','C':'4472C4','I':'A6A6A6',
}
ROLE_DESCS = {
    'R':  'Ejecuta y lidera esta responsabilidad.',
    'A':  'Aprueba y supervisa el resultado final.',
    'A/R':'Ejecuta y rinde cuentas del resultado final.',
    'S':  'Brinda apoyo y soporte activo.',
    'C':  'Aporta perspectiva experta como consultor.',
}

# Org-chart layout
SUBTITLE_TOP=660000; SUBTITLE_H=220000
L2_TOP=930000; L2_H=500000; L2_W=3200000
L3_TOP=1730000; L3_H=560000
DESC_TOP=L3_TOP+L3_H+50000; DESC_H_STD=2200000
HBAR_Y=L3_TOP-100000

# ── ET DATA ───────────────────────────────────────────────────

ET_DEPT = 'Estrategia y Transformación'
ET_DEPT_SHORT = 'ET'
ET_ROLE_SHORT = [
    'Dir. Estrat.','Gte Sr Estrat.','Gte Val. Delivery',
    'Resp. Proc. y Met.','Dir. Com. Srvs.'
]
ET_ROLE_FULL = [
    'Director Estrategia y Transformación OXXO Mx',
    'Gerente Sr Estrategia',
    'Gerente Value Delivery',
    'Responsable Procesos y Metodología',
    'Director Comercial de Servicios',
]

# (section, name, description, roles_dict {0..4: letter}, explanation)
ET_RASCI = [
    # ── Sección 1: Estrategia ──
    ('Director Estrategia y Transformación','Definir la visión y hoja de ruta estratégica de OXXO Mx',
     'Establece la dirección estratégica de la organización, definiendo prioridades y la hoja de ruta de transformación.',
     {0:'A/R',1:'S',2:'I',3:'I',4:'I'},
     'El Director lidera y rinde cuentas de la visión estratégica (A/R). El Gte. Sr. Estrategia apoya (S) con análisis e insumos.'),
    ('','Alinear prioridades estratégicas con los objetivos de negocio',
     'Garantiza que las iniciativas estratégicas estén alineadas con los objetivos corporativos y las necesidades del negocio.',
     {0:'A/R',1:'S',2:'C',3:'I',4:'C'},
     'El Director rinde cuentas de la alineación (A/R). El Gte. Sr. Estrategia apoya (S); Value Delivery y Com. Servicios consultan (C).'),
    ('','Evaluar y validar nuevas oportunidades de negocio',
     'Analiza y valida oportunidades de nuevos negocios, mercados o modelos que incrementen el valor de la empresa.',
     {0:'A',1:'R',2:'C',3:'I',4:'C'},
     'El Gte. Sr. Estrategia ejecuta el análisis (R). El Director aprueba (A). Value Delivery y Com. Servicios aportan perspectiva (C).'),
    ('','Gestionar el plan de largo plazo y escenarios estratégicos',
     'Administra el plan estratégico de largo plazo y construye escenarios para la toma de decisiones.',
     {0:'A',1:'R',2:'I',3:'I',4:'I'},
     'El Gte. Sr. Estrategia ejecuta la planificación (R). El Director aprueba y valida los escenarios (A).'),
    ('','Comunicar la estrategia a los diferentes niveles de la organización',
     'Asegura que la estrategia sea comprendida y adoptada en todos los niveles de la organización.',
     {0:'A/R',1:'S',2:'S',3:'I',4:'S'},
     'El Director lidera la comunicación estratégica (A/R). Gte. Estrategia, Value Delivery y Com. Servicios apoyan (S) en la cascada.'),
    ('','Monitorear indicadores clave de desempeño estratégico',
     'Supervisa los KPIs estratégicos para asegurar el avance hacia los objetivos definidos.',
     {0:'A',1:'R',2:'S',3:'I',4:'I'},
     'El Gte. Sr. Estrategia reporta los indicadores (R). El Director valida y actúa sobre las brechas (A). Value Delivery apoya (S).'),

    # ── Sección 2: Value Delivery ──
    ('Gerente Value Delivery','Definir y ejecutar el modelo de entrega de valor',
     'Establece el modelo operativo para garantizar que las iniciativas estratégicas generen valor medible.',
     {0:'A',1:'C',2:'R',3:'S',4:'I'},
     'El Gte. Value Delivery ejecuta el modelo (R). El Director aprueba (A). Metodología apoya en procesos (S). Estrategia consulta (C).'),
    ('','Identificar y eliminar barreras al valor en iniciativas en curso',
     'Detecta obstáculos que impiden la generación de valor en proyectos activos y los resuelve.',
     {0:'I',1:'C',2:'R',3:'S',4:'I'},
     'El Gte. Value Delivery lidera la remoción de obstáculos (R). Metodología apoya (S) con frameworks. Estrategia asesora (C).'),
    ('','Asegurar la entrega ágil y orientada a resultados',
     'Implementa prácticas ágiles y mecanismos de seguimiento para garantizar la entrega oportuna de resultados.',
     {0:'I',1:'I',2:'R',3:'A/R',4:'I'},
     'El Gte. Value Delivery dirige la ejecución ágil (R). Procesos y Metodología rinde cuentas de los marcos de trabajo (A/R).'),
    ('','Gestionar el portafolio de iniciativas estratégicas',
     'Administra el portafolio de proyectos estratégicos, priorizando recursos y asegurando avance.',
     {0:'A',1:'S',2:'R',3:'S',4:'I'},
     'El Gte. Value Delivery gestiona el portafolio (R). El Director aprueba prioridades (A). Estrategia y Metodología apoyan (S).'),
    ('','Medir y reportar el impacto de las iniciativas',
     'Define métricas de impacto y reporta el valor generado por las iniciativas estratégicas.',
     {0:'A',1:'S',2:'R',3:'I',4:'I'},
     'El Gte. Value Delivery reporta el impacto (R). El Director valida y aprueba los resultados (A). Estrategia apoya (S).'),
    ('','Facilitar la toma de decisiones basada en datos de valor',
     'Provee análisis y datos de impacto para apoyar la toma de decisiones de los líderes.',
     {0:'I',1:'C',2:'R',3:'S',4:'I'},
     'El Gte. Value Delivery facilita los análisis de decisión (R). Metodología apoya (S) con estructuras. Estrategia consulta (C).'),

    # ── Sección 3: Procesos y Metodología ──
    ('Responsable Procesos y Metodología','Diseñar y estandarizar metodologías de gestión de proyectos',
     'Desarrolla y estandariza los marcos de gestión de proyectos para toda la organización.',
     {0:'I',1:'C',2:'S',3:'R',4:'I'},
     'El Resp. Proc. y Met. diseña las metodologías (R). Value Delivery apoya (S) validando uso. Estrategia asesora (C).'),
    ('','Documentar y mantener procesos organizacionales',
     'Mapea, documenta y actualiza los procesos clave de la organización para garantizar su claridad y vigencia.',
     {0:'I',1:'I',2:'I',3:'R',4:'I'},
     'El Resp. Proc. y Met. es el único responsable de la documentación y gestión de procesos (R).'),
    ('','Auditar el cumplimiento de procesos y metodologías',
     'Verifica que las áreas cumplan con los procesos y metodologías definidos.',
     {0:'I',1:'I',2:'S',3:'R',4:'I'},
     'El Resp. Proc. y Met. ejecuta las auditorías (R). Value Delivery apoya (S) en la revisión de iniciativas.'),
    ('','Capacitar a los equipos en metodologías y herramientas',
     'Diseña e imparte programas de capacitación en metodologías de gestión y mejora continua.',
     {0:'I',1:'I',2:'S',3:'R',4:'I'},
     'El Resp. Proc. y Met. lidera la capacitación (R). Value Delivery apoya (S) identificando brechas en los equipos.'),
    ('','Proponer mejoras al modelo operativo',
     'Identifica oportunidades de mejora en el modelo operativo y propone cambios para incrementar la eficiencia.',
     {0:'A',1:'S',2:'S',3:'R',4:'I'},
     'El Resp. Proc. y Met. propone mejoras (R). El Director aprueba (A). Estrategia y Value Delivery apoyan (S).'),
    ('','Gestionar el cambio organizacional asociado a nuevos procesos',
     'Lidera la gestión del cambio para asegurar la adopción de nuevos procesos en la organización.',
     {0:'A',1:'S',2:'S',3:'R',4:'I'},
     'El Resp. Proc. y Met. conduce la gestión del cambio (R). El Director aprueba (A). Estrategia y Value Delivery apoyan (S).'),

    # ── Sección 4: Comercial Servicios ──
    ('Director Comercial de Servicios','Definir la estrategia comercial de servicios financieros y digitales',
     'Establece la estrategia comercial para los servicios financieros y digitales de la red de tiendas.',
     {0:'A',1:'I',2:'I',3:'I',4:'R'},
     'El Director Comercial de Servicios lidera la estrategia comercial (R). El Director ET aprueba la alineación estratégica (A).'),
    ('','Gestionar las relaciones con socios y alianzas comerciales',
     'Administra las relaciones con socios comerciales, proveedores y alianzas para potenciar la oferta de servicios.',
     {0:'I',1:'I',2:'I',3:'I',4:'R'},
     'El Director Comercial de Servicios gestiona directamente las alianzas y socios (R).'),
    ('','Supervisar la operación de corresponsalías y dispersión de efectivo',
     'Asegura la correcta operación de corresponsalías, dispersión de efectivo y servicios financieros en tienda.',
     {0:'I',1:'I',2:'I',3:'I',4:'R'},
     'El Director Comercial de Servicios rinde cuentas de la operación comercial (R).'),
    ('','Impulsar la adquisición y fidelización de clientes de servicios',
     'Diseña y ejecuta estrategias para captar nuevos clientes de servicios y retener a los existentes.',
     {0:'I',1:'C',2:'S',3:'I',4:'R'},
     'El Director Comercial de Servicios ejecuta la estrategia de adquisición (R). Value Delivery apoya (S). Estrategia asesora (C).'),
    ('','Definir el portafolio de productos y servicios de la red',
     'Determina el portafolio de productos y servicios ofrecidos a través de la red de tiendas.',
     {0:'A',1:'C',2:'I',3:'I',4:'R'},
     'El Director Comercial de Servicios define el portafolio (R). El Director ET aprueba (A). Estrategia asesora (C).'),
    ('','Monitorear KPIs comerciales y de servicios',
     'Supervisa los indicadores de desempeño comercial y de servicios para asegurar el cumplimiento de metas.',
     {0:'I',1:'I',2:'I',3:'I',4:'R'},
     'El Director Comercial de Servicios es el responsable directo del desempeño comercial (R).'),

    # ── Sección 5: Transversal ──
    ('Transversal','Gestionar la agenda de transformación organizacional',
     'Coordina y da seguimiento a los proyectos de transformación que impactan a toda la organización.',
     {0:'A/R',1:'S',2:'S',3:'S',4:'I'},
     'El Director ET es el dueño de la agenda de transformación (A/R). Estrategia, Value Delivery y Metodología apoyan (S).'),
    ('','Alinear iniciativas de transformación con la dirección corporativa',
     'Garantiza que los proyectos de transformación estén alineados con las directrices corporativas.',
     {0:'A/R',1:'S',2:'I',3:'I',4:'I'},
     'El Director ET alinea la transformación con el corporativo (A/R). El Gte. Sr. Estrategia apoya (S) con análisis.'),
    ('','Desarrollar capacidades organizacionales estratégicas',
     'Impulsa el desarrollo de las capacidades y talento necesarios para ejecutar la estrategia de largo plazo.',
     {0:'A/R',1:'S',2:'S',3:'S',4:'I'},
     'El Director ET lidera el desarrollo de capacidades (A/R). Todos los gerentes apoyan (S) en sus respectivos ámbitos.'),
    ('','Gestionar la comunicación e imagen institucional del área',
     'Administra la comunicación interna y externa del área de Estrategia y Transformación.',
     {0:'R',1:'S',2:'I',3:'I',4:'I'},
     'El Director ET gestiona la imagen y comunicación institucional (R). El Gte. Sr. Estrategia apoya (S).'),
    ('','Fomentar la cultura de innovación y mejora continua',
     'Promueve una cultura organizacional orientada a la innovación, el aprendizaje y la mejora continua.',
     {0:'A/R',1:'S',2:'S',3:'S',4:'S'},
     'El Director ET lidera la cultura de innovación (A/R). Todos los gerentes apoyan (S) en su ámbito.'),
    ('','Garantizar el uso eficiente de recursos del área',
     'Administra el presupuesto y recursos del área para maximizar el impacto de las iniciativas estratégicas.',
     {0:'A/R',1:'I',2:'I',3:'I',4:'I'},
     'El Director ET rinde cuentas del uso de recursos del área (A/R).'),
]

ET_L1_L2_OVERVIEW = {
    'l1_title': 'Director Estrategia y\nTransformación OXXO Mx',
    'l2s': [
        ('Gte Sr\nEstrategia','Planeación estratégica de largo plazo y evaluación de oportunidades de negocio.'),
        ('Gte Sr Comercial\nServicios Financieros','Gestión comercial de corresponsalías, dispersión de efectivo y productos financieros.'),
        ('Gte\nValue Delivery','Entrega de valor en iniciativas estratégicas bajo marcos ágiles.'),
        ('Gte Sr\nFullfilment','Operación y gestión del modelo de fulfillment y adquisición de clientes.'),
        ('Gte OXXO Cel','Estrategia y operación comercial de telefonía y chips en red de tiendas.'),
    ]
}

ET_L2_DATA = [
    {
        'subtitle':'OXXO Cel','l2_title':'Gerente\nOXXO Cel',
        'l3s':[
            {'title':'Analista\nDatos','bullets':[
                'Analizar datos de ventas y comportamiento de clientes de telefonía.',
                'Construir reportes y dashboards de desempeño de OXXO Cel.',
                'Identificar tendencias y oportunidades en el mercado de telefonía.',
                'Apoyar la toma de decisiones con análisis cuantitativos.',
            ]},
            {'title':'Analista Operación\nTeléfonos y Chips','bullets':[
                'Gestionar la operación diaria de venta de teléfonos y chips en tienda.',
                'Monitorear inventarios y asegurar disponibilidad de productos.',
                'Coordinar con proveedores la reposición y activación de chips.',
                'Documentar y mejorar los procesos operativos de OXXO Cel.',
            ]},
            {'title':'Coord. Comercial\nFidelización Clientes','bullets':[
                'Diseñar e implementar programas de fidelización para clientes de telefonía.',
                'Gestionar campañas de retención y upgrades de planes.',
                'Analizar el churn y desarrollar acciones para reducirlo.',
                'Coordinar con marketing las activaciones y promociones de telefonía.',
            ]},
            {'title':'Coord. Product Owner\nCX','bullets':[
                'Gestionar el backlog del producto de experiencia de cliente OXXO Cel.',
                'Priorizar funcionalidades y mejoras con base en necesidades del cliente.',
                'Coordinar con TI y operaciones para implementar mejoras de CX.',
                'Validar y aceptar entregables de desarrollos de producto.',
            ]},
            {'title':'Gte Área Adquisición\nCliente','bullets':[
                'Definir y ejecutar la estrategia de adquisición de nuevos clientes.',
                'Gestionar campañas de activación y alta de nuevos servicios.',
                'Analizar el desempeño de canales de adquisición.',
                'Coordinar con el equipo comercial las metas de nuevas altas.',
            ]},
        ]
    },
    {
        'subtitle':'Comercial Servicios Financieros','l2_title':'Gerente Sr Comercial\nServicios Financieros',
        'l3s':[
            {'title':'Gte Área Comercial\nCorresponsalías','bullets':[
                'Gestionar las relaciones comerciales con bancos y entidades corresponsales.',
                'Monitorear el desempeño de transacciones de corresponsalía en tienda.',
                'Desarrollar nuevas alianzas comerciales para ampliar la oferta de servicios.',
                'Asegurar el cumplimiento regulatorio de los servicios de corresponsalía.',
            ]},
            {'title':'Gte Área Comercial\nDispersión Efectivo','bullets':[
                'Gestionar la operación y crecimiento del negocio de dispersión de efectivo.',
                'Desarrollar relaciones con clientes empresariales de dispersión.',
                'Monitorear indicadores de volumen y calidad de la operación.',
                'Identificar nuevas oportunidades de negocio en dispersión.',
            ]},
            {'title':'Gte Área Operaciones\niCash','bullets':[
                'Supervisar la operación técnica y comercial de la plataforma iCash.',
                'Asegurar la disponibilidad y correcto funcionamiento del sistema.',
                'Gestionar incidencias y mejoras de la plataforma.',
                'Coordinar con TI el desarrollo de nuevas funcionalidades.',
            ]},
            {'title':'Gte Área\nProducto','bullets':[
                'Definir la estrategia de producto de servicios financieros en tienda.',
                'Gestionar el ciclo de vida de productos financieros (lanzamiento, crecimiento, retiro).',
                'Analizar el mercado y la competencia para identificar oportunidades.',
                'Coordinar con regulación y cumplimiento el lanzamiento de nuevos productos.',
            ]},
            {'title':'Resp. Comercial\nMedios de Pago','bullets':[
                'Gestionar las relaciones con operadores de medios de pago digitales.',
                'Monitorear el desempeño de transacciones de medios de pago en red.',
                'Desarrollar estrategias para incrementar la adopción de pagos digitales.',
                'Coordinar la habilitación técnica de nuevos medios de pago en POS.',
            ]},
        ]
    },
    {
        'subtitle':'Estrategia','l2_title':'Gerente Sr\nEstrategia',
        'l3s':[
            {'title':'Resp.\nEstrategia','bullets':[
                'Realizar análisis estratégicos y estudios de mercado para la dirección.',
                'Apoyar la elaboración del plan estratégico de largo plazo.',
                'Documentar y comunicar los resultados de los proyectos estratégicos.',
                'Monitorear tendencias del sector retail y servicios.',
            ]},
            {'title':'Resp.\nPlaneación','bullets':[
                'Coordinar el proceso de planeación anual y multianual del área.',
                'Consolidar y dar seguimiento a los planes e iniciativas estratégicas.',
                'Elaborar reportes de avance estratégico para la alta dirección.',
                'Gestionar el calendario y la agenda estratégica del área.',
            ]},
        ]
    },
    {
        'subtitle':'Fullfilment','l2_title':'Gerente Sr\nFullfilment',
        'l3s':[
            {'title':'Coord. Op.\nFullfilment','bullets':[
                'Coordinar la operación diaria del modelo de fulfillment.',
                'Monitorear indicadores operativos y de calidad del servicio.',
                'Gestionar incidencias y asegurar la continuidad operativa.',
                'Documentar y optimizar los procesos de fulfillment.',
            ]},
            {'title':'Coord. Product Owner\nServicios','bullets':[
                'Gestionar el backlog del producto de servicios de fulfillment.',
                'Priorizar funcionalidades con base en necesidades del cliente y del negocio.',
                'Coordinar con TI y operaciones la implementación de mejoras.',
                'Validar y aceptar entregables de desarrollos de producto.',
            ]},
            {'title':'Gte Área Comercial\nFullfilment','bullets':[
                'Desarrollar y gestionar las relaciones comerciales del modelo de fulfillment.',
                'Identificar nuevos clientes y oportunidades de expansión.',
                'Diseñar propuestas de valor y condiciones comerciales.',
                'Monitorear el desempeño comercial e indicadores clave.',
            ]},
            {'title':'Resp. MKT\nAdquisición Cliente','bullets':[
                'Diseñar e implementar estrategias de marketing para adquisición de clientes.',
                'Gestionar campañas digitales y de medios para el modelo de fulfillment.',
                'Medir el ROI de las acciones de marketing y optimizar inversión.',
                'Coordinar con el equipo comercial las activaciones y promociones.',
            ]},
            {'title':'Resp. Op. Dis.\nFullfilment','bullets':[
                'Gestionar la operación de distribución del modelo de fulfillment.',
                'Coordinar con proveedores logísticos la entrega de pedidos.',
                'Monitorear los tiempos de entrega y calidad del servicio.',
                'Optimizar rutas y procesos de distribución para reducir costos.',
            ]},
        ]
    },
    {
        'subtitle':'Value Delivery','l2_title':'Gerente\nValue Delivery',
        'l3s':[
            {'title':'Coord.\nValue Delivery','bullets':[
                'Coordinar el seguimiento de iniciativas estratégicas en el portafolio.',
                'Gestionar la comunicación entre equipos de proyecto y dirección.',
                'Actualizar tableros de seguimiento y reportes de avance.',
                'Identificar riesgos y escalar obstáculos para resolución oportuna.',
            ]},
            {'title':'Gte Área\nValue Delivery','bullets':[
                'Gestionar iniciativas de alto impacto del portafolio estratégico.',
                'Liderar equipos multidisciplinarios en la entrega de valor.',
                'Asegurar la alineación de objetivos de proyecto con la estrategia.',
                'Reportar el impacto y valor generado a la dirección.',
            ]},
            {'title':'Resp. Procesos\ny Metodología','bullets':[
                'Diseñar y estandarizar metodologías de gestión de proyectos del área.',
                'Documentar y mantener los procesos organizacionales del área ET.',
                'Capacitar a los equipos en metodologías y herramientas de gestión.',
                'Auditar el cumplimiento de procesos y metodologías definidos.',
            ]},
        ]
    },
]

ET_AGENDA = [
    '1. Macroestructura OXXO MX',
    '2. Mapa Interacciones Estrategia y Transformación',
    '3. Estructura Organizacional Estrategia y Transformación',
    '4. Matriz RASCI (Metodología)',
    '5. Detalle de la Matriz RASCI por Responsabilidad',
    '6. Estructura Organizacional OXXO Cel',
    '7. Estructura Organizacional Comercial Servicios Financieros',
    '8. Estructura Organizacional Estrategia',
    '9. Estructura Organizacional Fullfilment',
    '10. Estructura Organizacional Value Delivery',
    '11. Anexos Descripciones de Puesto',
]

ET_PROPOSITO = (
    'Diseñar y ejecutar la estrategia de largo plazo de OXXO México, asegurando '
    'la transformación organizacional, la entrega de valor en las iniciativas '
    'estratégicas y el desarrollo de servicios financieros y digitales que '
    'amplíen la propuesta de valor para los clientes y el negocio.'
)

# ── XML helpers ───────────────────────────────────────────────

def clear_slide(slide):
    spTree = slide.shapes._spTree
    keep = {f'{{{NS_P}}}nvGrpSpPr', f'{{{NS_P}}}grpSpPr'}
    for child in list(spTree):
        if child.tag not in keep:
            spTree.remove(child)


def set_fill_xml(spPr, hex_color=None, no_fill=False):
    for tag in [f'{{{NS_A}}}noFill', f'{{{NS_A}}}solidFill']:
        for el in spPr.findall(tag):
            spPr.remove(el)
    if no_fill:
        etree.SubElement(spPr, f'{{{NS_A}}}noFill')
    elif hex_color:
        sf = etree.SubElement(spPr, f'{{{NS_A}}}solidFill')
        etree.SubElement(sf, f'{{{NS_A}}}srgbClr').set('val', hex_color)


def set_border_xml(spPr, hex_color='000000', width_pt=0.75):
    ln = spPr.find(f'{{{NS_A}}}ln')
    if ln is None:
        ln = etree.SubElement(spPr, f'{{{NS_A}}}ln')
    ln.set('w', str(int(width_pt * 12700)))
    for child in list(ln):
        ln.remove(child)
    sf = etree.SubElement(ln, f'{{{NS_A}}}solidFill')
    etree.SubElement(sf, f'{{{NS_A}}}srgbClr').set('val', hex_color)


def set_no_border_xml(spPr):
    ln = spPr.find(f'{{{NS_A}}}ln')
    if ln is None:
        ln = etree.SubElement(spPr, f'{{{NS_A}}}ln')
    for child in list(ln):
        ln.remove(child)
    etree.SubElement(ln, f'{{{NS_A}}}noFill')


def make_sp(spTree, shape_id, name, left, top, width, height,
            fill_hex=None, no_fill=False, border_hex=None, no_border=False,
            border_pt=0.75, wrap='square', lIns=45720, rIns=45720, tIns=30000, bIns=30000):
    sp = etree.SubElement(spTree, f'{{{NS_P}}}sp')
    nvSpPr = etree.SubElement(sp, f'{{{NS_P}}}nvSpPr')
    cNvPr = etree.SubElement(nvSpPr, f'{{{NS_P}}}cNvPr')
    cNvPr.set('id', str(shape_id)); cNvPr.set('name', name)
    etree.SubElement(nvSpPr, f'{{{NS_P}}}cNvSpPr').set('txBox', '1')
    etree.SubElement(nvSpPr, f'{{{NS_P}}}nvPr')
    spPr = etree.SubElement(sp, f'{{{NS_P}}}spPr')
    xfrm = etree.SubElement(spPr, f'{{{NS_A}}}xfrm')
    off = etree.SubElement(xfrm, f'{{{NS_A}}}off')
    off.set('x', str(int(left))); off.set('y', str(int(top)))
    ext = etree.SubElement(xfrm, f'{{{NS_A}}}ext')
    ext.set('cx', str(int(width))); ext.set('cy', str(int(height)))
    pg = etree.SubElement(spPr, f'{{{NS_A}}}prstGeom'); pg.set('prst', 'rect')
    etree.SubElement(pg, f'{{{NS_A}}}avLst')
    set_fill_xml(spPr, hex_color=fill_hex, no_fill=no_fill)
    if no_border:
        set_no_border_xml(spPr)
    elif border_hex:
        set_border_xml(spPr, hex_color=border_hex, width_pt=border_pt)
    txBody = etree.SubElement(sp, f'{{{NS_P}}}txBody')
    bodyPr = etree.SubElement(txBody, f'{{{NS_A}}}bodyPr')
    bodyPr.set('wrap', wrap)
    bodyPr.set('lIns', str(lIns)); bodyPr.set('rIns', str(rIns))
    bodyPr.set('tIns', str(tIns)); bodyPr.set('bIns', str(bIns))
    etree.SubElement(txBody, f'{{{NS_A}}}lstStyle')
    return txBody


def add_para(txBody, text, size_pt, bold=False, italic=False, color='000000',
             align='l', spc_before_pt=0):
    para = etree.SubElement(txBody, f'{{{NS_A}}}p')
    pPr = etree.SubElement(para, f'{{{NS_A}}}pPr')
    pPr.set('algn', align)
    if spc_before_pt:
        spcBef = etree.SubElement(pPr, f'{{{NS_A}}}spcBef')
        etree.SubElement(spcBef, f'{{{NS_A}}}spcPts').set('val', str(int(spc_before_pt * 100)))
    r = etree.SubElement(para, f'{{{NS_A}}}r')
    rPr = etree.SubElement(r, f'{{{NS_A}}}rPr')
    rPr.set('sz', str(int(size_pt * 100)))
    rPr.set('b', '1' if bold else '0')
    rPr.set('i', '1' if italic else '0')
    sf = etree.SubElement(rPr, f'{{{NS_A}}}solidFill')
    etree.SubElement(sf, f'{{{NS_A}}}srgbClr').set('val', color)
    etree.SubElement(r, f'{{{NS_A}}}t').text = text
    return para


def add_title_ph(spTree, shape_id, text, font_pt=22):
    sp = etree.SubElement(spTree, f'{{{NS_P}}}sp')
    nvSpPr = etree.SubElement(sp, f'{{{NS_P}}}nvSpPr')
    cNvPr = etree.SubElement(nvSpPr, f'{{{NS_P}}}cNvPr')
    cNvPr.set('id', str(shape_id)); cNvPr.set('name', 'Title 1')
    cNvSpPr = etree.SubElement(nvSpPr, f'{{{NS_P}}}cNvSpPr')
    etree.SubElement(cNvSpPr, f'{{{NS_A}}}spLocks').set('noGrp', '1')
    nvPr = etree.SubElement(nvSpPr, f'{{{NS_P}}}nvPr')
    etree.SubElement(nvPr, f'{{{NS_P}}}ph').set('type', 'title')
    etree.SubElement(sp, f'{{{NS_P}}}spPr')
    txBody = etree.SubElement(sp, f'{{{NS_P}}}txBody')
    etree.SubElement(txBody, f'{{{NS_A}}}bodyPr')
    etree.SubElement(txBody, f'{{{NS_A}}}lstStyle')
    para = etree.SubElement(txBody, f'{{{NS_A}}}p')
    r = etree.SubElement(para, f'{{{NS_A}}}r')
    rPr = etree.SubElement(r, f'{{{NS_A}}}rPr')
    rPr.set('b', '1'); rPr.set('sz', str(int(font_pt * 100)))
    sf = etree.SubElement(rPr, f'{{{NS_A}}}solidFill')
    etree.SubElement(sf, f'{{{NS_A}}}srgbClr').set('val', '000000')
    etree.SubElement(r, f'{{{NS_A}}}t').text = text


def add_line(spTree, shape_id, x, y, length, vertical=False):
    if vertical:
        make_sp(spTree, shape_id, f'Line{shape_id}', x-6350, y, 12700, length,
                fill_hex='000000', no_border=True, wrap='none')
    else:
        make_sp(spTree, shape_id, f'Line{shape_id}', x, y-6350, length, 12700,
                fill_hex='000000', no_border=True, wrap='none')


def hide_slide(slide):
    slide._element.set('show', '0')


# ── Org chart slide builder ───────────────────────────────────

def build_org_chart_slide(slide, l2_info, start_id=200):
    spTree = slide.shapes._spTree
    clear_slide(slide)
    sid = start_id
    N = len(l2_info['l3s'])

    add_title_ph(spTree, sid, 'Estructura Organizacional', font_pt=22); sid += 1

    txb = make_sp(spTree, sid, 'Subtitle', 330200, SUBTITLE_TOP, 10000000, SUBTITLE_H,
                  no_fill=True, no_border=True, tIns=0, bIns=0); sid += 1
    add_para(txb, l2_info['subtitle'], 18, italic=True, color='CC0000', align='l')

    L2_L = (SLIDE_W - L2_W) // 2
    txb = make_sp(spTree, sid, 'L2Box', L2_L, L2_TOP, L2_W, L2_H,
                  fill_hex='FFFFFF', border_hex='000000', border_pt=1,
                  tIns=35000, bIns=35000); sid += 1
    first = True
    for line in l2_info['l2_title'].split('\n'):
        add_para(txb, line, 10, bold=True, color='000000', align='ctr',
                 spc_before_pt=0 if first else 2)
        first = False

    L2_CX = L2_L + L2_W // 2
    VERT_LEN = HBAR_Y - (L2_TOP + L2_H)
    add_line(spTree, sid, L2_CX, L2_TOP + L2_H, VERT_LEN, vertical=True); sid += 1

    GAP = 50000 if N <= 4 else 35000
    L3_W = (SLIDE_W - 2 * SIDE_M - (N - 1) * GAP) // N
    L3_LEFTS = [SIDE_M + i * (L3_W + GAP) for i in range(N)]
    L3_CXS = [L + L3_W // 2 for L in L3_LEFTS]

    if N > 1:
        add_line(spTree, sid, L3_CXS[0], HBAR_Y, L3_CXS[-1] - L3_CXS[0], vertical=False); sid += 1
    for cx in L3_CXS:
        add_line(spTree, sid, cx, HBAR_Y, L3_TOP - HBAR_Y, vertical=True); sid += 1

    title_pt = 10 if N <= 4 else 9
    bullet_pt = 9 if N <= 4 else 8

    for i, l3 in enumerate(l2_info['l3s']):
        L3_L = L3_LEFTS[i]
        txb = make_sp(spTree, sid, f'L3_{i}', L3_L, L3_TOP, L3_W, L3_H,
                      fill_hex='FFFFFF', border_hex='1F3864', border_pt=1,
                      tIns=25000, bIns=25000); sid += 1
        first = True
        for line in l3['title'].split('\n'):
            add_para(txb, line, title_pt, bold=True, color='1F3864', align='ctr',
                     spc_before_pt=0 if first else 1)
            first = False
        txb = make_sp(spTree, sid, f'Desc_{i}', L3_L, DESC_TOP, L3_W, DESC_H_STD,
                      fill_hex='F2F6FC', border_hex='BDD0E9', border_pt=0.75); sid += 1
        first = True
        for bullet in l3['bullets']:
            add_para(txb, '• ' + bullet, bullet_pt, color='1F3864',
                     align='l', spc_before_pt=0 if first else 5)
            first = False

    print(f'  Built org chart: {l2_info["subtitle"]} ({N} L3s)')


# ── L1→L2 overview slide builder ─────────────────────────────

def build_l1_l2_overview(slide, overview, dept_name):
    spTree = slide.shapes._spTree
    clear_slide(slide)
    sid = 100
    N = len(overview['l2s'])

    add_title_ph(spTree, sid, 'Estructura Organizacional', font_pt=22); sid += 1

    txb = make_sp(spTree, sid, 'Subtitle', 330200, SUBTITLE_TOP, 10000000, SUBTITLE_H,
                  no_fill=True, no_border=True, tIns=0, bIns=0); sid += 1
    add_para(txb, dept_name, 18, italic=True, color='CC0000', align='l')

    L1_W = 3400000; L1_H = 500000
    L1_L = (SLIDE_W - L1_W) // 2
    txb = make_sp(spTree, sid, 'L1Box', L1_L, L2_TOP, L1_W, L1_H,
                  fill_hex='1F3864', border_hex='1F3864', border_pt=1,
                  tIns=35000, bIns=35000); sid += 1
    first = True
    for line in overview['l1_title'].split('\n'):
        add_para(txb, line, 11, bold=True, color='FFFFFF', align='ctr',
                 spc_before_pt=0 if first else 2)
        first = False

    L1_CX = L1_L + L1_W // 2
    HBAR_Y_ = L3_TOP - 100000
    add_line(spTree, sid, L1_CX, L2_TOP + L1_H, HBAR_Y_ - (L2_TOP + L1_H), vertical=True); sid += 1

    GAP = 35000
    L2_W_ = (SLIDE_W - 2 * SIDE_M - (N - 1) * GAP) // N
    L2_LEFTS = [SIDE_M + i * (L2_W_ + GAP) for i in range(N)]
    L2_CXS = [L + L2_W_ // 2 for L in L2_LEFTS]

    if N > 1:
        add_line(spTree, sid, L2_CXS[0], HBAR_Y_, L2_CXS[-1] - L2_CXS[0], vertical=False); sid += 1
    for cx in L2_CXS:
        add_line(spTree, sid, cx, HBAR_Y_, L3_TOP - HBAR_Y_, vertical=True); sid += 1

    for i, (l2_title, l2_desc) in enumerate(overview['l2s']):
        L2_L_ = L2_LEFTS[i]
        txb = make_sp(spTree, sid, f'L2_{i}', L2_L_, L3_TOP, L2_W_, L3_H,
                      fill_hex='1F3864', border_hex='1F3864', border_pt=1,
                      tIns=20000, bIns=20000); sid += 1
        first = True
        for line in l2_title.split('\n'):
            add_para(txb, line, 9, bold=True, color='FFFFFF', align='ctr',
                     spc_before_pt=0 if first else 1)
            first = False
        txb = make_sp(spTree, sid, f'L2Desc_{i}', L2_L_, DESC_TOP, L2_W_, DESC_H_STD,
                      fill_hex='EEF3FB', border_hex='BDD0E9', border_pt=0.75); sid += 1
        add_para(txb, l2_desc, 8, color='1F3864', align='l')

    print(f'  Built L1→L2 overview ({dept_name})')


# ── RASCI detail slide builder ────────────────────────────────

def rasci_color(letter):
    base = letter.split('/')[0].strip()
    return BADGE_COLORS.get(letter, BADGE_COLORS.get(base, 'A6A6A6'))


def build_rasci_slide(slide, role_shorts, role_fulls, resp, dept_name, start_id=50):
    spTree = slide.shapes._spTree
    clear_slide(slide)
    sid = start_id
    N = len(role_shorts)

    add_title_ph(spTree, sid, f'Matriz RASCI - Detalle | {dept_name}', font_pt=20); sid += 1

    txb = make_sp(spTree, sid, 'RespName', MARGIN, 700000, CONTENT_W, 720000,
                  no_fill=True, no_border=True, lIns=0, rIns=0, tIns=0, bIns=0); sid += 1
    add_para(txb, resp[1], 14, bold=True, color='1F3864')
    add_para(txb, resp[2], 10, color='595959', spc_before_pt=4)

    col_w = CONTENT_W // N
    for i, short in enumerate(role_shorts):
        x = MARGIN + i * col_w
        txb = make_sp(spTree, sid, f'RHdr_{i}', x, 1530000, col_w - 4000, 270000,
                      fill_hex='1F3864', no_border=True,
                      lIns=10000, rIns=10000, tIns=15000, bIns=10000); sid += 1
        add_para(txb, short, 7, bold=True, color='FFFFFF', align='ctr')

        letter = resp[3].get(i, 'I')
        bg = rasci_color(letter)
        txb = make_sp(spTree, sid, f'Badge_{i}', x, 1810000, col_w - 4000, 290000,
                      fill_hex=bg, no_border=True,
                      lIns=10000, rIns=10000, tIns=15000, bIns=10000); sid += 1
        add_para(txb, letter, 10, bold=True, color='FFFFFF', align='ctr')

    txb = make_sp(spTree, sid, 'Expl', MARGIN, 2200000, CONTENT_W, 460000,
                  fill_hex='F5F5F5', border_hex='D0D0D0', border_pt=0.5,
                  lIns=60000, rIns=60000, tIns=40000, bIns=40000); sid += 1
    add_para(txb, 'Razonamiento: ', 9, bold=True, color='1F3864')
    add_para(txb, resp[4], 9, color='404040', spc_before_pt=2)

    BADGE_W = 650000
    current_y = 2780000
    for i, full in enumerate(role_fulls):
        letter = resp[3].get(i, 'I')
        if letter == 'I':
            continue
        if current_y + 420000 > SLIDE_H - 80000:
            break
        bg = rasci_color(letter)
        txb = make_sp(spTree, sid, f'RBadge_{i}', MARGIN, current_y, BADGE_W, 400000,
                      fill_hex=bg, no_border=True,
                      lIns=10000, rIns=10000, tIns=20000, bIns=20000); sid += 1
        add_para(txb, letter, 12, bold=True, color='FFFFFF', align='ctr')

        desc_text = ROLE_DESCS.get(letter, 'Involucrado en esta responsabilidad.')
        txb = make_sp(spTree, sid, f'RDesc_{i}', MARGIN + BADGE_W + 25000,
                      current_y, CONTENT_W - BADGE_W - 25000, 400000,
                      fill_hex='F8FAFF', border_hex='C9D8F0', border_pt=0.5,
                      lIns=55000, rIns=55000, tIns=20000, bIns=20000); sid += 1
        add_para(txb, f'{full}: {desc_text}', 9, color='1F3864')
        current_y += 425000


# ── Intro slide text helpers ──────────────────────────────────

def replace_text_runs(slide, old, new):
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)


def update_agenda_slide(slide, new_items):
    best = None; best_n = 0
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            txt = shape.text_frame.text
            n = len([p for p in shape.text_frame.paragraphs if p.text.strip()])
            if n > best_n and ('1.' in txt or '1 ' in txt):
                best_n = n; best = shape
    if not best:
        return
    txBody = best.text_frame._txBody
    size_val = None
    for rPr in txBody.iter(f'{{{NS_A}}}rPr'):
        size_val = rPr.get('sz')
        break
    for p in list(txBody.findall(f'{{{NS_A}}}p')):
        txBody.remove(p)
    for item in new_items:
        p = etree.SubElement(txBody, f'{{{NS_A}}}p')
        r = etree.SubElement(p, f'{{{NS_A}}}r')
        rPr = etree.SubElement(r, f'{{{NS_A}}}rPr')
        if size_val:
            rPr.set('sz', size_val)
        etree.SubElement(r, f'{{{NS_A}}}t').text = item


def update_proposito_slide(slide, dept_name, proposito_text):
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and 'Propósito' in shape.text_frame.text:
            txBody = shape.text_frame._txBody
            size_val = None
            for rPr in txBody.iter(f'{{{NS_A}}}rPr'):
                size_val = rPr.get('sz')
                break
            for p in list(txBody.findall(f'{{{NS_A}}}p')):
                txBody.remove(p)
            p = etree.SubElement(txBody, f'{{{NS_A}}}p')
            r = etree.SubElement(p, f'{{{NS_A}}}r')
            rPr = etree.SubElement(r, f'{{{NS_A}}}rPr')
            rPr.set('b', '1')
            if size_val: rPr.set('sz', size_val)
            etree.SubElement(r, f'{{{NS_A}}}t').text = f'Propósito {dept_name}'
            etree.SubElement(txBody, f'{{{NS_A}}}p')
            p = etree.SubElement(txBody, f'{{{NS_A}}}p')
            r = etree.SubElement(p, f'{{{NS_A}}}r')
            rPr = etree.SubElement(r, f'{{{NS_A}}}rPr')
            rPr.set('b', '0')
            if size_val: rPr.set('sz', size_val)
            etree.SubElement(r, f'{{{NS_A}}}t').text = proposito_text
            break


# ── Anexos slide builder ──────────────────────────────────────

def build_anexos_slide(slide, dept_name, l2_l3_data, start_id=50):
    spTree = slide.shapes._spTree
    clear_slide(slide)
    sid = start_id

    add_title_ph(spTree, sid, 'Anexos\nDescripciones de Puesto', font_pt=22); sid += 1

    txb = make_sp(spTree, sid, 'Anexos', MARGIN, 750000, CONTENT_W, 5800000,
                  no_fill=True, no_border=True, lIns=0, rIns=0, tIns=0, bIns=0); sid += 1

    add_para(txb, f'• Director {dept_name} OXXO Mx:', 9, bold=True, color='1F3864')
    for l2_name, l3_names in l2_l3_data:
        add_para(txb, f'  ○ {l2_name} (ET L2)', 9, color='1F3864', spc_before_pt=2)
        for l3 in l3_names:
            add_para(txb, f'      ▪ {l3} (ET L3)', 8, color='262626', spc_before_pt=1)


# ── Main ──────────────────────────────────────────────────────

def main():
    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    shutil.copy2(SRC, OUT)
    print(f'Copied to: {OUT}')

    prs = Presentation(OUT)
    total = len(prs.slides)
    print(f'Slides: {total}')

    # 1. Update intro slides (0-4)
    print('Updating intro slides...')
    for i in range(5):
        replace_text_runs(prs.slides[i], 'Finanzas', ET_DEPT)
        replace_text_runs(prs.slides[i], 'finanzas', 'estrategia y transformación')
    update_agenda_slide(prs.slides[2], ET_AGENDA)
    update_proposito_slide(prs.slides[3], ET_DEPT, ET_PROPOSITO)

    # 2. Rebuild slide 5 (index 5) as L1→L2 ET overview
    print('Building L1→L2 overview...')
    build_l1_l2_overview(prs.slides[5], ET_L1_L2_OVERVIEW, ET_DEPT)

    # 3. Update RASCI methodology slide (index 6)
    replace_text_runs(prs.slides[6], 'Finanzas', ET_DEPT)

    # 4. Build 30 RASCI detail slides (indices 7-36)
    print('Building RASCI slides...')
    for i, resp in enumerate(ET_RASCI):
        idx = 7 + i
        print(f'  RASCI {i+1}/30: {resp[1][:50]}')
        build_rasci_slide(prs.slides[idx], ET_ROLE_SHORT, ET_ROLE_FULL,
                          resp, ET_DEPT, start_id=50 + i * 30)

    # 5. Clear + hide extra RASCI slides (indices 37-52, Finanzas had 46 RASCI)
    print('Hiding extra RASCI slots...')
    for i in range(37, 53):
        clear_slide(prs.slides[i])
        hide_slide(prs.slides[i])

    # 6. Build 5 org chart slides (indices 53-57)
    print('Building org chart slides...')
    for i, l2_info in enumerate(ET_L2_DATA):
        build_org_chart_slide(prs.slides[53 + i], l2_info, start_id=200 + i * 50)

    # 7. Clear + hide extra org chart slots (indices 58-59, Finanzas had 7)
    for i in range(58, 60):
        clear_slide(prs.slides[i])
        hide_slide(prs.slides[i])

    # 8. Rebuild Anexos (index 66)
    print('Building Anexos...')
    et_l2_l3 = [(l2['subtitle'], [l3['title'].replace('\n', ' ') for l3 in l2['l3s']])
                for l2 in ET_L2_DATA]
    build_anexos_slide(prs.slides[66], ET_DEPT, et_l2_l3)

    # 9. Save
    prs.save(OUT)
    print(f'Saved: {OUT}')

    # 10. Verify ZIP
    with zipfile.ZipFile(OUT, 'r') as z:
        names = z.namelist()
        dupes = [n for n in set(names) if names.count(n) > 1]
        ct = z.read('[Content_Types].xml').decode()
        ct_slides = set(re.findall(r'PartName="/ppt/slides/(slide\d+\.xml)"', ct))
        zip_slides = set(n.replace('ppt/slides/', '') for n in names
                         if n.startswith('ppt/slides/slide') and '.rels' not in n)
        print(f'ZIP dupes: {dupes}')
        print(f'Orphans: {zip_slides - ct_slides}')

    prs2 = Presentation(OUT)
    print(f'Final slides: {len(prs2.slides)}')


if __name__ == '__main__':
    main()
