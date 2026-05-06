#!/usr/bin/env python3
"""Generate Finanzas Playbook PPTX from template."""

import copy
import shutil
import os
from lxml import etree
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

TEMPLATE_PATH = "/home/user/oxxo/Playbook Excelencia Operativa - playbook ejemplo.pptx"
OUTPUT_PATH = "/home/user/oxxo/Finanzas/Playbook/Finanzas Playbook.pptx"
RASCI_PATH = "/home/user/oxxo/Finanzas/Matriz RASCI - Finanzas.xlsx"

# ── helpers ────────────────────────────────────────────────────────────────────

def delete_slide(prs, slide_index):
    """Remove a slide by index using XML manipulation."""
    slides = prs.slides
    sldId_elem = slides._sldIdLst[slide_index]
    rId = sldId_elem.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    slides._sldIdLst.remove(sldId_elem)
    # Remove the relationship from the presentation part
    rels = prs.part.rels
    if rId in rels:
        rels.pop(rId)


def add_slide(prs, layout_index=11):
    """Add a new slide with the given layout index."""
    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)
    return slide


def set_text_frame(tf, text, font_size_pt=None, bold=None, color=None, alignment=None):
    """Clear and set simple text in a text frame."""
    tf.clear()
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text
    if font_size_pt:
        run.font.size = Pt(font_size_pt)
    if bold is not None:
        run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if alignment:
        para.alignment = alignment


def add_text_box(slide, left, top, width, height, text, font_size_pt=12, bold=False,
                  color=None, bg_color=None, alignment=PP_ALIGN.LEFT, word_wrap=True):
    """Add a text box with specified properties."""
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    para = tf.paragraphs[0]
    para.alignment = alignment
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_paragraph(tf, text, font_size_pt=12, bold=False, color=None, level=0, alignment=PP_ALIGN.LEFT):
    """Add a paragraph to a text frame."""
    from pptx.oxml.ns import qn
    from pptx.util import Pt
    para = tf.add_paragraph()
    para.level = level
    para.alignment = alignment
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return para


def set_shape_fill(shape, rgb_color):
    """Set solid fill color for a shape."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color


# ── RASCI data ─────────────────────────────────────────────────────────────────

ROLE_COLS = {
    3: "Director Adm. y Finanzas (DAF)",
    4: "Gte Sr Planeación Fin & CF",
    5: "Gte Sr Fiscal y Contraloría",
    6: "Gte Sr Transformación Procesos",
    7: "Gte Sr Servicios Compartidos",
    8: "Gte Sr Legal OXXO MX",
    9: "Gte Asuntos Corporativos",
    10: "Dir. Adm. Nacional OXXO",
    11: "Gte Sr Nac. Adm. CdS",
}

RASCI_COLORS = {
    'R': RGBColor(0xE5, 0x25, 0x2A),   # red
    'A': RGBColor(0xF5, 0xA6, 0x23),   # orange
    'S': RGBColor(0x00, 0x70, 0xC0),   # blue
    'C': RGBColor(0x70, 0xAD, 0x47),   # green
    'I': RGBColor(0x76, 0x76, 0x76),   # gray
}

RASCI_FULL = {
    'R': 'Responsible', 'A': 'Accountable', 'S': 'Support', 'C': 'Consulted', 'I': 'Informed'
}


def load_rasci_data():
    """Load RASCI data from Excel."""
    wb = openpyxl.load_workbook(RASCI_PATH)
    ws = wb['Hoja1']

    responsibilities = []
    for row in ws.iter_rows(min_row=7, values_only=True):
        # Row: col A(0)=section, B(1)=responsibility, C(2)=description, D-L(3-11)=roles, M(12)=explanation
        resp_name = row[1]
        if not resp_name:
            continue
        description = row[2] if len(row) > 2 else ""
        explanation = row[12] if len(row) > 12 else ""
        roles = {}
        for ci in range(3, 12):
            val = row[ci] if ci < len(row) else None
            if val and str(val).strip():
                roles[ci] = str(val).strip()
        responsibilities.append({
            'name': resp_name,
            'description': description or "",
            'explanation': explanation or "",
            'roles': roles,
        })
    return responsibilities


# ── Slide modification functions ────────────────────────────────────────────────

def modify_slide1(slide):
    """Slide 1 – Cover: Change title and subtitle."""
    print("  Modifying slide 1 (Cover)...")
    for shape in slide.shapes:
        if shape.name == 'Title 1' and shape.has_text_frame:
            tf = shape.text_frame
            # Clear and set new text preserving formatting
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ""
            tf.paragraphs[0].runs[0].text = "Playbook: Finanzas"
        elif shape.name == 'TextBox 16' and shape.has_text_frame:
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ""
            if tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = "Fit for Purpose"
            else:
                run = tf.paragraphs[0].add_run()
                run.text = "Fit for Purpose"


def modify_slide3(slide):
    """Slide 3 – Agenda: Replace TextBox 2440 content."""
    print("  Modifying slide 3 (Agenda)...")
    agenda_items = [
        "1. Macroestructura Oxxo MX",
        "2. Mapa Interacciones Finanzas",
        "3. Estructura Organizacional Finanzas",
        "4. Matriz RASCI (Metodología)",
        "5. Detalle de la Matriz RASCI por Responsabilidad",
        "6. Estructura Organizacional PF&A y Commercial Finance",
        "7. Estructura Organizacional Fiscal y Contraloría",
        "8. Estructura Organizacional Servicios Compartidos Administrativo",
        "9. Estructura Organizacional Legal y Regulatorio",
        "10. Estructura Organizacional Asuntos Corporativos",
        "11. Estructura Organizacional Administrativo Nacional OXXO",
        "12. Estructura Organizacional Administrativo CdS",
        "13. Anexos Descripciones de Puesto",
    ]
    for shape in slide.shapes:
        if shape.name == 'TextBox 2440' and shape.has_text_frame:
            tf = shape.text_frame
            # Get font properties from first paragraph/run
            first_run = None
            for para in tf.paragraphs:
                for run in para.runs:
                    first_run = run
                    break
                if first_run:
                    break

            orig_font_size = first_run.font.size if first_run else Pt(14)
            orig_font_color = None
            try:
                orig_font_color = first_run.font.color.rgb if first_run else None
            except:
                pass

            # Clear all paragraphs
            from pptx.oxml.ns import qn
            txPr = tf._txBody
            # Remove all 'a:p' elements
            for p in txPr.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}p'):
                txPr.remove(p)

            # Add new paragraphs for each agenda item
            for i, item in enumerate(agenda_items):
                para = tf.add_paragraph()
                run = para.add_run()
                run.text = item
                if orig_font_size:
                    run.font.size = orig_font_size
                if orig_font_color:
                    try:
                        run.font.color.rgb = orig_font_color
                    except:
                        pass
            break


def modify_slide4(slide):
    """Slide 4 – Macroestructura: Replace purpose text."""
    print("  Modifying slide 4 (Macroestructura)...")
    new_text = ("Propósito Finanzas\n\n"
                "Contribuir a la generación de valor y rentabilidad del negocio mediante la integración de "
                "capacidades financieras, fiscales, legales, administrativas y de control, asegurando "
                "información confiable y oportuna para la toma de decisiones, una gestión eficiente de "
                "procesos transaccionales y una postura sólida de cumplimiento, gobierno y riesgo que "
                "habilite el crecimiento sostenible de OXXO México.")

    for shape in slide.shapes:
        if shape.name == 'TextBox 2829' and shape.has_text_frame:
            tf = shape.text_frame
            # Get styling from first run
            first_run = None
            first_bold_run = None
            for para in tf.paragraphs:
                for run in para.runs:
                    if first_run is None:
                        first_run = run
                    if run.font.bold and first_bold_run is None:
                        first_bold_run = run

            orig_font_size = first_run.font.size if first_run else None

            # Clear existing content
            from pptx.oxml.ns import qn
            txBody = tf._txBody
            for p in txBody.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}p'):
                txBody.remove(p)

            lines = new_text.split('\n')
            for i, line in enumerate(lines):
                para = tf.add_paragraph()
                run = para.add_run()
                run.text = line
                if orig_font_size:
                    run.font.size = orig_font_size
                if i == 0:  # Title line bold
                    run.font.bold = True
            break


def modify_slide5(slide):
    """Slide 5 – Mapa de Interacciones: Change title."""
    print("  Modifying slide 5 (Mapa de Interacciones)...")
    for shape in slide.shapes:
        if shape.name == 'Title 1' and shape.has_text_frame:
            tf = shape.text_frame
            for para in tf.paragraphs:
                runs = para.runs
                # Find runs containing 'Excelencia' and 'Operativa'
                # Combine them into one run with 'Finanzas'
                excelencia_idx = None
                for ri, run in enumerate(runs):
                    if 'Excelencia' in run.text:
                        excelencia_idx = ri
                        break
                if excelencia_idx is not None:
                    # Replace the Excelencia run with Finanzas and clear subsequent
                    runs[excelencia_idx].text = 'Finanzas'
                    # Clear any runs after that in same para that are part of the old name
                    for ri in range(excelencia_idx + 1, len(runs)):
                        if 'Operativa' in runs[ri].text or runs[ri].text.strip() == '':
                            runs[ri].text = ''
            break


def modify_slide6(slide):
    """Slide 6 – Estructura Organizacional L2: Change title and 3 description boxes."""
    print("  Modifying slide 6 (Estructura Organizacional L2)...")

    box_texts = [
        {
            'title': "Director PF&A y Commercial Finance",
            'bullets': [
                "Funge como copiloto financiero del negocio, habilitando decisiones basadas en datos.",
                "Gestiona el proceso integral de planeación financiera, presupuesto anual y previsiones.",
                "Analiza rendimiento real vs. planificado e impulsa conversaciones de mejora.",
                "Maximiza rentabilidad de categorías y brinda soporte financiero a iniciativas estratégicas.",
            ]
        },
        {
            'title': "Gerente Sr Fiscal y Contraloría / Gerente Sr Servicios Compartidos",
            'bullets': [
                "Fiscal: Define y dirige la estrategia fiscal, cumplimiento normativo y control interno.",
                "Fiscal: Gestiona requerimientos de autoridades, auditorías y prevención de lavado de dinero.",
                "SSC: Lidera la operación centralizada de P2P, O2C y R2R con altos estándares de servicio.",
                "SSC: Define modelos de servicios centrales y garantiza cierres contables precisos y oportunos.",
            ]
        },
        {
            'title': "Gerente Sr Legal / Gerente Asuntos Corporativos / Área Administrativa",
            'bullets': [
                "Legal: Gestiona cumplimiento, litigios y gobernanza contractual para proteger la operación.",
                "Asuntos Corp: Lidera la estrategia de relacionamiento externo y manejo de riesgos reputacionales.",
                "Dir. Adm. OXXO: Asegura procesos administrativos, inventarios y control en todas las plazas.",
                "Gte. Adm. CdS: Impulsa control administrativo y productividad en CEDIS y distribución.",
            ]
        },
    ]

    # Change title
    for shape in slide.shapes:
        if shape.name == 'Title 1' and shape.has_text_frame:
            tf = shape.text_frame
            for para in tf.paragraphs:
                runs = para.runs
                for ri, run in enumerate(runs):
                    if 'Excelencia' in run.text:
                        run.text = run.text.replace('Excelencia', 'Finanzas')
                    elif 'Operativa' in run.text:
                        run.text = run.text.replace('Operativa', '')

    # Find and update the Rectángulo 2 shapes
    rect_shapes = [s for s in slide.shapes if s.name == 'Rectángulo 2']

    for i, shape in enumerate(rect_shapes[:3]):
        if i < len(box_texts):
            box_data = box_texts[i]
            tf = shape.text_frame
            tf.word_wrap = True

            # Clear existing content
            from pptx.oxml.ns import qn
            txBody = tf._txBody
            for p in txBody.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}p'):
                txBody.remove(p)

            # Add title paragraph
            para = tf.add_paragraph()
            run = para.add_run()
            run.text = box_data['title']
            run.font.size = Pt(9)
            run.font.bold = True

            # Add bullet paragraphs
            for bullet in box_data['bullets']:
                para = tf.add_paragraph()
                run = para.add_run()
                run.text = "• " + bullet
                run.font.size = Pt(8)
                run.font.bold = False


def add_rasci_detail_slide(prs, resp_data):
    """Add a single RASCI detail slide."""
    slide = add_slide(prs, layout_index=11)

    # Slide dimensions
    slide_width = prs.slide_width   # EMU
    slide_height = prs.slide_height

    margin = Inches(0.4)
    content_width = slide_width - 2 * margin

    # ── Title ──
    title_shape = None
    for shape in slide.placeholders:
        title_shape = shape
        break

    if title_shape:
        tf = title_shape.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = "Matriz RASCI - Detalle | Finanzas"
        run.font.size = Pt(22)
        run.font.bold = True
    else:
        # Add title text box
        txb = slide.shapes.add_textbox(margin, Inches(0.15), content_width, Inches(0.55))
        tf = txb.text_frame
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = "Matriz RASCI - Detalle | Finanzas"
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)

    # ── Responsibility header box ──
    resp_top = Inches(0.85)
    resp_height = Inches(0.7)
    resp_box = slide.shapes.add_textbox(margin, resp_top, content_width, resp_height)
    resp_tf = resp_box.text_frame
    resp_tf.word_wrap = True
    # (no background fill needed for resp_box)

    # Responsibility name (bold)
    para1 = resp_tf.paragraphs[0]
    run1 = para1.add_run()
    run1.text = resp_data['name']
    run1.font.size = Pt(14)
    run1.font.bold = True
    run1.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)

    # Description on next paragraph if available
    if resp_data.get('description'):
        para2 = resp_tf.add_paragraph()
        run2 = para2.add_run()
        run2.text = resp_data['description']
        run2.font.size = Pt(10)
        run2.font.bold = False
        run2.font.color.rgb = RGBColor(0x40, 0x40, 0x40)

    # ── Role assignment table ──
    roles = resp_data['roles']
    role_names = list(ROLE_COLS.values())

    # Build the role assignments display
    assignments_top = Inches(1.65)
    assignments_height = Inches(0.5)

    # Create a simple table showing role -> RASCI letter
    n_roles_with_assignment = len(roles)

    if n_roles_with_assignment > 0:
        # Table: header row + assignment row
        col_count = len(ROLE_COLS) + 1  # +1 for "Responsabilidad" column
        table_width = content_width
        col_width = table_width // col_count

        table = slide.shapes.add_table(
            2, col_count,
            margin, assignments_top,
            table_width, assignments_height
        ).table

        # Header row
        table.cell(0, 0).text = "Responsabilidad"
        for ci, (col_idx, role_name) in enumerate(ROLE_COLS.items()):
            # Short role name
            short = role_name.replace("Director ", "Dir. ").replace("Gerente ", "Gte. ").replace("Administración", "Adm.")
            table.cell(0, ci + 1).text = short

        # Assignment row
        table.cell(1, 0).text = resp_data['name'][:40]
        for ci, col_idx in enumerate(ROLE_COLS.keys()):
            val = roles.get(col_idx, "")
            table.cell(1, ci + 1).text = val

        # Style table cells
        from pptx.oxml.ns import qn
        for ri in range(2):
            for ci in range(col_count):
                cell = table.cell(ri, ci)
                tf = cell.text_frame
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(7)
                        if ri == 0:
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                        else:
                            run.font.bold = False
                    para.alignment = PP_ALIGN.CENTER

                # Header background
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                from lxml import etree
                solidFill = etree.SubElement(tcPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
                srgbClr = etree.SubElement(solidFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                if ri == 0:
                    srgbClr.set('val', '1F3964')
                else:
                    # Color by RASCI value if single letter
                    val = roles.get(list(ROLE_COLS.keys())[ci - 1] if ci > 0 else -1, "")
                    letter = val[0] if val else ""
                    color_map = {'R': 'E5252A', 'A': 'F5A623', 'S': '0070C0', 'C': '70AD47', 'I': '767676'}
                    bg = color_map.get(letter, 'FFFFFF')
                    srgbClr.set('val', bg)
                    if ci > 0 and val:
                        for para in tf.paragraphs:
                            for run in para.runs:
                                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # ── RASCI explanation ──
    explanation = resp_data.get('explanation', '')
    if not explanation:
        explanation = "Ver detalle de roles y responsabilidades en la Matriz RASCI."

    explanation_top = Inches(2.25)
    explanation_height = Inches(0.8)

    exp_box = slide.shapes.add_textbox(margin, explanation_top, content_width, explanation_height)
    exp_tf = exp_box.text_frame
    exp_tf.word_wrap = True

    # Label
    para_label = exp_tf.paragraphs[0]
    run_label = para_label.add_run()
    run_label.text = "Explicación RASCI:"
    run_label.font.size = Pt(10)
    run_label.font.bold = True
    run_label.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)

    para_exp = exp_tf.add_paragraph()
    run_exp = para_exp.add_run()
    run_exp.text = explanation
    run_exp.font.size = Pt(10)
    run_exp.font.bold = False
    run_exp.font.color.rgb = RGBColor(0x40, 0x40, 0x40)

    # ── Per-role description boxes ──
    role_top = Inches(3.15)
    role_box_height = Inches(0.52)
    role_spacing = Inches(0.02)
    badge_width = Inches(0.7)
    desc_width = content_width - badge_width - Inches(0.1)

    ROLE_DESC_TEMPLATES = {
        'R': "{role}: Ejecuta y lidera esta responsabilidad, tomando las decisiones necesarias para garantizar su cumplimiento y calidad de resultados.",
        'A/R': "{role}: Ejecuta y lidera esta responsabilidad, tomando las decisiones necesarias para garantizar su cumplimiento y calidad de resultados.",
        'A': "{role}: Aprueba y supervisa esta responsabilidad, siendo responsable final ante la organización por sus resultados.",
        'S': "{role}: Brinda apoyo activo en esta responsabilidad, aportando su expertise especializado cuando se requiere.",
        'C': "{role}: Aporta perspectiva experta como área consultada, contribuyendo con información y criterio antes de las decisiones clave.",
        'I': "{role}: Recibe información del avance y resultados para anticipar impactos en su área y tomar acciones preventivas.",
    }

    current_top = role_top
    for col_idx, assignment in roles.items():
        role_name = ROLE_COLS.get(col_idx, f"Role {col_idx}")
        letter = assignment.strip()

        # RASCI badge
        badge_color = RASCI_COLORS.get(letter[0] if letter else 'I', RGBColor(0x76, 0x76, 0x76))

        badge_box = slide.shapes.add_textbox(margin, current_top, badge_width, role_box_height)
        badge_tf = badge_box.text_frame
        badge_tf.word_wrap = False
        badge_para = badge_tf.paragraphs[0]
        badge_para.alignment = PP_ALIGN.CENTER
        badge_run = badge_para.add_run()
        badge_run.text = letter
        badge_run.font.size = Pt(13)
        badge_run.font.bold = True
        badge_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Set badge background color via XML
        sp = badge_box._element
        spPr = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
        if spPr is None:
            spPr = etree.SubElement(sp, '{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
        solidFill = etree.SubElement(spPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        srgbClr = etree.SubElement(solidFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        srgbClr.set('val', str(badge_color))

        # Description text box
        desc_left = margin + badge_width + Inches(0.1)

        # Use template description (explanation already shown above, so here use role-specific template)
        template = ROLE_DESC_TEMPLATES.get(letter, ROLE_DESC_TEMPLATES['I'])
        desc_text = template.format(role=role_name)

        desc_box = slide.shapes.add_textbox(desc_left, current_top, desc_width, role_box_height)
        desc_tf = desc_box.text_frame
        desc_tf.word_wrap = True
        desc_para = desc_tf.paragraphs[0]
        desc_run = desc_para.add_run()
        desc_run.text = desc_text
        desc_run.font.size = Pt(9)
        # Make role name bold
        desc_run.font.bold = False

        current_top += role_box_height + role_spacing

        # Check if we're going off screen
        if current_top > slide_height - Inches(0.3):
            break

    return slide


def add_l3_structure_slide(prs, l2_name, l3_roles):
    """Add an L3 structure slide."""
    slide = add_slide(prs, layout_index=11)

    margin = Inches(0.4)
    slide_width = prs.slide_width
    content_width = slide_width - 2 * margin

    # Title placeholder
    title_set = False
    for shape in slide.placeholders:
        tf = shape.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = f"Estructura Organizacional\n{l2_name}"
        run.font.size = Pt(22)
        run.font.bold = True
        title_set = True
        break

    if not title_set:
        txb = slide.shapes.add_textbox(margin, Inches(0.15), content_width, Inches(0.7))
        tf = txb.text_frame
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = f"Estructura Organizacional | {l2_name}"
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)

    # Content box with L3 roles
    content_top = Inches(1.1)
    content_height = Inches(5.2)

    content_box = slide.shapes.add_textbox(margin, content_top, content_width, content_height)
    tf = content_box.text_frame
    tf.word_wrap = True

    # Header
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = f"Reportan al/a la {l2_name}:"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)

    for role in l3_roles:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = f"  • {role}"
        r.font.size = Pt(12)
        r.font.bold = False
        r.font.color.rgb = RGBColor(0x26, 0x26, 0x26)

    return slide


def add_annexes_slide(prs):
    """Add the Anexos slide."""
    print("  Adding Anexos slide...")
    slide = add_slide(prs, layout_index=11)

    margin = Inches(0.4)
    slide_width = prs.slide_width
    content_width = slide_width - 2 * margin

    # Title
    for shape in slide.placeholders:
        tf = shape.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = "Anexos\nDescripciones de Puesto"
        run.font.size = Pt(22)
        run.font.bold = True
        break

    content_top = Inches(1.1)
    content_height = Inches(5.4)
    content_box = slide.shapes.add_textbox(margin, content_top, content_width, content_height)
    tf = content_box.text_frame
    tf.word_wrap = True

    annexes_data = [
        ("• Director Administración y Finanzas:", 0, True),
        ("  ○ Director PF&A y Commercial Finance (FIN L2)", 1, False),
        ("      ▪ Gerente BP Finanzas Plataformas RH, TI (FIN L3)", 2, False),
        ("      ▪ Gerente Commercial Finance (FIN L3)", 2, False),
        ("      ▪ Gerente Planeación Financiera (FIN L3)", 2, False),
        ("  ○ Gerente Sr Fiscal y Contraloría (FIN L2)", 1, False),
        ("      ▪ Gerente Sr Contraloría (FIN L3)", 2, False),
        ("      ▪ Gerente Control Fiscal (FIN L3)", 2, False),
        ("      ▪ Gerente Área Normatividad y Control (FIN L3)", 2, False),
        ("      ▪ Gerente Gestión Riesgos SI (FIN L3)", 2, False),
        ("      ▪ Responsable Auditorías Fiscales (FIN L3)", 2, False),
        ("  ○ Gerente Sr Servicios Compartidos Administrativo (FIN L2)", 1, False),
        ("      ▪ Gerente Purchase to Pay (FIN L3)", 2, False),
        ("      ▪ Gerente Record to Report (FIN L3)", 2, False),
        ("      ▪ Gerente Order to Cash (FIN L3)", 2, False),
        ("      ▪ Responsable Gobierno Datos y Cond. (FIN L3)", 2, False),
        ("      ▪ Coordinador Información y Gestión CSF (FIN L3)", 2, False),
        ("  ○ Gerente Sr Legal OXXO MX (FIN L2)", 1, False),
        ("      ▪ Gerente Sr Legal Operaciones (FIN L3)", 2, False),
        ("      ▪ Gerente Área Legal Comercial (FIN L3)", 2, False),
        ("      ▪ Gerente Área Cumplimiento (FIN L3)", 2, False),
        ("      ▪ Gerente Legal CdS (FIN L3)", 2, False),
        ("      ▪ Responsable Legal (i78) (FIN L3)", 2, False),
        ("  ○ Gerente Asuntos Corporativos (FIN L2)", 1, False),
        ("      ▪ Gerente Área Atención Autoridades (FIN L3)", 2, False),
        ("      ▪ Gerente Área Relaciones Gobierno (FIN L3)", 2, False),
        ("      ▪ Gerente Relaciones Gobierno (FIN L3)", 2, False),
        ("      ▪ Coordinador Asuntos Corporativos (FIN L3)", 2, False),
        ("      ▪ Analista Comunicación Externa (FIN L3)", 2, False),
        ("  ○ Director Administrativo Nacional OXXO (FIN L2)", 1, False),
        ("      ▪ Gerente Área Admvo Neg Dllo (FIN L3)", 2, False),
        ("      ▪ Gerente Área BP Finanzas Plataformas Expansión (FIN L3)", 2, False),
        ("      ▪ Gerente Gestión Control Zona (FIN L3)", 2, False),
        ("      ▪ Responsable Gestión y Control (FIN L3)", 2, False),
        ("      ▪ Responsable Traslado de Valores (FIN L3)", 2, False),
        ("  ○ Gerente Sr Nacional Administrativo CdS (FIN L2)", 1, False),
        ("      ▪ Gerente Área Admin Control CEDIS (FIN L3)", 2, False),
        ("      ▪ Gerente Área Admin y Finanzas Dist. (FIN L3)", 2, False),
        ("      ▪ Responsable Gestión y Control CEDIS (FIN L3)", 2, False),
        ("      ▪ Responsable Gestoría (FIN L3)", 2, False),
    ]

    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = annexes_data[0][0]
    run.font.size = Pt(9)
    run.font.bold = annexes_data[0][2]
    run.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)

    for text, level, bold in annexes_data[1:]:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = text
        r.font.size = Pt(9)
        r.font.bold = bold
        if bold:
            r.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)
        else:
            r.font.color.rgb = RGBColor(0x26, 0x26, 0x26)

    return slide


def add_closing_slide(prs):
    """Add the closing slide."""
    print("  Adding closing slide...")
    slide = add_slide(prs, layout_index=11)

    margin = Inches(0.4)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    content_width = slide_width - 2 * margin

    # Title
    for shape in slide.placeholders:
        tf = shape.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = "Gracias"
        run.font.size = Pt(40)
        run.font.bold = True
        break

    # Date text box
    date_box = slide.shapes.add_textbox(
        margin, slide_height - Inches(1.0), content_width, Inches(0.5)
    )
    tf = date_box.text_frame
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = "Mayo 2026"
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x40, 0x40, 0x40)

    return slide


# ── Main ────────────────────────────────────────────────────────────────────────

def main():
    print("Loading template...")
    prs = Presentation(TEMPLATE_PATH)

    print(f"Template has {len(prs.slides)} slides")

    # ── Step 1: Modify existing slides 1-7 ──
    print("\nModifying existing slides...")

    modify_slide1(prs.slides[0])   # Slide 1 – Cover
    # Slide 2 – Propósito: no changes
    print("  Slide 2 (Propósito): no changes")
    modify_slide3(prs.slides[2])   # Slide 3 – Agenda
    modify_slide4(prs.slides[3])   # Slide 4 – Macroestructura
    modify_slide5(prs.slides[4])   # Slide 5 – Mapa Interacciones
    modify_slide6(prs.slides[5])   # Slide 6 – L2 Structure
    # Slide 7 – Metodología RASCI: no changes
    print("  Slide 7 (Metodología RASCI): no changes")

    # ── Step 2: Delete slides 8-20 (indices 7-19) from back to front ──
    print("\nDeleting slides 8 through 20...")
    # We need to delete slides 8 onwards (indices 7+)
    # Delete from the back to avoid index shifting
    total = len(prs.slides)
    for i in range(total - 1, 6, -1):
        delete_slide(prs, i)
        print(f"  Deleted slide at index {i}")

    print(f"Slides remaining: {len(prs.slides)}")

    # ── Step 3: Load RASCI data ──
    print("\nLoading RASCI data...")
    rasci_data = load_rasci_data()
    print(f"  Loaded {len(rasci_data)} responsibilities")

    # ── Step 4: Add RASCI detail slides ──
    print("\nAdding RASCI detail slides...")
    for i, resp in enumerate(rasci_data):
        print(f"  Adding RASCI slide {i+1}/46: {resp['name'][:50]}")
        add_rasci_detail_slide(prs, resp)

    # ── Step 5: Add L3 Structure slides ──
    print("\nAdding L3 structure slides...")

    l3_data = [
        ("PF&A y Commercial Finance", [
            "Gerente BP Finanzas Plataformas RH, TI",
            "Gerente Commercial Finance",
            "Gerente Planeación Financiera",
        ]),
        ("Fiscal y Contraloría", [
            "Responsable Auditorías Fiscales",
            "Gerente Control Fiscal",
            "Gerente Sr Contraloría",
            "Gerente Área Normatividad y Control",
            "Gerente Gestión Riesgos SI",
        ]),
        ("Servicios Compartidos Administrativo", [
            "Gerente Purchase to Pay",
            "Gerente Record to Report",
            "Gerente Order to Cash",
            "Responsable Gobierno Datos y Cond.",
            "Coordinador Información y Gestión CSF",
        ]),
        ("Legal y Regulatorio", [
            "Gerente Sr Legal Operaciones",
            "Gerente Área Legal Comercial",
            "Gerente Área Cumplimiento",
            "Gerente Legal CdS",
            "Responsable Legal (i78)",
        ]),
        ("Asuntos Corporativos", [
            "Gerente Área Atención Autoridades",
            "Gerente Área Relaciones Gobierno",
            "Gerente Relaciones Gobierno",
            "Coordinador Asuntos Corporativos",
            "Analista Comunicación Externa",
        ]),
        ("Administrativo Nacional OXXO", [
            "Gerente Área Admvo Neg Dllo",
            "Gerente Área BP Finanzas Plataformas Expansión",
            "Gerente Gestión Control Zona (x4)",
            "Responsable Gestión y Control",
            "Responsable Traslado de Valores",
        ]),
        ("Administrativo CdS", [
            "Gerente Área Admin Control CEDIS",
            "Gerente Área Admin y Finanzas Dist.",
            "Responsable Gestión y Control CEDIS",
            "Responsable Gestoría",
        ]),
    ]

    for l2_name, l3_roles in l3_data:
        print(f"  Adding L3 slide: {l2_name}")
        add_l3_structure_slide(prs, l2_name, l3_roles)

    # ── Step 6: Add Annexes and Closing slides ──
    print("\nAdding final slides...")
    add_annexes_slide(prs)
    add_closing_slide(prs)

    # ── Step 7: Save ──
    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    print(f"\nSaving to: {OUTPUT_PATH}")
    prs.save(OUTPUT_PATH)

    total_slides = len(prs.slides)
    print(f"\n{'='*60}")
    print(f"DONE! Total slides generated: {total_slides}")
    print(f"Breakdown:")
    print(f"  Slides 1-7:  Original template slides (modified)")
    print(f"  Slides 8-{7+len(rasci_data)}:  RASCI detail slides ({len(rasci_data)} responsibilities)")
    l3_start = 8 + len(rasci_data)
    print(f"  Slides {l3_start}-{l3_start+len(l3_data)-1}: L3 structure slides ({len(l3_data)} areas)")
    print(f"  Slide {l3_start+len(l3_data)}:  Anexos slide")
    print(f"  Slide {l3_start+len(l3_data)+1}:  Closing slide")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()
