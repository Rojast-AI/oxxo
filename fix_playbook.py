#!/usr/bin/env python3
"""
Comprehensive fix script for Finanzas Playbook.pptx
Applies Fix 1, 2, 3, and 4 in order.
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from lxml import etree
import copy

PPTX_PATH = '/home/user/oxxo/Finanzas/Playbook/Finanzas Playbook.pptx'

NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'


def get_spPr(shape):
    """Get the spPr element from a shape, searching under both namespaces."""
    sp = shape._element
    # Try p: namespace first (presentation shapes use p:spPr)
    spPr = sp.find(f'{{{NS_P}}}spPr')
    if spPr is not None:
        return spPr
    spPr = sp.find(f'{{{NS_A}}}spPr')
    if spPr is not None:
        return spPr
    # Deep search
    for ns in [NS_P, NS_A]:
        spPr = sp.find(f'.//{{{ns}}}spPr')
        if spPr is not None:
            return spPr
    return None


def set_shape_fill(shape, hex_color):
    """Set solid fill color on a shape."""
    spPr = get_spPr(shape)
    if spPr is None:
        return
    # Remove noFill and existing solidFill
    for nf in spPr.findall(f'{{{NS_A}}}noFill'):
        spPr.remove(nf)
    for sf in spPr.findall(f'{{{NS_A}}}solidFill'):
        spPr.remove(sf)
    solidFill = etree.SubElement(spPr, f'{{{NS_A}}}solidFill')
    srgbClr = etree.SubElement(solidFill, f'{{{NS_A}}}srgbClr')
    srgbClr.set('val', hex_color)


def add_border(shape, hex_color, width_pt=1):
    """Add a solid border to a shape."""
    spPr = get_spPr(shape)
    if spPr is None:
        return
    ln = spPr.find(f'{{{NS_A}}}ln')
    if ln is None:
        ln = etree.SubElement(spPr, f'{{{NS_A}}}ln')
    ln.set('w', str(int(width_pt * 12700)))
    solidFill = ln.find(f'{{{NS_A}}}solidFill')
    if solidFill is None:
        solidFill = etree.SubElement(ln, f'{{{NS_A}}}solidFill')
    srgbClr = solidFill.find(f'{{{NS_A}}}srgbClr')
    if srgbClr is None:
        srgbClr = etree.SubElement(solidFill, f'{{{NS_A}}}srgbClr')
    srgbClr.set('val', hex_color)


def reposition_via_sp(shape, new_left=None, new_top=None, new_width=None, new_height=None):
    """Reposition/resize a shape by directly modifying its xfrm element."""
    sp = shape._element
    # Search for xfrm in spPr under p: namespace
    xfrm = None
    for spPr_tag in [f'{{{NS_P}}}spPr', f'{{{NS_A}}}spPr']:
        spPr = sp.find(spPr_tag)
        if spPr is not None:
            xfrm = spPr.find(f'{{{NS_A}}}xfrm')
            if xfrm is not None:
                break
    if xfrm is None:
        # Try deeper search
        xfrm = sp.find(f'.//{{{NS_A}}}xfrm')
    if xfrm is None:
        return
    off = xfrm.find(f'{{{NS_A}}}off')
    ext = xfrm.find(f'{{{NS_A}}}ext')
    if off is not None:
        if new_left is not None:
            off.set('x', str(int(new_left)))
        if new_top is not None:
            off.set('y', str(int(new_top)))
    if ext is not None:
        if new_width is not None:
            ext.set('cx', str(int(new_width)))
        if new_height is not None:
            ext.set('cy', str(int(new_height)))


# ─────────────────────────────────────────────
# FIX 1: Slide 4 (index 3) — Change ALL org chart text to black
# ─────────────────────────────────────────────
def fix1_slide4_text_black(prs):
    print("Applying Fix 1: Slide 4 - Set all org chart text to black...")
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if shape.name == 'Title 1':
            # Keep title red (CC0000) — skip
            continue
        for rPr in shape._element.iter(f'{{{NS_A}}}rPr'):
            # Remove existing solidFill
            for sf in rPr.findall(f'{{{NS_A}}}solidFill'):
                rPr.remove(sf)
            # Add black solidFill
            solidFill = etree.SubElement(rPr, f'{{{NS_A}}}solidFill')
            srgbClr = etree.SubElement(solidFill, f'{{{NS_A}}}srgbClr')
            srgbClr.set('val', '000000')
    print("  Fix 1 done.")


# ─────────────────────────────────────────────
# FIX 2: Slide 5 (index 4) — Replace image with placeholder text
# ─────────────────────────────────────────────
def fix2_slide5_image_placeholder(prs):
    print("Applying Fix 2: Slide 5 - Replace image with placeholder text...")
    slide5 = prs.slides[4]

    # 1. Find and remove the PICTURE shape
    pic_shape = None
    for shape in slide5.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            pic_shape = shape
            break

    if pic_shape:
        pic_shape._element.getparent().remove(pic_shape._element)
        print(f"  Removed picture shape.")
    else:
        print("  No picture shape found.")

    # 2. Add new placeholder text box
    from pptx.util import Emu, Pt
    txBox = slide5.shapes.add_textbox(
        left=Emu(2000000),
        top=Emu(1200000),
        width=Emu(8000000),
        height=Emu(4000000)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    # Center vertically
    from pptx.enum.text import MSO_ANCHOR
    tf.auto_size = None
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Añadir imagen"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    # 3. Find TextBox 1 and replace text
    for shape in slide5.shapes:
        if shape.name == 'TextBox 1':
            tf2 = shape.text_frame
            # Clear all paragraphs
            for para in tf2.paragraphs:
                for run in para.runs:
                    run.text = ''
            # Set first paragraph
            p2 = tf2.paragraphs[0]
            p2.clear()
            run2 = p2.add_run()
            run2.text = "Falta información"
            run2.font.size = Pt(16)
            run2.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
            run2.font.italic = True
            break

    print("  Fix 2 done.")


# ─────────────────────────────────────────────
# FIX 3: Slide 6 (index 5) — Replace 3 grouped boxes with 7 individual L2 boxes
# ─────────────────────────────────────────────
def fix3_slide6_l2_boxes(prs):
    print("Applying Fix 3: Slide 6 - Replace grouped boxes with 7 individual L2 boxes...")
    slide6 = prs.slides[5]

    # Remove shapes named 'Rectángulo 2' and the 3 arrow connectors
    to_remove = []
    for shape in slide6.shapes:
        if shape.name == 'Rectángulo 2':
            to_remove.append(shape)
        elif shape.name in ['Straight Arrow Connector 21', 'Straight Arrow Connector 22', 'Straight Arrow Connector 23']:
            to_remove.append(shape)

    for shape in to_remove:
        shape._element.getparent().remove(shape._element)
    print(f"  Removed {len(to_remove)} shapes.")

    # Layout constants
    SLIDE_WIDTH = 12192000
    SIDE_MARGIN = 400000
    ROW1_TOP = 3700000
    ROW1_HEIGHT = 1450000
    ROW2_TOP = 5250000
    ROW2_HEIGHT = 1450000
    GAP = 30000

    # Row 1: 4 boxes
    usable_width = SLIDE_WIDTH - 2 * SIDE_MARGIN
    row1_box_w = (usable_width - 3 * GAP) // 4  # 2723250

    # Row 2: 3 boxes
    row2_box_w = (usable_width - 2 * GAP) // 3  # 3644000

    # Box data: (title, [bullets])
    boxes_data = [
        # Row 1
        (
            "Dir. PF&A y Commercial Finance",
            [
                "Copiloto financiero del negocio, habilitando decisiones basadas en datos.",
                "Gestiona planeación financiera, presupuesto anual y previsiones continuas.",
                "Analiza rendimiento real vs. planificado e impulsa mejoras.",
                "Maximiza rentabilidad y brinda soporte a iniciativas estratégicas.",
            ]
        ),
        (
            "Gte Sr Fiscal y Contraloría",
            [
                "Define y dirige la estrategia fiscal y cumplimiento normativo.",
                "Gestiona auditorías, requerimientos de autoridades y control interno.",
                "Diseña sistemas de verificación y prevención de lavado de dinero.",
                "Garantiza alineación de políticas a IFRS y normas corporativas.",
            ]
        ),
        (
            "Gte Sr Servicios Compartidos",
            [
                "Lidera la operación centralizada de P2P, O2C y R2R.",
                "Define modelos de servicios centrales y esquemas de gestión.",
                "Garantiza cierres contables precisos y oportunos.",
                "Adopta mejores prácticas de servicios centralizados.",
            ]
        ),
        (
            "Gte Sr Legal y Regulatorio",
            [
                "Gestiona cumplimiento, litigios y gobernanza contractual.",
                "Diseña formatos de contratos y opiniones legales.",
                "Revisa contratos de expansión y capacita en riesgos legales.",
                "Genera opiniones para dar factibilidad a proyectos comerciales.",
            ]
        ),
        # Row 2
        (
            "Gte Asuntos Corporativos",
            [
                "Lidera la estrategia de relacionamiento externo y manejo de crisis.",
                "Mantiene contacto con actores externos e identifica riesgos regulatorios.",
                "Define la estrategia de comunicación y posicionamiento en medios.",
                "Gestiona proyectos con consultores en políticas públicas.",
            ]
        ),
        (
            "Dir. Adm. Nacional OXXO",
            [
                "Asegura procesos administrativos, financieros y de control en plazas.",
                "Garantiza el proceso de inventarios en tiendas mediante tecnología.",
                "Evalúa y aprueba inversiones en tiendas nuevas y remodeladas.",
                "Define el plan anual de normatividad basado en criticidad y riesgos.",
            ]
        ),
        (
            "Gte Sr Nac. Adm. CdS",
            [
                "Impulsa el buen desempeño de la cadena de suministro y distribución.",
                "Diseña KPIs y herramientas de gestión en CEDIS.",
                "Gestiona el capital de trabajo en conjunto con Comercial y Abasto.",
                "Asegura estabilidad normativa y control en la red de CdS.",
            ]
        ),
    ]

    for i, (title, bullets) in enumerate(boxes_data):
        if i < 4:  # Row 1
            left = SIDE_MARGIN + i * (row1_box_w + GAP)
            top = ROW1_TOP
            width = row1_box_w
            height = ROW1_HEIGHT
        else:  # Row 2
            j = i - 4
            left = SIDE_MARGIN + j * (row2_box_w + GAP)
            top = ROW2_TOP
            width = row2_box_w
            height = ROW2_HEIGHT

        # Add rectangle text box
        txBox = slide6.shapes.add_textbox(
            left=Emu(left),
            top=Emu(top),
            width=Emu(width),
            height=Emu(height)
        )

        # Set fill (light gray F2F2F2) and border (595959, 1pt) via XML
        sp = txBox._element
        spPr = get_spPr(txBox)

        # Remove any existing fills/noFill
        for nf in spPr.findall(f'{{{NS_A}}}noFill'):
            spPr.remove(nf)
        for sf in spPr.findall(f'{{{NS_A}}}solidFill'):
            spPr.remove(sf)
        # Add light gray fill
        solidFill = etree.SubElement(spPr, f'{{{NS_A}}}solidFill')
        srgbClr = etree.SubElement(solidFill, f'{{{NS_A}}}srgbClr')
        srgbClr.set('val', 'F2F2F2')

        # Add border
        ln = spPr.find(f'{{{NS_A}}}ln')
        if ln is None:
            ln = etree.SubElement(spPr, f'{{{NS_A}}}ln')
        ln.set('w', str(int(1 * 12700)))  # 1pt
        bf = ln.find(f'{{{NS_A}}}solidFill')
        if bf is None:
            bf = etree.SubElement(ln, f'{{{NS_A}}}solidFill')
        bc = bf.find(f'{{{NS_A}}}srgbClr')
        if bc is None:
            bc = etree.SubElement(bf, f'{{{NS_A}}}srgbClr')
        bc.set('val', '595959')

        # Populate text
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = None

        # Set some margin inside the box
        tf.margin_left = Emu(45720)
        tf.margin_right = Emu(45720)
        tf.margin_top = Emu(45720)
        tf.margin_bottom = Emu(45720)

        # Title paragraph
        p_title = tf.paragraphs[0]
        p_title.alignment = PP_ALIGN.LEFT
        run_title = p_title.add_run()
        run_title.text = title
        run_title.font.bold = True
        run_title.font.size = Pt(9)
        run_title.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

        # Bullet paragraphs
        for bullet in bullets:
            from pptx.oxml.ns import qn
            para = tf.add_paragraph()
            para.alignment = PP_ALIGN.LEFT
            run_b = para.add_run()
            run_b.text = "• " + bullet
            run_b.font.bold = False
            run_b.font.size = Pt(8)
            run_b.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    print("  Fix 3 done.")


# ─────────────────────────────────────────────
# FIX 4: RASCI Detail slides (indices 7 to 52)
# ─────────────────────────────────────────────
RASCI_COLORS = {
    'R': '70AD47',
    'A': 'FFC000',
    'A/R': 'FFC000',
    'R/A': '70AD47',
    'S': 'FF7070',
    'C': '4472C4',
    'I': 'A6A6A6',
}

BADGE_NAMES = {'TextBox 5', 'TextBox 7', 'TextBox 9', 'TextBox 11', 'TextBox 13', 'TextBox 15', 'TextBox 17', 'TextBox 19'}
DESC_NAMES = {'TextBox 6', 'TextBox 8', 'TextBox 10', 'TextBox 12', 'TextBox 14', 'TextBox 16', 'TextBox 18', 'TextBox 20'}


def set_text_run_color(shape, hex_color, font_size_pt=None, bold=None, italic=None, align=None):
    """Set all text runs in a shape to the specified color and optional size."""
    if not hasattr(shape, 'text_frame'):
        return
    tf = shape.text_frame
    for para in tf.paragraphs:
        if align is not None:
            para.alignment = align
        for run in para.runs:
            run.font.color.rgb = RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
            if font_size_pt is not None:
                run.font.size = Pt(font_size_pt)
            if bold is not None:
                run.font.bold = bold
            if italic is not None:
                run.font.italic = italic


def reposition_shape(shape, new_left=None, new_top=None, new_width=None, new_height=None):
    """Reposition/resize a shape."""
    spPr = get_spPr(shape)
    if spPr is None:
        return
    xfrm = spPr.find(f'{{{NS_A}}}xfrm')
    if xfrm is None:
        return
    off = xfrm.find(f'{{{NS_A}}}off')
    ext = xfrm.find(f'{{{NS_A}}}ext')
    if off is not None:
        if new_left is not None:
            off.set('x', str(int(new_left)))
        if new_top is not None:
            off.set('y', str(int(new_top)))
    if ext is not None:
        if new_width is not None:
            ext.set('cx', str(int(new_width)))
        if new_height is not None:
            ext.set('cy', str(int(new_height)))


def fix4_rasci_slides(prs):
    print("Applying Fix 4: RASCI slides (7-52) - Restyle badge/description rows...")
    count = 0

    for slide_idx in range(7, min(53, len(prs.slides))):
        slide = prs.slides[slide_idx]

        # Collect badge/desc shapes
        badge_shapes = {}  # name -> shape
        desc_shapes = {}   # name -> shape
        textbox4 = None

        for shape in slide.shapes:
            if shape.name == 'TextBox 4':
                textbox4 = shape
            elif shape.name in BADGE_NAMES:
                badge_shapes[shape.name] = shape
            elif shape.name in DESC_NAMES:
                desc_shapes[shape.name] = shape

        # 1. Remove TextBox 4
        if textbox4 is not None:
            textbox4._element.getparent().remove(textbox4._element)

        # 2. Determine number of roles
        n_roles = len(badge_shapes)
        if n_roles == 0:
            continue

        # Font sizing
        if n_roles <= 5:
            badge_font = 14
            desc_font = 11
        elif n_roles <= 7:
            badge_font = 12
            desc_font = 10
        else:
            badge_font = 10
            desc_font = 9

        # Layout
        available_height = 6858000 - 2050000 - 150000  # 4658000
        row_height = min(available_height // n_roles, 700000)
        BADGE_WIDTH = 640080
        DESC_WIDTH = 10728960
        START_TOP = 2100000

        # RASCI badge name ordering
        ordered_badges = ['TextBox 5', 'TextBox 7', 'TextBox 9', 'TextBox 11', 'TextBox 13', 'TextBox 15', 'TextBox 17', 'TextBox 19']
        ordered_descs = ['TextBox 6', 'TextBox 8', 'TextBox 10', 'TextBox 12', 'TextBox 14', 'TextBox 16', 'TextBox 18', 'TextBox 20']

        badge_left = 365760
        desc_left = 1097280

        for i, (bname, dname) in enumerate(zip(ordered_badges, ordered_descs)):
            if bname not in badge_shapes:
                break

            badge = badge_shapes[bname]
            desc = desc_shapes.get(dname)

            row_top = START_TOP + i * row_height

            # Reposition badge
            reposition_shape(badge,
                new_left=badge_left,
                new_top=row_top,
                new_width=BADGE_WIDTH,
                new_height=row_height
            )

            # Get RASCI letter from badge text
            badge_text = badge.text_frame.paragraphs[0].text.strip() if badge.text_frame.paragraphs else ''
            rasci_letter = badge_text.split('\n')[0].strip()
            color = RASCI_COLORS.get(rasci_letter, 'A6A6A6')

            # Set badge fill color
            set_shape_fill(badge, color)

            # Style badge text
            set_text_run_color(badge, 'FFFFFF', font_size_pt=badge_font, bold=True, align=PP_ALIGN.CENTER)

            # Reposition description
            if desc:
                reposition_shape(desc,
                    new_left=desc_left,
                    new_top=row_top,
                    new_width=DESC_WIDTH,
                    new_height=row_height
                )
                # Style description text
                set_text_run_color(desc, '000000', font_size_pt=desc_font)
                # Add colored left border to description
                add_border(desc, color, width_pt=3)

        count += 1

    print(f"  Fix 4 done. Processed {count} RASCI slides.")


# ─────────────────────────────────────────────
# Main execution
# ─────────────────────────────────────────────
def main():
    print(f"Loading: {PPTX_PATH}")
    prs = Presentation(PPTX_PATH)
    print(f"  Total slides: {len(prs.slides)}")

    fix1_slide4_text_black(prs)
    fix2_slide5_image_placeholder(prs)
    fix3_slide6_l2_boxes(prs)
    fix4_rasci_slides(prs)

    prs.save(PPTX_PATH)
    print(f"\nSaved to: {PPTX_PATH}")
    # Verify
    prs2 = Presentation(PPTX_PATH)
    print(f"Total slides after save: {len(prs2.slides)}")


if __name__ == '__main__':
    main()
