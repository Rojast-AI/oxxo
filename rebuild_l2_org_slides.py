#!/usr/bin/env python3
"""
Replace slide 6 in Finanzas Playbook with 7 L2 org chart slides (one per L2).
Each slide shows the L2 at top, L3 roles in boxes below, with description bullets.
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree
import copy

PPTX_PATH = '/home/user/oxxo/Finanzas/Playbook/Finanzas Playbook.pptx'
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'

# ─────────────────────────────────────────────────────────────
# L2 → L3 data extracted from Descripciones de Puesto files
# ─────────────────────────────────────────────────────────────
L2_DATA = [
    {
        'subtitle': 'PF&A y Commercial Finance',
        'l2_title': 'Gte Sr Planeación Financiera\ny Análisis & Commercial Finance',
        'l3s': [
            {
                'title': 'Gerente BP Finanzas\nPlataformas RH, TI',
                'bullets': [
                    'Elaborar y dar seguimiento al presupuesto anual y las previsiones de las plataformas RH y TI.',
                    'Coordinar el análisis de variaciones del desempeño financiero de RH y TI.',
                    'Generar información y asesoramiento sobre el impacto financiero de las decisiones de RH y TI.',
                    'Elaborar reportes financieros periódicos de las plataformas RH y TI.',
                ]
            },
            {
                'title': 'Gerente Commercial\nFinance',
                'bullets': [
                    'Elaborar el modelamiento de trade-offs económicos para respaldar decisiones comerciales.',
                    'Monitorear y analizar el rendimiento del margen bruto, la rentabilidad y sus drivers.',
                    'Identificar y reportar riesgos financieros con impacto potencial en la rentabilidad.',
                    'Elaborar reportes de visibilidad unificada de ingresos y márgenes para los socios comerciales.',
                ]
            },
            {
                'title': 'Gerente Planeación\nFinanciera',
                'bullets': [
                    'Coordinar el proceso de presupuestación anual y previsiones continuas de las funciones asignadas.',
                    'Elaborar el análisis de variaciones del desempeño financiero real versus lo planeado.',
                    'Desarrollar modelos financieros para la planeación, análisis de escenarios y pruebas de sensibilidad.',
                    'Generar los reportes financieros periódicos para los responsables de función y la dirección.',
                ]
            },
        ]
    },
    {
        'subtitle': 'Fiscal y Contraloría',
        'l2_title': 'Gte Sr Fiscal y Contraloría',
        'l3s': [
            {
                'title': 'Gerente Área\nNormatividad y Control',
                'bullets': [
                    'Diseñar e implementar el plan anual de normatividad y control interno.',
                    'Coordinar la ejecución de revisiones de control interno en las unidades de negocio.',
                    'Elaborar y reportar los resultados del sistema de control interno.',
                    'Homologar criterios y políticas contables alineados con FEMSA e IFRS.',
                ]
            },
            {
                'title': 'Gerente Control\nFiscal',
                'bullets': [
                    'Ejecutar el cumplimiento de las obligaciones fiscales en materia de impuestos federales, estatales y locales.',
                    'Implementar y mantener las herramientas de documentación y respaldo de los procesos fiscales.',
                    'Coordinar con Auditorías Fiscales la preparación de información para revisiones de autoridades.',
                    'Elaborar análisis de los cambios en la legislación fiscal y su aplicación en la operación.',
                ]
            },
            {
                'title': 'Gerente Gestión\nRiesgos SI',
                'bullets': [
                    'Diseñar e implementar la metodología de prevención de lavado de dinero y financiamiento al terrorismo.',
                    'Desarrollar y mantener los modelos de detección y reducción de fraude monetario.',
                    'Coordinar la implementación del sistema integrado de gestión de riesgos en las unidades de negocio.',
                    'Elaborar reportes de exposición al riesgo y efectividad de los controles de prevención.',
                ]
            },
            {
                'title': 'Gerente Sr\nContraloría',
                'bullets': [
                    'Supervisar la emisión de información financiera confiable y oportuna bajo normas IFRS.',
                    'Coordinar la representación del negocio ante auditoría interna y externa.',
                    'Implementar y monitorear el sistema de control interno orientado a la integridad de reportes financieros.',
                    'Coordinar la consolidación financiera de OXXO México.',
                ]
            },
            {
                'title': 'Responsable\nAuditorías Fiscales',
                'bullets': [
                    'Coordinar la atención y seguimiento a requerimientos de autoridades fiscales.',
                    'Ejecutar el seguimiento a las auditorías externas e internas en materia fiscal.',
                    'Elaborar análisis de cumplimiento de las obligaciones fiscales de OXXO México.',
                    'Coordinar la recuperación de saldos de impuestos a favor.',
                ]
            },
        ]
    },
    {
        'subtitle': 'Servicios Compartidos',
        'l2_title': 'Gte Sr Servicios Compartidos\nAdministrativo',
        'l3s': [
            {
                'title': 'Coordinador\nInformación y Gestión CSF',
                'bullets': [
                    'Coordinar la consolidación y reporte de los indicadores de desempeño del Centro de Servicios.',
                    'Dar seguimiento al cumplimiento de los acuerdos de nivel de servicio con los clientes internos.',
                    'Coordinar los rituales de gestión del Centro de Servicios.',
                    'Apoyar la gestión de la relación con los clientes internos del Centro de Servicios.',
                ]
            },
            {
                'title': 'Gerente Order\nto Cash',
                'bullets': [
                    'Ejecutar y supervisar el ciclo Order to Cash end-to-end.',
                    'Coordinar la gestión de la cartera de cuentas por cobrar.',
                    'Implementar y mantener los controles del proceso O2C.',
                    'Elaborar y reportar los indicadores de desempeño del proceso O2C.',
                ]
            },
            {
                'title': 'Gerente Purchase\nto Pay',
                'bullets': [
                    'Ejecutar y supervisar el ciclo Purchase to Pay end-to-end.',
                    'Coordinar la optimización de los flujos de aprobación, condiciones de pago y gestión de caja.',
                    'Implementar y mantener los controles del proceso P2P.',
                    'Gestionar la relación operativa con proveedores en el ámbito del proceso P2P.',
                ]
            },
            {
                'title': 'Gerente Record\nto Report',
                'bullets': [
                    'Ejecutar y supervisar el ciclo Record to Report end-to-end.',
                    'Coordinar el proceso de cierre contable mensual y anual.',
                    'Implementar y mantener los controles del proceso R2R.',
                    'Coordinar con Contraloría la alineación de criterios contables y el proceso de auditoría.',
                ]
            },
            {
                'title': 'Responsable\nGobierno Datos y Cond',
                'bullets': [
                    'Implementar y mantener los estándares de gobierno de datos del Centro de Servicios.',
                    'Coordinar el mantenimiento y actualización de los catálogos y condiciones comerciales.',
                    'Identificar y reportar desviaciones en la calidad de datos que impacten los procesos transaccionales.',
                    'Coordinar con las áreas usuarias la resolución de inconsistencias en datos y condiciones.',
                ]
            },
        ]
    },
    {
        'subtitle': 'Legal y Regulatorio',
        'l2_title': 'Gte Sr Legal y Regulatorio',
        'l3s': [
            {
                'title': 'Gerente Área\nCumplimiento',
                'bullets': [
                    'Diseñar e implementar el programa de cumplimiento legal y regulatorio de OXXO México.',
                    'Coordinar el monitoreo del cumplimiento de las obligaciones legales y regulatorias en las unidades de negocio.',
                    'Elaborar y reportar los indicadores de desempeño del programa de cumplimiento.',
                    'Coordinar con las áreas del negocio la implementación de acciones correctivas derivadas de hallazgos.',
                ]
            },
            {
                'title': 'Gerente Área\nLegal Comercial',
                'bullets': [
                    'Ejecutar la revisión y análisis legal de contratos y acuerdos comerciales de OXXO.',
                    'Gestionar las contingencias y litigios legales del área comercial.',
                    'Elaborar y actualizar los formatos de contratos comerciales de OXXO.',
                    'Asesorar al área comercial en procesos legales, negociaciones y contingencias.',
                ]
            },
            {
                'title': 'Gerente Legal CdS',
                'bullets': [
                    'Ejecutar la revisión y análisis legal de contratos y acuerdos de la Cadena de Suministro.',
                    'Gestionar las contingencias y litigios legales de la Cadena de Suministro.',
                    'Coordinar con el área de Gestoría la atención de trámites legales y regulatorios de CEDIS y Distribución.',
                    'Asesorar al área de Cadena de Suministro en procesos legales, negociaciones y contingencias.',
                ]
            },
            {
                'title': 'Gerente Sr Legal\nOperaciones',
                'bullets': [
                    'Ejecutar el análisis y atención de los requerimientos legales de la operación de tiendas y plazas.',
                    'Gestionar los litigios y procesos legales operativos de OXXO a nivel nacional.',
                    'Elaborar y revisar los documentos legales requeridos por la operación de tiendas.',
                    'Coordinar con los Gerentes de Plaza y Región la atención de contingencias legales operativas.',
                ]
            },
            {
                'title': 'Responsable\nLegal (i78)',
                'bullets': [
                    'Ejecutar el análisis y atención de los requerimientos legales de la plataforma i78.',
                    'Elaborar y revisar los documentos legales requeridos por la operación de i78.',
                    'Coordinar con las áreas operativas de i78 la atención de contingencias legales.',
                    'Dar seguimiento a los litigios y procesos legales activos de la plataforma i78.',
                ]
            },
        ]
    },
    {
        'subtitle': 'Asuntos Corporativos',
        'l2_title': 'Gte Sr Asuntos Corporativos',
        'l3s': [
            {
                'title': 'Analista\nComunicación Externa',
                'bullets': [
                    'Elaborar contenidos y materiales de comunicación externa para medios y audiencias clave.',
                    'Coordinar el monitoreo de medios de comunicación y redes sociales.',
                    'Elaborar reportes de monitoreo de medios y análisis de cobertura.',
                    'Apoyar la implementación de la estrategia de comunicación proactiva y reactiva del área.',
                ]
            },
            {
                'title': 'Coordinador\nAsuntos Corporativos',
                'bullets': [
                    'Elaborar análisis de riesgo e informes políticos basados en tendencias regulatorias del entorno.',
                    'Coordinar el monitoreo sistemático del entorno político, económico y social.',
                    'Elaborar reportes de la agenda de riesgos del entorno para la toma de decisiones.',
                    'Coordinar la gestión de proyectos con consultores especializados en políticas públicas.',
                ]
            },
            {
                'title': 'Gerente Área\nAtención Autoridades',
                'bullets': [
                    'Gestionar y dar seguimiento a los requerimientos de autoridades gubernamentales hacia OXXO.',
                    'Elaborar las respuestas y documentación requerida ante requerimientos de autoridades.',
                    'Coordinar con las áreas operativas y de soporte la atención de requerimientos de autoridades.',
                    'Monitorear y reportar el estatus de los requerimientos de autoridades en curso.',
                ]
            },
            {
                'title': 'Gerente Área\nRelaciones Gobierno',
                'bullets': [
                    'Ejecutar los planes de relacionamiento con autoridades gubernamentales en la región o ámbito asignado.',
                    'Monitorear el entorno político-regulatorio local y reportar tendencias relevantes para el negocio.',
                    'Participar en organismos empresariales y foros de influencia en representación de OXXO.',
                    'Ejecutar proyectos de cabildeo e incidencia en política pública en el ámbito asignado.',
                ]
            },
            {
                'title': 'Gerente\nRelaciones Gobierno',
                'bullets': [
                    'Coordinar la implementación de la estrategia de relacionamiento con autoridades gubernamentales.',
                    'Elaborar y dar seguimiento a los planes de relacionamiento con autoridades y actores clave.',
                    'Coordinar el monitoreo del entorno regulatorio y político relevante para el negocio.',
                    'Gestionar proyectos de cabildeo y relacionamiento institucional con consultores especializados.',
                ]
            },
        ]
    },
    {
        'subtitle': 'Administración Nacional OXXO',
        'l2_title': 'Director Nacional\nAdmvo OXXO',
        'l3s': [
            {
                'title': 'Gerente Área\nAdmvo Neg Dllo',
                'bullets': [
                    'Implementar los procesos administrativos y financieros de los negocios en desarrollo.',
                    'Coordinar la generación y entrega de información financiera de los negocios en desarrollo.',
                    'Elaborar el seguimiento presupuestal de los negocios en desarrollo.',
                    'Ejecutar los controles de normatividad y control interno aplicables a los negocios en desarrollo.',
                ]
            },
            {
                'title': 'Gerente Área BP\nFinanzas Plataformas Expansión',
                'bullets': [
                    'Elaborar los modelos de evaluación financiera de proyectos de expansión.',
                    'Coordinar el seguimiento presupuestal de los proyectos de expansión en curso.',
                    'Generar reportes financieros periódicos de la cartera de expansión.',
                    'Implementar los controles administrativos y financieros definidos para los procesos de expansión.',
                ]
            },
            {
                'title': 'Gerente Gestión\nControl Zona',
                'bullets': [
                    'Ejecutar el plan de normatividad y control interno en las plazas de la zona asignada.',
                    'Implementar y dar seguimiento a los procedimientos contables y administrativos en tiendas de la zona.',
                    'Generar y reportar los indicadores de gestión y control de la zona.',
                    'Coordinar con los Gerentes de Plaza y Región la implementación de acciones correctivas.',
                ]
            },
            {
                'title': 'Responsable\nGestión y Control',
                'bullets': [
                    'Ejecutar los procesos de control administrativo y financiero en las plazas asignadas.',
                    'Generar y consolidar los reportes de gestión e indicadores operativos de las plazas.',
                    'Coordinar la ejecución de revisiones de control interno en plazas y tiendas.',
                    'Implementar y dar seguimiento a las herramientas de gestión y KPIs en las plazas.',
                ]
            },
            {
                'title': 'Responsable\nTraslado de Valores',
                'bullets': [
                    'Coordinar la operación del modelo de traslado y resguardo de valores en las tiendas.',
                    'Implementar y dar seguimiento al modelo I-Cash y su caso de negocio.',
                    'Coordinar la relación con los proveedores de servicios de traslado de valores.',
                    'Elaborar reportes de desempeño y costos del modelo de traslado de valores.',
                ]
            },
        ]
    },
    {
        'subtitle': 'Administración Nacional CdS',
        'l2_title': 'Gte Sr Nac. Administrativo\nCentros de Suministro',
        'l3s': [
            {
                'title': 'Gerente Área\nAdmin Control CEDIS',
                'bullets': [
                    'Implementar y mantener los procedimientos contables y administrativos en CEDIS.',
                    'Ejecutar el plan anual de normatividad y control interno en CEDIS.',
                    'Coordinar la generación de información financiera y operativa de CEDIS.',
                    'Elaborar y dar seguimiento a los indicadores de productividad de CEDIS.',
                ]
            },
            {
                'title': 'Gerente Área Admin\ny Finanzas Dist',
                'bullets': [
                    'Implementar y mantener los procedimientos contables y administrativos en Distribución.',
                    'Coordinar la generación de información financiera de rentabilidad en Distribución.',
                    'Elaborar y dar seguimiento a los indicadores de productividad de Distribución.',
                    'Ejecutar el plan de normatividad y control interno en Distribución.',
                ]
            },
            {
                'title': 'Responsable Gestión\ny Control CEDIS',
                'bullets': [
                    'Ejecutar los procesos de control administrativo en las unidades de CEDIS asignadas.',
                    'Generar y consolidar los reportes de gestión e indicadores de productividad de CEDIS.',
                    'Coordinar la ejecución de revisiones de control interno en CEDIS.',
                    'Implementar los mecanismos de seguimiento al capital de trabajo en CEDIS.',
                ]
            },
            {
                'title': 'Responsable\nGestoría',
                'bullets': [
                    'Ejecutar y dar seguimiento a los trámites y gestiones administrativas de CEDIS y Distribución.',
                    'Coordinar la gestión documental y archivo de expedientes normativos y regulatorios.',
                    'Identificar y reportar vencimientos y renovaciones de permisos operativos de CEDIS y Distribución.',
                    'Coordinar con el área Legal los requerimientos de gestoría que impliquen procesos jurídicos.',
                ]
            },
        ]
    },
]


# ─────────────────────────────────────────────────────────────
# XML helpers
# ─────────────────────────────────────────────────────────────
def get_spPr(shape):
    sp = shape._element
    spPr = sp.find(f'{{{NS_P}}}spPr')
    if spPr is not None:
        return spPr
    return sp.find(f'.//{{{NS_A}}}spPr')


def set_fill(shape, hex_color=None, no_fill=False):
    spPr = get_spPr(shape)
    if spPr is None:
        return
    for tag in [f'{{{NS_A}}}noFill', f'{{{NS_A}}}solidFill']:
        for el in spPr.findall(tag):
            spPr.remove(el)
    if no_fill:
        etree.SubElement(spPr, f'{{{NS_A}}}noFill')
    elif hex_color:
        sf = etree.SubElement(spPr, f'{{{NS_A}}}solidFill')
        clr = etree.SubElement(sf, f'{{{NS_A}}}srgbClr')
        clr.set('val', hex_color)


def set_border(shape, hex_color='000000', width_pt=1):
    spPr = get_spPr(shape)
    if spPr is None:
        return
    ln = spPr.find(f'{{{NS_A}}}ln')
    if ln is None:
        ln = etree.SubElement(spPr, f'{{{NS_A}}}ln')
    ln.set('w', str(int(width_pt * 12700)))
    sf = ln.find(f'{{{NS_A}}}solidFill')
    if sf is None:
        sf = etree.SubElement(ln, f'{{{NS_A}}}solidFill')
    clr = sf.find(f'{{{NS_A}}}srgbClr')
    if clr is None:
        clr = etree.SubElement(sf, f'{{{NS_A}}}srgbClr')
    clr.set('val', hex_color)


def set_no_border(shape):
    spPr = get_spPr(shape)
    if spPr is None:
        return
    ln = spPr.find(f'{{{NS_A}}}ln')
    if ln is None:
        ln = etree.SubElement(spPr, f'{{{NS_A}}}ln')
    for sf in ln.findall(f'{{{NS_A}}}solidFill'):
        ln.remove(sf)
    etree.SubElement(ln, f'{{{NS_A}}}noFill')


def add_text_box(slide, left, top, width, height, text, font_size, bold=False,
                 italic=False, color='000000', align=PP_ALIGN.LEFT,
                 word_wrap=True, v_anchor=MSO_ANCHOR.TOP, margin=45720):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    tf.margin_left = Emu(margin)
    tf.margin_right = Emu(margin)
    tf.margin_top = Emu(margin // 2)
    tf.margin_bottom = Emu(margin // 2)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = RGBColor(
        int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
    )
    return txBox


def add_line(slide, x1, y1, length, vertical=False, color='000000', width_pt=0.75):
    """Add a thin rectangle as a connector line."""
    if vertical:
        shape = slide.shapes.add_textbox(Emu(x1 - 6350), Emu(y1), Emu(12700), Emu(length))
    else:
        shape = slide.shapes.add_textbox(Emu(x1), Emu(y1 - 6350), Emu(length), Emu(12700))
    set_fill(shape, hex_color=color)
    set_no_border(shape)
    return shape


def add_box_with_text(slide, left, top, width, height, title_lines, bullets=None,
                      title_size=9, bullet_size=8, fill='FFFFFF', border='000000'):
    """Add a role box: title (bold) + optional bullet points."""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    set_fill(txBox, hex_color=fill)
    set_border(txBox, hex_color=border, width_pt=0.75)

    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Emu(45720)
    tf.margin_right = Emu(45720)
    tf.margin_top = Emu(30000)
    tf.margin_bottom = Emu(30000)

    # Title paragraph(s)
    first_para = True
    for line in title_lines.split('\n'):
        if first_para:
            p = tf.paragraphs[0]
            first_para = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.bold = True
        run.font.size = Pt(title_size)
        run.font.color.rgb = RGBColor(0, 0, 0)

    # Bullet paragraphs
    if bullets:
        for bullet in bullets:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = '• ' + bullet
            run.font.bold = False
            run.font.size = Pt(bullet_size)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    return txBox


# ─────────────────────────────────────────────────────────────
# Build one L2 org chart slide
# ─────────────────────────────────────────────────────────────
def build_l2_slide(prs, l2_info):
    layout = prs.slide_layouts[11]  # 1_Title Only
    slide = prs.slides.add_slide(layout)

    SLIDE_W = 12192000
    SLIDE_H = 6858000
    SIDE_MARGIN = 304800

    subtitle_text = l2_info['subtitle']
    l2_title = l2_info['l2_title']
    l3s = l2_info['l3s']
    N = len(l3s)

    # ── Title placeholder ──────────────────────────────────────
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = 'Estructura Organizacional'
            run.font.bold = True
            run.font.size = Pt(22)
            run.font.color.rgb = RGBColor(0, 0, 0)
            break

    # ── Subtitle (red italic) ──────────────────────────────────
    sub_box = add_text_box(
        slide, left=330200, top=530000, width=10000000, height=300000,
        text=subtitle_text, font_size=18, italic=True, color='CC0000',
        align=PP_ALIGN.LEFT
    )
    set_fill(sub_box, no_fill=True)
    set_no_border(sub_box)

    # ── L2 box (centered) ──────────────────────────────────────
    L2_W = 3200000
    L2_H = 520000
    L2_LEFT = (SLIDE_W - L2_W) // 2
    L2_TOP = 920000
    l2_box = add_box_with_text(
        slide, L2_LEFT, L2_TOP, L2_W, L2_H,
        title_lines=l2_title,
        title_size=10, fill='FFFFFF', border='000000'
    )

    # ── Connectors ────────────────────────────────────────────
    L2_BOTTOM = L2_TOP + L2_H
    L2_CENTER_X = L2_LEFT + L2_W // 2

    # Layout for L3 boxes
    USABLE_W = SLIDE_W - 2 * SIDE_MARGIN
    GAP = 60000 if N <= 4 else 40000
    L3_W = (USABLE_W - (N - 1) * GAP) // N
    L3_TOP = 1780000
    L3_H = 680000

    # Vertical line from L2 bottom to horizontal bar
    HBAR_Y = L3_TOP - 120000
    vert_len = HBAR_Y - L2_BOTTOM
    add_line(slide, L2_CENTER_X, L2_BOTTOM, vert_len, vertical=True)

    # Horizontal bar
    L3_CENTERS = [SIDE_MARGIN + i * (L3_W + GAP) + L3_W // 2 for i in range(N)]
    H_LEFT = L3_CENTERS[0]
    H_RIGHT = L3_CENTERS[-1]
    if N > 1:
        add_line(slide, H_LEFT, HBAR_Y, H_RIGHT - H_LEFT, vertical=False)

    # Vertical drops from bar to each L3 box
    DROP_LEN = L3_TOP - HBAR_Y
    for cx in L3_CENTERS:
        add_line(slide, cx, HBAR_Y, DROP_LEN, vertical=True)

    # ── L3 boxes ──────────────────────────────────────────────
    DESC_TOP = L3_TOP + L3_H + 40000
    DESC_H = SLIDE_H - DESC_TOP - 80000

    # Font sizes based on content
    if N <= 3:
        b_title_size = 10
        b_bullet_size = 9
    elif N <= 4:
        b_title_size = 10
        b_bullet_size = 8
    else:
        b_title_size = 9
        b_bullet_size = 7

    for i, l3 in enumerate(l3s):
        L3_LEFT = SIDE_MARGIN + i * (L3_W + GAP)

        # L3 header box
        add_box_with_text(
            slide, L3_LEFT, L3_TOP, L3_W, L3_H,
            title_lines=l3['title'],
            title_size=b_title_size,
            fill='FFFFFF', border='000000'
        )

        # Description box with bullets
        desc_box = slide.shapes.add_textbox(
            Emu(L3_LEFT), Emu(DESC_TOP), Emu(L3_W), Emu(DESC_H)
        )
        set_fill(desc_box, hex_color='F5F5F5')
        set_border(desc_box, hex_color='BFBFBF', width_pt=0.5)

        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        tf.margin_left = Emu(45720)
        tf.margin_right = Emu(45720)
        tf.margin_top = Emu(45720)
        tf.margin_bottom = Emu(30000)

        first = True
        for bullet in l3['bullets']:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            # Space between bullets via paragraph spacing
            pPr = p._p.get_or_add_pPr()
            spcBef = etree.SubElement(pPr, f'{{{NS_A}}}spcBef')
            spcPts = etree.SubElement(spcBef, f'{{{NS_A}}}spcPts')
            spcPts.set('val', '100')  # 10pt space before each bullet
            run = p.add_run()
            run.text = '• ' + bullet
            run.font.size = Pt(b_bullet_size)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    print(f'  Built slide: {subtitle_text} ({N} L3s)')
    return slide


# ─────────────────────────────────────────────────────────────
# Delete a slide by index
# ─────────────────────────────────────────────────────────────
def delete_slide(prs, slide_idx):
    xml_slides = prs.slides._sldIdLst
    slide = prs.slides[slide_idx]
    # Remove from sldIdLst
    rId_elem = xml_slides[slide_idx]
    xml_slides.remove(rId_elem)
    print(f'  Deleted slide at index {slide_idx}')


# ─────────────────────────────────────────────────────────────
# Move a slide from old_idx to new_idx
# ─────────────────────────────────────────────────────────────
def move_slide(prs, old_idx, new_idx):
    xml_slides = prs.slides._sldIdLst
    sldId = xml_slides[old_idx]
    xml_slides.remove(sldId)
    xml_slides.insert(new_idx, sldId)


# ─────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────
def main():
    print(f'Loading: {PPTX_PATH}')
    prs = Presentation(PPTX_PATH)
    print(f'  Slides before: {len(prs.slides)}')

    # Delete current slide 6 (index 5) - the 7-box L2 structure slide
    print('Deleting old slide 6 (index 5)...')
    delete_slide(prs, 5)
    print(f'  Slides after deletion: {len(prs.slides)}')

    # Build 7 new L2 org chart slides (appended at the end)
    print('Building 7 new L2 org chart slides...')
    for l2_info in L2_DATA:
        build_l2_slide(prs, l2_info)

    total = len(prs.slides)
    print(f'  Slides after building new: {total}')

    # Move the 7 new slides from the end to position 6 (index 5)
    # They were appended as indices (total-7) to (total-1)
    # Move them one by one to index 5, 6, 7, 8, 9, 10, 11
    print('Repositioning new slides to position 6...')
    for i in range(7):
        # After each move, the next slide to move is still at (total-7+i) - i
        # But since we inserted at position (5+i), the tail shifts
        current_tail_idx = len(prs.slides) - (7 - i)
        move_slide(prs, current_tail_idx, 5 + i)
        print(f'  Moved slide to position {5 + i + 1}')

    print(f'  Final slide count: {len(prs.slides)}')

    prs.save(PPTX_PATH)
    print(f'\nSaved: {PPTX_PATH}')

    # Verify
    prs2 = Presentation(PPTX_PATH)
    print(f'Verification - total slides: {len(prs2.slides)}')
    for i in range(min(15, len(prs2.slides))):
        s = prs2.slides[i]
        title_txt = ''
        for shape in s.shapes:
            if shape.name == 'Title 1' and hasattr(shape, 'text_frame'):
                title_txt = shape.text_frame.text[:60]
                break
        print(f'  Slide {i+1}: {title_txt!r}')


if __name__ == '__main__':
    main()
