#!/usr/bin/env python3
"""
Rebuild L2 org chart slides without deleting any slides.
Strategy:
  - Clear slide 6 content in-place (no ZIP change, just XML rewrite)
  - Add 6 new slides via add_slide() → clean filenames, no conflicts
  - Reorder new slides to positions 7-12
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import copy, zipfile, re

SRC   = '/tmp/Finanzas Playbook_working.pptx'  # last known good (62 slides)
DEST  = '/home/user/oxxo/Finanzas/Playbook/Finanzas Playbook.pptx'

NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'

# ─────────────────────────────────────────────────────────────
# L2 → L3 data
# ─────────────────────────────────────────────────────────────
L2_DATA = [
    {
        'subtitle': 'PF&A y Commercial Finance',
        'l2_title': 'Gte Sr Planeación Financiera\ny Análisis & Commercial Finance',
        'l3s': [
            {'title': 'Gerente BP Finanzas\nPlataformas RH, TI',
             'bullets': ['Elaborar y dar seguimiento al presupuesto anual y las previsiones de las plataformas RH y TI.',
                         'Coordinar el análisis de variaciones del desempeño financiero de RH y TI.',
                         'Generar información y asesoramiento sobre el impacto financiero de las decisiones de RH y TI.',
                         'Elaborar reportes financieros periódicos de las plataformas RH y TI.']},
            {'title': 'Gerente Commercial\nFinance',
             'bullets': ['Elaborar el modelamiento de trade-offs económicos para respaldar decisiones comerciales.',
                         'Monitorear y analizar el rendimiento del margen bruto, la rentabilidad y sus drivers.',
                         'Identificar y reportar riesgos financieros con impacto potencial en la rentabilidad.',
                         'Elaborar reportes de visibilidad unificada de ingresos y márgenes para los socios comerciales.']},
            {'title': 'Gerente Planeación\nFinanciera',
             'bullets': ['Coordinar el proceso de presupuestación anual y previsiones continuas de las funciones asignadas.',
                         'Elaborar el análisis de variaciones del desempeño financiero real versus lo planeado.',
                         'Desarrollar modelos financieros para la planeación, análisis de escenarios y pruebas de sensibilidad.',
                         'Generar los reportes financieros periódicos para los responsables de función y la dirección.']},
        ]
    },
    {
        'subtitle': 'Fiscal y Contraloría',
        'l2_title': 'Gte Sr Fiscal y Contraloría',
        'l3s': [
            {'title': 'Gerente Área\nNormatividad y Control',
             'bullets': ['Diseñar e implementar el plan anual de normatividad y control interno.',
                         'Coordinar la ejecución de revisiones de control interno en las unidades de negocio.',
                         'Elaborar y reportar los resultados del sistema de control interno.',
                         'Homologar criterios y políticas contables alineados con FEMSA e IFRS.']},
            {'title': 'Gerente Control\nFiscal',
             'bullets': ['Ejecutar el cumplimiento de las obligaciones fiscales en materia de impuestos federales, estatales y locales.',
                         'Implementar y mantener las herramientas de documentación y respaldo de los procesos fiscales.',
                         'Coordinar con Auditorías Fiscales la preparación de información para revisiones de autoridades.',
                         'Elaborar análisis de los cambios en la legislación fiscal y su aplicación en la operación.']},
            {'title': 'Gerente Gestión\nRiesgos SI',
             'bullets': ['Diseñar e implementar la metodología de prevención de lavado de dinero y financiamiento al terrorismo.',
                         'Desarrollar y mantener los modelos de detección y reducción de fraude monetario.',
                         'Coordinar la implementación del sistema integrado de gestión de riesgos en las unidades de negocio.',
                         'Elaborar reportes de exposición al riesgo y efectividad de los controles de prevención.']},
            {'title': 'Gerente Sr\nContraloría',
             'bullets': ['Supervisar la emisión de información financiera confiable y oportuna bajo normas IFRS.',
                         'Coordinar la representación del negocio ante auditoría interna y externa.',
                         'Implementar y monitorear el sistema de control interno orientado a la integridad de reportes financieros.',
                         'Coordinar la consolidación financiera de OXXO México.']},
            {'title': 'Responsable\nAuditorías Fiscales',
             'bullets': ['Coordinar la atención y seguimiento a requerimientos de autoridades fiscales.',
                         'Ejecutar el seguimiento a las auditorías externas e internas en materia fiscal.',
                         'Elaborar análisis de cumplimiento de las obligaciones fiscales de OXXO México.',
                         'Coordinar la recuperación de saldos de impuestos a favor.']},
        ]
    },
    {
        'subtitle': 'Servicios Compartidos',
        'l2_title': 'Gte Sr Servicios Compartidos\nAdministrativo',
        'l3s': [
            {'title': 'Coordinador\nInformación y Gestión CSF',
             'bullets': ['Coordinar la consolidación y reporte de los indicadores de desempeño del Centro de Servicios.',
                         'Dar seguimiento al cumplimiento de los acuerdos de nivel de servicio con los clientes internos.',
                         'Coordinar los rituales de gestión del Centro de Servicios.',
                         'Apoyar la gestión de la relación con los clientes internos del Centro de Servicios.']},
            {'title': 'Gerente Order\nto Cash',
             'bullets': ['Ejecutar y supervisar el ciclo Order to Cash end-to-end.',
                         'Coordinar la gestión de la cartera de cuentas por cobrar.',
                         'Implementar y mantener los controles del proceso O2C.',
                         'Elaborar y reportar los indicadores de desempeño del proceso O2C.']},
            {'title': 'Gerente Purchase\nto Pay',
             'bullets': ['Ejecutar y supervisar el ciclo Purchase to Pay end-to-end.',
                         'Coordinar la optimización de los flujos de aprobación, condiciones de pago y gestión de caja.',
                         'Implementar y mantener los controles del proceso P2P.',
                         'Gestionar la relación operativa con proveedores en el ámbito del proceso P2P.']},
            {'title': 'Gerente Record\nto Report',
             'bullets': ['Ejecutar y supervisar el ciclo Record to Report end-to-end.',
                         'Coordinar el proceso de cierre contable mensual y anual.',
                         'Implementar y mantener los controles del proceso R2R.',
                         'Coordinar con Contraloría la alineación de criterios contables y el proceso de auditoría.']},
            {'title': 'Responsable\nGobierno Datos y Cond',
             'bullets': ['Implementar y mantener los estándares de gobierno de datos del Centro de Servicios.',
                         'Coordinar el mantenimiento y actualización de los catálogos y condiciones comerciales.',
                         'Identificar y reportar desviaciones en la calidad de datos que impacten los procesos transaccionales.',
                         'Coordinar con las áreas usuarias la resolución de inconsistencias en datos y condiciones.']},
        ]
    },
    {
        'subtitle': 'Legal y Regulatorio',
        'l2_title': 'Gte Sr Legal y Regulatorio',
        'l3s': [
            {'title': 'Gerente Área\nCumplimiento',
             'bullets': ['Diseñar e implementar el programa de cumplimiento legal y regulatorio de OXXO México.',
                         'Coordinar el monitoreo del cumplimiento de las obligaciones legales y regulatorias en las unidades de negocio.',
                         'Elaborar y reportar los indicadores de desempeño del programa de cumplimiento.',
                         'Coordinar con las áreas del negocio la implementación de acciones correctivas derivadas de hallazgos.']},
            {'title': 'Gerente Área\nLegal Comercial',
             'bullets': ['Ejecutar la revisión y análisis legal de contratos y acuerdos comerciales de OXXO.',
                         'Gestionar las contingencias y litigios legales del área comercial.',
                         'Elaborar y actualizar los formatos de contratos comerciales de OXXO.',
                         'Asesorar al área comercial en procesos legales, negociaciones y contingencias.']},
            {'title': 'Gerente Legal CdS',
             'bullets': ['Ejecutar la revisión y análisis legal de contratos y acuerdos de la Cadena de Suministro.',
                         'Gestionar las contingencias y litigios legales de la Cadena de Suministro.',
                         'Coordinar con el área de Gestoría la atención de trámites legales y regulatorios de CEDIS y Distribución.',
                         'Asesorar al área de Cadena de Suministro en procesos legales, negociaciones y contingencias.']},
            {'title': 'Gerente Sr Legal\nOperaciones',
             'bullets': ['Ejecutar el análisis y atención de los requerimientos legales de la operación de tiendas y plazas.',
                         'Gestionar los litigios y procesos legales operativos de OXXO a nivel nacional.',
                         'Elaborar y revisar los documentos legales requeridos por la operación de tiendas.',
                         'Coordinar con los Gerentes de Plaza y Región la atención de contingencias legales operativas.']},
            {'title': 'Responsable\nLegal (i78)',
             'bullets': ['Ejecutar el análisis y atención de los requerimientos legales de la plataforma i78.',
                         'Elaborar y revisar los documentos legales requeridos por la operación de i78.',
                         'Coordinar con las áreas operativas de i78 la atención de contingencias legales.',
                         'Dar seguimiento a los litigios y procesos legales activos de la plataforma i78.']},
        ]
    },
    {
        'subtitle': 'Asuntos Corporativos',
        'l2_title': 'Gte Sr Asuntos Corporativos',
        'l3s': [
            {'title': 'Analista\nComunicación Externa',
             'bullets': ['Elaborar contenidos y materiales de comunicación externa para medios y audiencias clave.',
                         'Coordinar el monitoreo de medios de comunicación y redes sociales.',
                         'Elaborar reportes de monitoreo de medios y análisis de cobertura.',
                         'Apoyar la implementación de la estrategia de comunicación proactiva y reactiva del área.']},
            {'title': 'Coordinador\nAsuntos Corporativos',
             'bullets': ['Elaborar análisis de riesgo e informes políticos basados en tendencias regulatorias del entorno.',
                         'Coordinar el monitoreo sistemático del entorno político, económico y social.',
                         'Elaborar reportes de la agenda de riesgos del entorno para la toma de decisiones.',
                         'Coordinar la gestión de proyectos con consultores especializados en políticas públicas.']},
            {'title': 'Gerente Área\nAtención Autoridades',
             'bullets': ['Gestionar y dar seguimiento a los requerimientos de autoridades gubernamentales hacia OXXO.',
                         'Elaborar las respuestas y documentación requerida ante requerimientos de autoridades.',
                         'Coordinar con las áreas operativas y de soporte la atención de requerimientos de autoridades.',
                         'Monitorear y reportar el estatus de los requerimientos de autoridades en curso.']},
            {'title': 'Gerente Área\nRelaciones Gobierno',
             'bullets': ['Ejecutar los planes de relacionamiento con autoridades gubernamentales en la región o ámbito asignado.',
                         'Monitorear el entorno político-regulatorio local y reportar tendencias relevantes para el negocio.',
                         'Participar en organismos empresariales y foros de influencia en representación de OXXO.',
                         'Ejecutar proyectos de cabildeo e incidencia en política pública en el ámbito asignado.']},
            {'title': 'Gerente\nRelaciones Gobierno',
             'bullets': ['Coordinar la implementación de la estrategia de relacionamiento con autoridades gubernamentales.',
                         'Elaborar y dar seguimiento a los planes de relacionamiento con autoridades y actores clave.',
                         'Coordinar el monitoreo del entorno regulatorio y político relevante para el negocio.',
                         'Gestionar proyectos de cabildeo y relacionamiento institucional con consultores especializados.']},
        ]
    },
    {
        'subtitle': 'Administración Nacional OXXO',
        'l2_title': 'Director Nacional\nAdmvo OXXO',
        'l3s': [
            {'title': 'Gerente Área\nAdmvo Neg Dllo',
             'bullets': ['Implementar los procesos administrativos y financieros de los negocios en desarrollo.',
                         'Coordinar la generación y entrega de información financiera de los negocios en desarrollo.',
                         'Elaborar el seguimiento presupuestal de los negocios en desarrollo.',
                         'Ejecutar los controles de normatividad y control interno aplicables a los negocios en desarrollo.']},
            {'title': 'Gerente Área BP\nFinanzas Plataformas Expansión',
             'bullets': ['Elaborar los modelos de evaluación financiera de proyectos de expansión.',
                         'Coordinar el seguimiento presupuestal de los proyectos de expansión en curso.',
                         'Generar reportes financieros periódicos de la cartera de expansión.',
                         'Implementar los controles administrativos y financieros para los procesos de expansión.']},
            {'title': 'Gerente Gestión\nControl Zona',
             'bullets': ['Ejecutar el plan de normatividad y control interno en las plazas de la zona asignada.',
                         'Implementar y dar seguimiento a los procedimientos contables y administrativos en tiendas de la zona.',
                         'Generar y reportar los indicadores de gestión y control de la zona.',
                         'Coordinar con los Gerentes de Plaza y Región la implementación de acciones correctivas.']},
            {'title': 'Responsable\nGestión y Control',
             'bullets': ['Ejecutar los procesos de control administrativo y financiero en las plazas asignadas.',
                         'Generar y consolidar los reportes de gestión e indicadores operativos de las plazas.',
                         'Coordinar la ejecución de revisiones de control interno en plazas y tiendas.',
                         'Implementar y dar seguimiento a las herramientas de gestión y KPIs en las plazas.']},
            {'title': 'Responsable\nTraslado de Valores',
             'bullets': ['Coordinar la operación del modelo de traslado y resguardo de valores en las tiendas.',
                         'Implementar y dar seguimiento al modelo I-Cash y su caso de negocio.',
                         'Coordinar la relación con los proveedores de servicios de traslado de valores.',
                         'Elaborar reportes de desempeño y costos del modelo de traslado de valores.']},
        ]
    },
    {
        'subtitle': 'Administración Nacional CdS',
        'l2_title': 'Gte Sr Nac. Administrativo\nCentros de Suministro',
        'l3s': [
            {'title': 'Gerente Área\nAdmin Control CEDIS',
             'bullets': ['Implementar y mantener los procedimientos contables y administrativos en CEDIS.',
                         'Ejecutar el plan anual de normatividad y control interno en CEDIS.',
                         'Coordinar la generación de información financiera y operativa de CEDIS.',
                         'Elaborar y dar seguimiento a los indicadores de productividad de CEDIS.']},
            {'title': 'Gerente Área Admin\ny Finanzas Dist',
             'bullets': ['Implementar y mantener los procedimientos contables y administrativos en Distribución.',
                         'Coordinar la generación de información financiera de rentabilidad en Distribución.',
                         'Elaborar y dar seguimiento a los indicadores de productividad de Distribución.',
                         'Ejecutar el plan de normatividad y control interno en Distribución.']},
            {'title': 'Responsable Gestión\ny Control CEDIS',
             'bullets': ['Ejecutar los procesos de control administrativo en las unidades de CEDIS asignadas.',
                         'Generar y consolidar los reportes de gestión e indicadores de productividad de CEDIS.',
                         'Coordinar la ejecución de revisiones de control interno en CEDIS.',
                         'Implementar los mecanismos de seguimiento al capital de trabajo en CEDIS.']},
            {'title': 'Responsable\nGestoría',
             'bullets': ['Ejecutar y dar seguimiento a los trámites y gestiones administrativas de CEDIS y Distribución.',
                         'Coordinar la gestión documental y archivo de expedientes normativos y regulatorios.',
                         'Identificar y reportar vencimientos y renovaciones de permisos operativos de CEDIS y Distribución.',
                         'Coordinar con el área Legal los requerimientos de gestoría que impliquen procesos jurídicos.']},
        ]
    },
]

# ─────────────────────────────────────────────────────────────
# XML helpers
# ─────────────────────────────────────────────────────────────
def get_spPr(shape):
    sp = shape._element
    spPr = sp.find(f'{{{NS_P}}}spPr')
    if spPr is not None:
        return spPr
    return sp.find(f'.//{{{NS_A}}}spPr')

def set_fill(spPr, hex_color=None, no_fill=False):
    for tag in [f'{{{NS_A}}}noFill', f'{{{NS_A}}}solidFill']:
        for el in spPr.findall(tag):
            spPr.remove(el)
    if no_fill:
        etree.SubElement(spPr, f'{{{NS_A}}}noFill')
    elif hex_color:
        sf = etree.SubElement(spPr, f'{{{NS_A}}}solidFill')
        clr = etree.SubElement(sf, f'{{{NS_A}}}srgbClr')
        clr.set('val', hex_color)

def set_border(spPr, hex_color='000000', width_pt=0.75):
    ln = spPr.find(f'{{{NS_A}}}ln')
    if ln is None:
        ln = etree.SubElement(spPr, f'{{{NS_A}}}ln')
    ln.set('w', str(int(width_pt * 12700)))
    for sf in ln.findall(f'{{{NS_A}}}solidFill'):
        ln.remove(sf)
    for nf in ln.findall(f'{{{NS_A}}}noFill'):
        ln.remove(nf)
    sf = etree.SubElement(ln, f'{{{NS_A}}}solidFill')
    clr = etree.SubElement(sf, f'{{{NS_A}}}srgbClr')
    clr.set('val', hex_color)

def set_no_border(spPr):
    ln = spPr.find(f'{{{NS_A}}}ln')
    if ln is None:
        ln = etree.SubElement(spPr, f'{{{NS_A}}}ln')
    for child in list(ln):
        ln.remove(child)
    etree.SubElement(ln, f'{{{NS_A}}}noFill')

def make_sp_xml(shape_id, name, left, top, width, height,
                fill_hex=None, no_fill=False, border_hex=None, no_border=False,
                border_pt=0.75):
    """Build a p:sp XML element directly — no python-pptx shape API."""
    sp = etree.Element(f'{{{NS_P}}}sp')

    nvSpPr = etree.SubElement(sp, f'{{{NS_P}}}nvSpPr')
    cNvPr  = etree.SubElement(nvSpPr, f'{{{NS_P}}}cNvPr')
    cNvPr.set('id', str(shape_id)); cNvPr.set('name', name)
    cNvSpPr = etree.SubElement(nvSpPr, f'{{{NS_P}}}cNvSpPr')
    cNvSpPr.set('txBox', '1')
    etree.SubElement(nvSpPr, f'{{{NS_P}}}nvPr')

    spPr = etree.SubElement(sp, f'{{{NS_P}}}spPr')
    xfrm = etree.SubElement(spPr, f'{{{NS_A}}}xfrm')
    off  = etree.SubElement(xfrm, f'{{{NS_A}}}off')
    off.set('x', str(int(left))); off.set('y', str(int(top)))
    ext  = etree.SubElement(xfrm, f'{{{NS_A}}}ext')
    ext.set('cx', str(int(width))); ext.set('cy', str(int(height)))
    pg = etree.SubElement(spPr, f'{{{NS_A}}}prstGeom')
    pg.set('prst', 'rect')
    etree.SubElement(pg, f'{{{NS_A}}}avLst')

    set_fill(spPr, hex_color=fill_hex, no_fill=no_fill)

    if no_border:
        set_no_border(spPr)
    elif border_hex:
        set_border(spPr, hex_color=border_hex, width_pt=border_pt)

    txBody = etree.SubElement(sp, f'{{{NS_P}}}txBody')
    bodyPr = etree.SubElement(txBody, f'{{{NS_A}}}bodyPr')
    bodyPr.set('wrap', 'square')
    bodyPr.set('lIns', '45720'); bodyPr.set('rIns', '45720')
    bodyPr.set('tIns', '30000'); bodyPr.set('bIns', '30000')
    etree.SubElement(txBody, f'{{{NS_A}}}lstStyle')

    return sp, txBody

def add_run(para, text, size_pt, bold=False, italic=False, color='000000', align=None):
    if align:
        pPr = etree.SubElement(para, f'{{{NS_A}}}pPr')
        pPr.set('algn', align)
    r = etree.SubElement(para, f'{{{NS_A}}}r')
    rPr = etree.SubElement(r, f'{{{NS_A}}}rPr')
    rPr.set('sz', str(int(size_pt * 100)))
    rPr.set('b', '1' if bold else '0')
    rPr.set('i', '1' if italic else '0')
    sf = etree.SubElement(rPr, f'{{{NS_A}}}solidFill')
    clr = etree.SubElement(sf, f'{{{NS_A}}}srgbClr')
    clr.set('val', color)
    t = etree.SubElement(r, f'{{{NS_A}}}t')
    t.text = text
    return r

def add_para_with_run(txBody, text, size_pt, bold=False, italic=False,
                       color='000000', align='l', space_before_pts=0):
    para = etree.SubElement(txBody, f'{{{NS_A}}}p')
    pPr  = etree.SubElement(para, f'{{{NS_A}}}pPr')
    pPr.set('algn', align)
    if space_before_pts:
        spcBef = etree.SubElement(pPr, f'{{{NS_A}}}spcBef')
        spcPts = etree.SubElement(spcBef, f'{{{NS_A}}}spcPts')
        spcPts.set('val', str(space_before_pts * 100))
    r = etree.SubElement(para, f'{{{NS_A}}}r')
    rPr = etree.SubElement(r, f'{{{NS_A}}}rPr')
    rPr.set('sz', str(int(size_pt * 100)))
    rPr.set('b', '1' if bold else '0')
    rPr.set('i', '1' if italic else '0')
    sf = etree.SubElement(rPr, f'{{{NS_A}}}solidFill')
    clr = etree.SubElement(sf, f'{{{NS_A}}}srgbClr')
    clr.set('val', color)
    t = etree.SubElement(r, f'{{{NS_A}}}t')
    t.text = text
    return para


# ─────────────────────────────────────────────────────────────
# Clear slide shapes (keep only nvGrpSpPr / grpSpPr)
# ─────────────────────────────────────────────────────────────
def clear_slide(slide):
    spTree = slide.shapes._spTree
    keep_tags = {f'{{{NS_P}}}nvGrpSpPr', f'{{{NS_P}}}grpSpPr'}
    for child in list(spTree):
        if child.tag not in keep_tags:
            spTree.remove(child)


# ─────────────────────────────────────────────────────────────
# Build org chart content into a slide's spTree
# ─────────────────────────────────────────────────────────────
def build_org_chart(slide, l2_info, start_id=100):
    spTree = slide.shapes._spTree
    SLIDE_W  = 12192000
    SLIDE_H  = 6858000
    SIDE_M   = 304800

    subtitle = l2_info['subtitle']
    l2_title = l2_info['l2_title']
    l3s      = l2_info['l3s']
    N        = len(l3s)
    sid      = start_id

    # ── Title placeholder ──────────────────────────────────────
    # Re-add a title placeholder (ph type=title)
    sp_title = etree.SubElement(spTree, f'{{{NS_P}}}sp')
    nvSpPr_t = etree.SubElement(sp_title, f'{{{NS_P}}}nvSpPr')
    cNvPr_t  = etree.SubElement(nvSpPr_t, f'{{{NS_P}}}cNvPr')
    cNvPr_t.set('id', str(sid)); cNvPr_t.set('name', 'Title 1'); sid += 1
    cNvSpPr_t = etree.SubElement(nvSpPr_t, f'{{{NS_P}}}cNvSpPr')
    etree.SubElement(cNvSpPr_t, f'{{{NS_A}}}spLocks').set('noGrp', '1')
    nvPr_t = etree.SubElement(nvSpPr_t, f'{{{NS_P}}}nvPr')
    etree.SubElement(nvPr_t, f'{{{NS_P}}}ph').set('type', 'title')
    spPr_t = etree.SubElement(sp_title, f'{{{NS_P}}}spPr')
    txBody_t = etree.SubElement(sp_title, f'{{{NS_P}}}txBody')
    etree.SubElement(txBody_t, f'{{{NS_A}}}bodyPr')
    etree.SubElement(txBody_t, f'{{{NS_A}}}lstStyle')
    p_t = etree.SubElement(txBody_t, f'{{{NS_A}}}p')
    rPr_t = etree.SubElement(etree.SubElement(p_t, f'{{{NS_A}}}r'), f'{{{NS_A}}}rPr')
    rPr_t.set('b', '1'); rPr_t.set('sz', '2200')
    sf_t = etree.SubElement(rPr_t, f'{{{NS_A}}}solidFill')
    etree.SubElement(sf_t, f'{{{NS_A}}}srgbClr').set('val', '000000')
    etree.SubElement(rPr_t.getparent(), f'{{{NS_A}}}t').text = 'Estructura Organizacional'

    # ── Red italic subtitle ────────────────────────────────────
    sp_sub, txb_sub = make_sp_xml(sid, 'Subtitle', 330200, 530000, 10000000, 280000,
                                   no_fill=True, no_border=True); sid += 1
    txb_sub.find(f'{{{NS_A}}}bodyPr').set('wrap', 'square')
    add_para_with_run(txb_sub, subtitle, 18, italic=True, color='CC0000', align='l')
    spTree.append(sp_sub)

    # ── L2 box (centered) ──────────────────────────────────────
    L2_W = 3200000; L2_H = 520000
    L2_L = (SLIDE_W - L2_W) // 2; L2_T = 900000
    sp_l2, txb_l2 = make_sp_xml(sid, 'L2Box', L2_L, L2_T, L2_W, L2_H,
                                  fill_hex='FFFFFF', border_hex='000000', border_pt=0.75); sid += 1
    for line in l2_title.split('\n'):
        add_para_with_run(txb_l2, line, 10, bold=True, color='000000', align='ctr')
    spTree.append(sp_l2)

    # ── Connector layout ──────────────────────────────────────
    L2_CX = L2_L + L2_W // 2
    L2_BOTTOM = L2_T + L2_H
    GAP  = 50000 if N <= 4 else 35000
    L3_W = (SLIDE_W - 2 * SIDE_M - (N - 1) * GAP) // N
    L3_T = 1730000; L3_H = 660000

    L3_LEFTS = [SIDE_M + i * (L3_W + GAP) for i in range(N)]
    L3_CXS   = [L + L3_W // 2 for L in L3_LEFTS]

    HBAR_Y   = L3_T - 100000
    VERT_LEN = HBAR_Y - L2_BOTTOM

    # Vertical from L2 to horizontal bar
    sp_v, txb_v = make_sp_xml(sid, 'ConnV', L2_CX - 6350, L2_BOTTOM, 12700, VERT_LEN,
                               fill_hex='000000', no_border=True); sid += 1
    spTree.append(sp_v)

    # Horizontal bar
    if N > 1:
        h_left = L3_CXS[0]; h_right = L3_CXS[-1]
        sp_h, txb_h = make_sp_xml(sid, 'ConnH', h_left, HBAR_Y - 6350,
                                   h_right - h_left, 12700,
                                   fill_hex='000000', no_border=True); sid += 1
        spTree.append(sp_h)

    # Vertical drops to each L3
    DROP_LEN = L3_T - HBAR_Y
    for cx in L3_CXS:
        sp_d, _ = make_sp_xml(sid, f'ConnD{sid}', cx - 6350, HBAR_Y,
                               12700, DROP_LEN, fill_hex='000000', no_border=True); sid += 1
        spTree.append(sp_d)

    # ── L3 boxes ──────────────────────────────────────────────
    DESC_T = L3_T + L3_H + 40000
    DESC_H = SLIDE_H - DESC_T - 60000

    title_pt  = 10 if N <= 4 else 9
    bullet_pt = 9  if N <= 4 else 7

    for i, l3 in enumerate(l3s):
        L3_L = L3_LEFTS[i]

        # L3 header box
        sp_l3, txb_l3 = make_sp_xml(sid, f'L3_{i}', L3_L, L3_T, L3_W, L3_H,
                                      fill_hex='FFFFFF', border_hex='000000', border_pt=0.75); sid += 1
        for line in l3['title'].split('\n'):
            add_para_with_run(txb_l3, line, title_pt, bold=True, color='000000', align='ctr')
        spTree.append(sp_l3)

        # Description box
        sp_d2, txb_d2 = make_sp_xml(sid, f'Desc_{i}', L3_L, DESC_T, L3_W, DESC_H,
                                      fill_hex='F5F5F5', border_hex='BFBFBF', border_pt=0.5); sid += 1
        first = True
        for bullet in l3['bullets']:
            spc = 0 if first else 6
            add_para_with_run(txb_d2, '• ' + bullet, bullet_pt, color='333333',
                               align='l', space_before_pts=spc)
            first = False
        spTree.append(sp_d2)

    print(f'  Built: {subtitle} ({N} L3s, {sid - start_id} shapes)')


# ─────────────────────────────────────────────────────────────
# Move a slide sldId from old_idx to new_idx
# ─────────────────────────────────────────────────────────────
def move_slide(prs, old_idx, new_idx):
    sldIdLst = prs.slides._sldIdLst
    sldId    = sldIdLst[old_idx]
    sldIdLst.remove(sldId)
    sldIdLst.insert(new_idx, sldId)


# ─────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────
def main():
    print(f'Loading working version: {SRC}')
    prs = Presentation(SRC)
    print(f'  Slides: {len(prs.slides)}')  # 62

    layout = prs.slide_layouts[11]  # 1_Title Only

    # ── Slide 6 (index 5): clear and reuse in-place ───────────
    print('Rebuilding slide 6 in-place...')
    slide6 = prs.slides[5]
    clear_slide(slide6)
    build_org_chart(slide6, L2_DATA[0], start_id=200)

    # ── Add 6 more slides (appended → no filename conflicts) ──
    print('Adding 6 new slides...')
    for l2_info in L2_DATA[1:]:
        new_slide = prs.slides.add_slide(layout)
        clear_slide(new_slide)
        build_org_chart(new_slide, l2_info, start_id=300)

    print(f'  Slides after adding: {len(prs.slides)}')  # 68

    # ── Move new slides (currently at end) to positions 7-12 ──
    print('Reordering slides...')
    for i in range(6):
        tail_idx = len(prs.slides) - (6 - i)
        move_slide(prs, tail_idx, 6 + i)
        print(f'  Moved to position {7 + i}')

    prs.save(DEST)
    print(f'\nSaved: {DEST}')

    # ── Verify ZIP cleanliness ────────────────────────────────
    with zipfile.ZipFile(DEST, 'r') as z:
        names = z.namelist()
        dupes = [n for n in names if names.count(n) > 1]
        print(f'ZIP duplicates: {set(dupes)}')
        ct_slides = set(re.findall(r'PartName="/ppt/slides/(slide\d+\.xml)"',
                                    z.read('[Content_Types].xml').decode()))
        zip_slides = set(n.replace('ppt/slides/', '') for n in names
                         if n.startswith('ppt/slides/slide') and '.rels' not in n)
        orphans = zip_slides - ct_slides
        print(f'Orphaned slides: {sorted(orphans)}')

    # ── Verify with python-pptx ───────────────────────────────
    prs2 = Presentation(DEST)
    print(f'Final slide count: {len(prs2.slides)}')
    for i in range(13):
        s = prs2.slides[i]
        for shape in s.shapes:
            if 'Title' in shape.name and hasattr(shape, 'text_frame'):
                t = shape.text_frame.text[:55]
                print(f'  Slide {i+1}: {t!r}')
                break


if __name__ == '__main__':
    main()
