#!/usr/bin/env python3
"""Generate TI Playbook PPTX from Finanzas Playbook base."""

import shutil, os, zipfile, re
from lxml import etree
from pptx import Presentation

SRC = '/home/user/oxxo/Finanzas/Playbook/Finanzas Playbook.pptx'
OUT = '/home/user/oxxo/TI/Playbook/TI Playbook.pptx'

NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
SLIDE_W = 12192000; SLIDE_H = 6858000
MARGIN = 365760; CONTENT_W = SLIDE_W - 2 * MARGIN
SIDE_M = 304800

BADGE_COLORS = {
    'R':'70AD47','A':'FFC000','A/R':'FFC000','S':'FF7070','C':'4472C4','I':'A6A6A6',
}

def rasci_color(letter):
    return BADGE_COLORS.get(letter, 'A6A6A6')
ROLE_DESCS = {
    'R':  'Ejecuta directamente esta responsabilidad como responsable operativo principal.',
    'A':  'Aprueba y responde ante la organización por los resultados de esta responsabilidad.',
    'A/R':'Ejecuta y lidera esta responsabilidad, tomando las decisiones necesarias para garantizar su cumplimiento y calidad de resultados.',
    'S':  'Brinda apoyo activo en esta responsabilidad, aportando su expertise especializado cuando se requiere.',
    'C':  'Aporta perspectiva experta como área consultada, contribuyendo con información y criterio antes de las decisiones clave.',
    'I':  'Recibe información del avance y resultados para anticipar impactos en su área y tomar acciones preventivas.',
}

# Org-chart layout (same as rebuild_final.py)
SUBTITLE_TOP=660000; SUBTITLE_H=220000
L2_TOP=1900000; L2_H=500000; L2_W=3200000
L3_TOP=2900000; L3_H=560000
DESC_TOP=L3_TOP+L3_H+50000; DESC_H_STD=2850000
HBAR_Y=L3_TOP-100000

# ── TI DATA ───────────────────────────────────────────────────

TI_DEPT = 'TI'
TI_ROLE_SHORT = ['Dir. TI','Gte Sr Tienda','Gte Sr CdS','Gte Cust/Dig',
                 'Gte BO','Gte Arquit.','Gte Plan.','Gte Dllo.','Gte Ciberseg.']
TI_ROLE_FULL = [
    'Director TI OXXO Mx','Gte Sr Tienda y Comercial','Gte Sr Solucs. CdS',
    'Gte Sr Customer y Digital','Gte Solucs. Back Office','Gte Arquitectura',
    'Gte Planeación Proyectos','Gte Desarrollo Soluciones','Gte Sr Ciberseguridad',
]

# (section, name, description, roles_dict {0..8: letter}, explanation)
TI_RASCI = [
    ('Director TI OXXO Mx','Liderar estrategia tecnológica y medir valor entregado',
     'Lidera la estrategia tecnológica alineando prioridades con las necesidades del negocio y midiendo el valor entregado.',
     {0:'A/R',1:'C',2:'C',3:'C',4:'I',5:'S',6:'S',7:'I',8:'I'},
     'El Director rinde cuentas y ejecuta la visión de alto nivel (A/R). Los Gtes. de Soluciones son consultados (C) para asegurar que la estrategia sirva al negocio.'),
    ('','Definir arquitectura, políticas y gobernanza ERP/POS',
     'Define la arquitectura empresarial, políticas de integración y gobernanza de plataformas críticas (ERP, POS, WMS).',
     {0:'A',1:'I',2:'I',3:'I',4:'I',5:'R',6:'I',7:'S',8:'S'},
     'El Gte. de Arquitectura es el autor de las reglas técnicas (R). El Director aprueba (A) y Desarrollo apoya (S) en la viabilidad técnica.'),
    ('','Asegurar diseño de sistemas robustos y escalables',
     'Asegura el diseño y construcción de sistemas robustos, escalables y con alta disponibilidad.',
     {0:'A',1:'S',2:'S',3:'S',4:'S',5:'S',6:'I',7:'R',8:'I'},
     'El Gte. de Desarrollo es el responsable de que el código funcione bien (R). Arquitectura y los Gtes. de Soluciones apoyan (S) en el diseño funcional.'),
    ('','Garantizar protección de datos y cumplimiento regulatorio',
     'Garantiza la protección de datos y el cumplimiento regulatorio en conjunto con TI Corporativo.',
     {0:'A',1:'I',2:'I',3:'I',4:'I',5:'C',6:'I',7:'I',8:'R'},
     'El Gte. de Ciberseguridad es el ejecutor de las salvaguardas (R). El Director rinde cuentas (A) y Arquitectura es consultada (C) para el diseño seguro.'),
    ('','Promover analítica, omnicanalidad y gob. de datos',
     'Promueve capacidades de analítica, omnicanalidad y fomenta una cultura de gobernanza de datos.',
     {0:'A/R',1:'S',2:'S',3:'S',4:'I',5:'C',6:'I',7:'I',8:'I'},
     'El Director lidera la iniciativa cultural y estratégica (A/R). Arquitectura consulta (C) y los Gtes. de Soluciones apoyan (S) con los casos de uso.'),
    ('','Administrar CAPEX/OPEX y contratos con proveedores',
     'Administra el presupuesto (CAPEX/OPEX) y negocia contratos estratégicos con proveedores tecnológicos.',
     {0:'A/R',1:'S',2:'S',3:'S',4:'S',5:'I',6:'C',7:'I',8:'I'},
     'El Director gestiona el presupuesto (A/R). Planeación de Proyectos apoya (S) con el control de los gastos por iniciativa.'),
    ('Gerente Sr Soluciones Tienda y Comercial','Roadmap de soluciones digitales para tienda y comercial',
     'Establece el roadmap de soluciones digitales para los procesos de tienda e iniciativas comerciales.',
     {0:'I',1:'A/R',2:'I',3:'I',4:'I',5:'S',6:'S',7:'C',8:'I'},
     'El Gte. de Tienda/Comercial define el futuro de su vertical (A/R). Arquitectura y Planeación apoyan (S) para que sea posible ejecutarlo.'),
    ('','Criterios de valor/costo para priorización de backlog',
     'Establece criterios de valor, costo y riesgo para maximizar el impacto de las iniciativas comerciales.',
     {0:'C',1:'A/R',2:'I',3:'I',4:'I',5:'I',6:'S',7:'I',8:'I'},
     'El Gte. de Tienda/Comercial decide qué va primero según el negocio (A/R). Planeación apoya (S) calculando los recursos necesarios.'),
    ('','Analizar puntos de fricción en operación de tienda',
     'Analiza puntos de fricción en la operación de tienda para proponer soluciones que agilicen la labor del colaborador.',
     {0:'I',1:'A/R',2:'I',3:'I',4:'I',5:'I',6:'I',7:'S',8:'I'},
     'El Gte. de Tienda/Comercial identifica los problemas del usuario (A/R). Desarrollo apoya (S) proponiendo cómo arreglarlos con software.'),
    ('','Garantizar estabilidad y time-to-market ágil',
     'Garantiza que las soluciones propuestas sean estables, seguras y tengan un tiempo de entrega ágil.',
     {0:'I',1:'A/R',2:'I',3:'I',4:'I',5:'S',6:'S',7:'R',8:'I'},
     'El Gte. de Tienda/Comercial rinde cuentas del éxito (A), pero Desarrollo es quien ejecuta la entrega rápida (R).'),
    ('','Lineamientos para replicabilidad en red de tiendas',
     'Desarrolla lineamientos para que las soluciones exitosas puedan replicarse en toda la red de tiendas.',
     {0:'I',1:'A/R',2:'I',3:'I',4:'I',5:'C',6:'I',7:'I',8:'I'},
     'El Gte. de Tienda/Comercial asegura que lo que funcione en una, funcione en todas (A/R). Arquitectura valida (C) el estándar.'),
    ('Gerente Sr Soluciones CdS (Cadena de Suministro)','Roadmap de soluciones Abasto, CEDIS y Distribución',
     'Define el roadmap de soluciones para Abasto, CEDIS y Distribución (end-to-end).',
     {0:'I',1:'I',2:'A/R',3:'I',4:'I',5:'S',6:'S',7:'C',8:'I'},
     'El Gte. de CdS es el dueño de la estrategia logística-tecnológica (A/R). Desarrollo es consultado (C) para estimar tiempos.'),
    ('','Identificar tecnologías para eficiencia en CdS',
     'Identifica oportunidades tecnológicas que impacten directamente en la eficiencia de la cadena de suministro.',
     {0:'I',1:'I',2:'A/R',3:'I',4:'I',5:'S',6:'I',7:'S',8:'I'},
     'El Gte. de CdS busca cómo mejorar la operación (A/R). Arquitectura y Desarrollo apoyan (S) evaluando la tecnología.'),
    ('','Seguimiento a inversión y ROI en soluciones de CdS',
     'Brinda seguimiento periódico a la inversión en soluciones de CdS para asegurar el retorno de inversión.',
     {0:'A',1:'I',2:'R',3:'I',4:'I',5:'I',6:'S',7:'I',8:'I'},
     'El Gte. de CdS vigila que el dinero rinda (R) bajo la supervisión del Director (A). Planeación apoya (S) con el tracking.'),
    ('','Negociar SLAs competitivos con socios logísticos',
     'Negocia acuerdos y asegura niveles de servicio competitivos con socios tecnológicos logísticos.',
     {0:'C',1:'I',2:'A/R',3:'I',4:'I',5:'I',6:'I',7:'I',8:'I'},
     'El Gte. de CdS negocia con los proveedores técnicos de la cadena (A/R). El Director es consultado (C) en contratos grandes.'),
    ('','Evaluar herramientas de industria para ventaja en logística',
     'Evalúa herramientas y metodologías de industria para mantener una ventaja competitiva en logística.',
     {0:'I',1:'I',2:'A/R',3:'I',4:'I',5:'C',6:'I',7:'I',8:'I'},
     'El Gte. de CdS es el experto en el mercado logístico (A/R). Arquitectura consulta (C) para asegurar compatibilidad.'),
    ('Gerente Sr Soluciones TI Customer y Digital','Estrategias cliente-marca mediante habilitadores digitales',
     'Define estrategias para mejorar la relación cliente-marca mediante habilitadores digitales.',
     {0:'I',1:'I',2:'I',3:'A/R',4:'I',5:'S',6:'S',7:'C',8:'I'},
     'El Gte. Customer/Digital lidera la visión del cliente final (A/R). Desarrollo es consultado (C) para ver si la app o portal es posible.'),
    ('','Estrategias de calidad de software y estabilidad apps',
     'Crea estrategias de desarrollo y calidad que optimicen el uso de recursos y la estabilidad de las apps.',
     {0:'I',1:'S',2:'S',3:'A/R',4:'I',5:'I',6:'I',7:'R',8:'I'},
     'El Gte. Customer/Digital responde por la app (A), pero Desarrollo ejecuta las pruebas y calidad (R).'),
    ('','Revisar tendencias para evitar obsolescencia digital',
     'Revisa tendencias de industria para evitar la obsolescencia y habilitar nuevas oportunidades de negocio.',
     {0:'I',1:'I',2:'I',3:'A/R',4:'I',5:'C',6:'I',7:'S',8:'I'},
     'El Gte. Customer/Digital vigila que la tecnología no sea vieja (A/R). Arquitectura apoya (C) con el radar tecnológico.'),
    ('','Anticipar impactos y riesgos en soluciones digitales',
     'Analiza y anticipa posibles impactos en soluciones digitales para facilitar la toma de decisiones.',
     {0:'I',1:'I',2:'I',3:'A/R',4:'I',5:'I',6:'S',7:'I',8:'S'},
     'El Gte. Customer/Digital identifica riesgos de usuario (A/R). Ciberseguridad apoya (S) evaluando riesgos de hackeo.'),
    ('','Diseño de automatización para plataformas digitales',
     'Diseña soluciones que simplifiquen procesos y reduzcan tiempos de ejecución en plataformas digitales.',
     {0:'I',1:'I',2:'I',3:'A/R',4:'I',5:'S',6:'I',7:'R',8:'I'},
     'El Gte. Customer/Digital rinde cuentas (A), pero Desarrollo construye la automatización (R).'),
    ('Gerente Soluciones Back Office','Soluciones para áreas soporte alineadas a objetivos',
     'Define soluciones para áreas de soporte administrativo alineadas a los objetivos de la empresa.',
     {0:'I',1:'I',2:'I',3:'I',4:'A/R',5:'S',6:'S',7:'C',8:'I'},
     'El Gte. de Back Office define qué necesitan RH, Finanzas, etc. (A/R). Desarrollo consulta (C) el alcance.'),
    ('','Administrar servidores, redes y almacenamiento',
     'Administra servidores, redes y almacenamiento, garantizando disponibilidad y seguridad.',
     {0:'I',1:'I',2:'I',3:'I',4:'A/R',5:'S',6:'I',7:'S',8:'S'},
     "El Gte. de Back Office es el responsable de los 'fierros' e infraestructura (A/R). Ciberseguridad apoya (S) con la protección."),
    ('','Incrementar eficiencia interna mediante tecnología',
     'Identifica áreas de mejora en procesos internos para incrementar eficiencia mediante la tecnología.',
     {0:'I',1:'S',2:'S',3:'S',4:'A/R',5:'I',6:'I',7:'R',8:'I'},
     'El Gte. de Back Office rinde cuentas del ahorro interno (A). Desarrollo ejecuta (R) las herramientas de eficiencia.'),
    ('','Cumplimiento normativo y regulatorio de Back Office',
     'Revisa políticas internas y regulaciones externas para mantener la integridad operativa del Back Office.',
     {0:'I',1:'I',2:'I',3:'I',4:'A/R',5:'I',6:'I',7:'I',8:'S'},
     'El Gte. de Back Office asegura que los sistemas internos cumplan la ley (A/R). Ciberseguridad apoya (S) con auditorías.'),
    ('','Mantenimiento de sistemas de soporte de negocio',
     'Lidera el mantenimiento de sistemas que soportan los procesos de negocio actuales y futuros.',
     {0:'I',1:'I',2:'I',3:'I',4:'A/R',5:'I',6:'I',7:'R',8:'I'},
     'El Gte. de Back Office rinde cuentas de la vida del sistema (A). Desarrollo ejecuta los cambios y parches (R).'),
    ('Gerente Arquitectura','Evolucionar estrategia de Arquitectura Empresarial',
     'Evoluciona la Estrategia de Arquitectura Empresarial estableciendo marcos de trabajo modernos.',
     {0:'A',1:'I',2:'I',3:'I',4:'I',5:'R',6:'C',7:'S',8:'I'},
     'El Gte. de Arquitectura es el autor de los planos maestros (R). El Director aprueba la dirección estratégica (A).'),
    ('','Asegurar que diseños cumplan lineamientos corporativos',
     'Colabora con equipos de TI para asegurar que los diseños de solución cumplan con los lineamientos corporativos.',
     {0:'I',1:'C',2:'C',3:'C',4:'C',5:'A/R',6:'I',7:'S',8:'S'},
     'El Gte. de Arquitectura es el juez que valida los proyectos (A/R). Los Gtes. de Soluciones son consultados (C) por funcionalidad.'),
    ('','Incorporar nuevas tecnologías al ecosistema actual',
     'Adopta estándares de industria para incorporar nuevas tecnologías al ecosistema actual de forma fluida.',
     {0:'I',1:'S',2:'S',3:'S',4:'S',5:'A/R',6:'I',7:'S',8:'I'},
     'El Gte. de Arquitectura lidera la integración de nuevas piezas (A/R). Desarrollo apoya (S) en la implementación piloto.'),
    ('','Recomendaciones técnicas en procesos de RFP',
     'Explora nuevas tecnologías y brinda recomendaciones en procesos de licitación (RFP).',
     {0:'I',1:'S',2:'S',3:'S',4:'S',5:'A/R',6:'I',7:'I',8:'I'},
     'El Gte. de Arquitectura evalúa técnicamente a los proveedores (A/R). Los Gtes. de Soluciones apoyan (S) con requerimientos.'),
    ('','Visibilidad de capacidades tecnológicas del negocio',
     'Mantiene actualizados los artefactos que dan visibilidad sobre las capacidades tecnológicas del negocio.',
     {0:'I',1:'I',2:'I',3:'I',4:'I',5:'A/R',6:'I',7:'I',8:'I'},
     'El Gte. de Arquitectura mantiene el catálogo de lo que TI puede hacer (A/R). Es un rol puramente informativo para el resto.'),
    ('Gerente Área Planeación de Proyectos','Estándares y herramientas para gestión de proyectos',
     'Establece los estándares, metodologías y herramientas para la gestión de proyectos de TI.',
     {0:'A',1:'I',2:'I',3:'I',4:'I',5:'S',6:'R',7:'I',8:'I'},
     "El Gte. de Planeación define el 'cómo' se gestionan los proyectos (R). El Director aprueba la metodología (A)."),
    ('','Supervisar asignación de cargas de trabajo (Staffing)',
     'Supervisa la asignación de cargas de trabajo para asegurar eficiencia y cumplimiento de plazos.',
     {0:'I',1:'S',2:'S',3:'S',4:'S',5:'I',6:'A/R',7:'I',8:'I'},
     'El Gte. de Planeación vigila que nadie esté saturado (A/R). Los Gtes. de Soluciones apoyan (S) reportando sus necesidades.'),
    ('','Priorización transversal alineada al impacto negocio',
     'Crea espacios de colaboración para alinear la ejecución de proyectos con el impacto esperado en el negocio.',
     {0:'A',1:'C',2:'C',3:'C',4:'C',5:'I',6:'R',7:'I',8:'I'},
     'El Gte. de Planeación lidera el foro de priorización global (R) bajo la responsabilidad del Director (A).'),
    ('','Oportunidades para madurez en gestión de proyectos',
     'Identifica oportunidades para incrementar la madurez y eficiencia en la gestión de proyectos.',
     {0:'I',1:'I',2:'I',3:'I',4:'I',5:'S',6:'A/R',7:'I',8:'I'},
     'El Gte. de Planeación busca mejorar la oficina de proyectos (A/R). Arquitectura apoya (S) con herramientas de visibilidad.'),
    ('','Reportar KPIs del portafolio (costo, tiempo, calidad)',
     'Reporta indicadores clave (costos, tiempos, calidad) para garantizar el éxito del portafolio.',
     {0:'I',1:'I',2:'I',3:'I',4:'I',5:'I',6:'A/R',7:'I',8:'I'},
     'El Gte. de Planeación es el responsable de decir si los proyectos van a tiempo y presupuesto (A/R).'),
    ('Gerente Desarrollo Soluciones','Estrategia de construcción E2E bajo marcos ágiles',
     'Define la estrategia de construcción E2E bajo marcos de trabajo ágiles para optimizar el tiempo de entrega.',
     {0:'I',1:'I',2:'I',3:'I',4:'I',5:'S',6:'S',7:'A/R',8:'I'},
     'El Gte. de Desarrollo define cómo se programa (A/R). Arquitectura y Planeación apoyan (S) con el marco de trabajo.'),
    ('','Alineación de valor durante desarrollo de software',
     'Asegura la correcta alineación de la propuesta de valor durante todo el proceso de desarrollo de software.',
     {0:'I',1:'S',2:'S',3:'S',4:'I',5:'I',6:'S',7:'A/R',8:'I'},
     'El Gte. de Desarrollo asegura que lo programado sirva (A/R). Los Gtes. de Soluciones apoyan (S) validando que sea lo que pidieron.'),
    ('','Políticas y roles para planeación técnica',
     'Define políticas y roles claros para la planeación y seguimiento técnico de los desarrollos.',
     {0:'I',1:'I',2:'I',3:'I',4:'I',5:'S',6:'S',7:'A/R',8:'I'},
     'El Gte. de Desarrollo define quién hace qué en el equipo técnico (A/R). Planeación apoya (S) con la estructura.'),
    ('','Calidad y eficiencia de equipos de desarrollo',
     'Implementa modelos de industria para mejorar la calidad y eficiencia de los equipos de desarrollo.',
     {0:'I',1:'I',2:'I',3:'I',4:'I',5:'S',6:'I',7:'A/R',8:'I'},
     'El Gte. de Desarrollo es el dueño de la productividad técnica (A/R). Arquitectura apoya (S) con estándares de código.'),
    ('Gerente Sr Ciberseguridad','Estrategia para salvaguardar activos digitales',
     'Define y ejecuta la estrategia para salvaguardar la información y activos digitales de la empresa.',
     {0:'A',1:'I',2:'I',3:'I',4:'S',5:'C',6:'I',7:'I',8:'R'},
     'El Gte. de Ciberseguridad ejecuta la defensa (R). El Director rinde cuentas ante el negocio (A).'),
    ('','Respuesta ante amenazas o vulnerabilidades',
     'Lidera la respuesta ante amenazas o vulnerabilidades detectadas en el ecosistema tecnológico.',
     {0:'A',1:'S',2:'S',3:'S',4:'S',5:'I',6:'I',7:'I',8:'R'},
     'El Gte. de Ciberseguridad lidera el equipo de respuesta a incidentes (R). El Director responde por el impacto (A).'),
    ('','Cumplimiento de protocolos de seguridad en soluciones',
     'Asegura que todas las soluciones (Tienda, Digital, CdS) cumplan con los protocolos de seguridad.',
     {0:'I',1:'C',2:'C',3:'C',4:'C',5:'S',6:'I',7:'S',8:'A/R'},
     'El Gte. de Ciberseguridad rinde cuentas del cumplimiento (A/R). Los Gtes. de Soluciones son consultados (C) para aplicar las reglas.'),
]

TI_L1_L2_OVERVIEW = {
    'l1_title': 'Director TI\nOXXO Mx',
    'l2s': [
        ('Gte Sr Tienda\ny Comercial','Soluciones digitales para tiendas e iniciativas comerciales.'),
        ('Gte Sr\nSolucs. CdS','Soluciones para Abasto, CEDIS y Distribución end-to-end.'),
        ('Gte Sr Solucs.\nCustomer y Digital','Estrategias digitales de relación cliente-marca.'),
        ('Gte Solucs.\nBack Office','Soluciones de soporte administrativo e infraestructura.'),
        ('Gte\nArquitectura','Arquitectura empresarial, estándares y gobernanza técnica.'),
        ('Gte Planeación\nProyectos','Gestión de portafolio, PMO y metodologías de proyectos.'),
        ('Gte Desarrollo\nSoluciones','Construcción E2E, DevOps y calidad del software.'),
        ('Gte Sr\nCiberseguridad','Protección de activos digitales y cumplimiento de seguridad.'),
    ]
}

TI_L2_DATA = [
    {
        'subtitle':'Desarrollo y Soluciones','l2_title':'Gerente Desarrollo\ny Soluciones',
        'l3s':[
            {'title':'Resp. Arquitectura\nde Solución','bullets':[
                'Diseñar los diagramas de arquitectura técnica y flujos de datos.',
                'Validar la viabilidad técnica e integración de nuevas soluciones.',
                'Asegurar el cumplimiento de atributos de calidad técnica.',
                'Elaborar y mantener el catálogo de estándares de arquitectura.',
            ]},
            {'title':'Resp. DevOps\ny Herramientas','bullets':[
                'Diseñar e implementar los pipelines de CI/CD.',
                'Gestionar y optimizar la infraestructura en nube y contenedores.',
                'Monitorear proactivamente la salud de los entornos productivos.',
                'Administrar los entornos de desarrollo, pruebas y producción.',
            ]},
            {'title':'Resp. Ingeniería\nSoftware Backend','bullets':[
                'Desarrollar y supervisar microservicios y lógica de negocio del servidor.',
                'Optimizar bases de datos y modelos de datos.',
                'Implementar protocolos de seguridad en la capa de servidor.',
                'Diseñar y mantener las APIs de integración entre sistemas.',
            ]},
            {'title':'Resp. Ingeniería\nFrontend / Mobile','bullets':[
                'Implementar diseños UI/UX en plataformas iOS, Android y Web.',
                'Asegurar la compatibilidad multiplataforma de las interfaces.',
                'Optimizar el rendimiento del lado del cliente.',
                'Mantener y evolucionar el UI Kit estándar de la organización.',
            ]},
        ]
    },
    {
        'subtitle':'Planeación de Proyectos','l2_title':'Gerente Planeación\nProyectos',
        'l3s':[
            {'title':'Resp. Gestión\nPilotos EVA','bullets':[
                'Diseñar el plan de cada piloto tecnológico.',
                'Monitorear el desempeño técnico y operativo.',
                'Diseñar e impartir materiales de capacitación.',
                'Elaborar reportes de resultados del piloto.',
            ]},
            {'title':'Resp. PMO\nOXXO México','bullets':[
                'Mantener la visibilidad centralizada del portafolio.',
                'Identificar y resolver proactivamente dependencias.',
                'Consolidar y analizar los indicadores de avance.',
                'Validar que los entregables de los proyectos cumplan criterios.',
            ]},
            {'title':'Resp. Procesos\nde Solución','bullets':[
                'Definir y documentar los procesos estándar de TI.',
                'Identificar y diagnosticar puntos de fricción.',
                'Diseñar e implementar mecanismos de gobernanza metodológica.',
                'Implementar el plan de gestión del cambio técnico.',
            ]},
        ]
    },
    {
        'subtitle':'Soluciones Cadena de Suministro','l2_title':'Gerente Soluciones\nCadena de Suministro',
        'l3s':[
            {'title':'Resp. Solucs.\nAbasto y Planeación','bullets':[
                'Documentar y mantener el backlog del dominio Abasto y Planeación.',
                'Identificar las necesidades tecnológicas del equipo de Abasto y Compras.',
                'Ejecutar el seguimiento del ciclo de vida de las soluciones.',
                'Elaborar la evaluación de ajustes en parámetros de sistemas de resurtido.',
            ]},
            {'title':'Resp. Solucs.\nAlmacén (WMS)','bullets':[
                'Documentar y mantener el backlog del dominio WMS.',
                'Identificar las necesidades tecnológicas de los equipos de operación CEDIS.',
                'Ejecutar el seguimiento del ciclo de vida de las soluciones WMS.',
                'Administrar y configurar el sistema WMS para optimizar procesos de CEDIS.',
            ]},
            {'title':'Resp. Solucs.\nTransporte y Dist.','bullets':[
                'Documentar y mantener el backlog del dominio Transporte y Distribución.',
                'Identificar las necesidades tecnológicas de Transporte y Tráfico.',
                'Ejecutar el seguimiento del ciclo de vida de las soluciones TMS.',
                'Implementar herramientas de telemetría y monitoreo de flota.',
            ]},
            {'title':'Resp. Visibilidad\ny Analítica de Red','bullets':[
                'Documentar y mantener el backlog del dominio Visibilidad y Analítica.',
                'Identificar las necesidades de visibilidad de las áreas logísticas.',
                'Diseñar y desarrollar la Torre de Control de la cadena de suministro.',
                'Ejecutar el seguimiento del ciclo de vida de las soluciones analíticas.',
            ]},
            {'title':'TPO Core\nAdministrativo','bullets':[
                'Documentar y mantener el backlog del dominio Core Administrativo.',
                'Identificar las necesidades tecnológicas de Finanzas CdS y BackOffice.',
                'Ejecutar el seguimiento del ciclo de vida de las soluciones administrativas.',
                'Elaborar la evaluación y validación de soluciones para el dominio Core Adm.',
            ]},
            {'title':'TPO EDI','bullets':[
                'Documentar y mantener el backlog de iniciativas de la plataforma EDI.',
                'Identificar las necesidades tecnológicas de Abastecimiento y CxP.',
                'Ejecutar el seguimiento del ciclo de vida de las soluciones EDI.',
                'Mantener los estándares técnicos de comunicación electrónica (EDIFact, XML, AS2).',
            ]},
            {'title':'TPO\nMercaderías','bullets':[
                'Documentar y mantener el backlog del dominio Mercaderías.',
                'Identificar las necesidades tecnológicas de Comercial y Operaciones CEDIS.',
                'Ejecutar el seguimiento del ciclo de vida de las soluciones de datos maestros.',
                'Elaborar la evaluación de soluciones para la gestión de catálogos.',
            ]},
        ]
    },
    {
        'subtitle':'Soluciones Customer y Digital','l2_title':'Gerente Soluciones\nCustomer y Digital',
        'l3s':[
            {'title':'Resp. Analytics y\nCustomer Insight','bullets':[
                'Documentar y mantener el backlog del dominio Analytics.',
                'Identificar las necesidades de información y analítica de las áreas comerciales.',
                'Diseñar y mantener modelos de datos y dashboards de métricas clave del cliente.',
                'Ejecutar el seguimiento del ciclo de vida de las soluciones analíticas.',
            ]},
            {'title':'Resp. Canales\nDigitales App y Web','bullets':[
                'Documentar y mantener el backlog de canales digitales.',
                'Identificar las necesidades tecnológicas de los equipos de Diseño UX.',
                'Ejecutar el seguimiento del ciclo de vida de la App móvil y plataformas web.',
                'Elaborar la evaluación y validación de nuevas funcionalidades y stack tecnológico.',
            ]},
            {'title':'Resp. Pagos\nDigitales y Servicios','bullets':[
                'Documentar y mantener el backlog del dominio Pagos Digitales.',
                'Identificar las necesidades tecnológicas de Servicios Financieros y Tesorería.',
                'Ejecutar el seguimiento del ciclo de vida de las soluciones de pagos digitales.',
                'Implementar la integración de pasarelas de pago y billeteras digitales.',
            ]},
            {'title':'Resp. Solucs.\nLoyalty y CRM','bullets':[
                'Documentar y mantener el backlog del dominio Loyalty y CRM.',
                'Identificar las necesidades tecnológicas de Marketing y Lealtad.',
                'Ejecutar el seguimiento del ciclo de vida de las soluciones de Loyalty y CRM.',
                'Elaborar la evaluación y validación de funcionalidades del programa de lealtad.',
            ]},
        ]
    },
    {
        'subtitle':'Soluciones Backoffice','l2_title':'Gerente de Soluciones\nBackoffice',
        'l3s':[
            {'title':'TPO\nAutomatización','bullets':[
                'Identificar y documentar oportunidades de automatización en procesos backoffice.',
                'Diseñar soluciones de automatización (RPA, integraciones, workflows digitales).',
                'Implementar y dar seguimiento a las soluciones de automatización.',
                'Ejecutar la validación de las automatizaciones desarrolladas.',
            ]},
            {'title':'TPO\nFAST','bullets':[
                'Documentar y mantener el backlog de iniciativas de la plataforma FAST.',
                'Identificar las necesidades tecnológicas de las áreas usuarias sobre FAST.',
                'Ejecutar el seguimiento del ciclo de vida de las soluciones sobre FAST.',
                'Elaborar la evaluación y validación de desarrollos y configuraciones sobre FAST.',
            ]},
            {'title':'TPO Finanzas\nIngresos','bullets':[
                'Documentar y mantener el backlog del dominio Finanzas Ingresos.',
                'Identificar las necesidades tecnológicas de las áreas de Finanzas.',
                'Ejecutar el seguimiento del ciclo de vida de las soluciones de Finanzas Ingresos.',
                'Elaborar la evaluación y validación de soluciones para el dominio financiero.',
            ]},
            {'title':'TPO Fiscal\nCumplimiento','bullets':[
                'Documentar y mantener el backlog del dominio Fiscal Cumplimiento.',
                'Identificar las necesidades tecnológicas del área Fiscal.',
                'Ejecutar el seguimiento del ciclo de vida de las soluciones tecnológicas fiscales.',
                'Elaborar la evaluación y validación de soluciones y desarrollos fiscales.',
            ]},
            {'title':'TPO\nRH','bullets':[
                'Documentar y mantener el backlog del dominio RH.',
                'Identificar las necesidades tecnológicas del área de RH.',
                'Ejecutar el seguimiento del ciclo de vida de las soluciones tecnológicas de RH.',
                'Elaborar la evaluación y validación de soluciones para el dominio RH.',
            ]},
        ]
    },
    {
        'subtitle':'Soluciones Tienda y Comercial','l2_title':'Gerente de Soluciones\nTienda y Comercial',
        'l3s':[
            {'title':'Gte Área Solucs.\nTI Mercaderías','bullets':[
                'Documentar y mantener el backlog del dominio Mercaderías.',
                'Identificar las necesidades tecnológicas de Comercial y Cadena de Suministro.',
                'Ejecutar el seguimiento del ciclo de vida de las soluciones de mercaderías.',
                'Diseñar y evolucionar las herramientas digitales de control de inventarios.',
            ]},
            {'title':'Gte Área Solucs.\nTI POS','bullets':[
                'Documentar y mantener el backlog del sistema POS.',
                'Identificar las necesidades tecnológicas de las áreas usuarias del POS.',
                'Ejecutar el seguimiento del ciclo de vida de las aplicaciones del POS.',
                'Implementar nuevas funcionalidades transaccionales en el POS.',
            ]},
            {'title':'Gte Área Solucs.\nTI Servicios','bullets':[
                'Documentar y mantener el backlog del dominio Servicios.',
                'Identificar las necesidades tecnológicas de Servicios Financieros.',
                'Ejecutar el seguimiento del ciclo de vida de las soluciones de servicios en tienda.',
                'Diseñar y optimizar los flujos de registro y pago de servicios en el POS.',
            ]},
            {'title':'Resp. Solucs.\nOperaciones','bullets':[
                'Documentar y mantener el backlog del dominio Operaciones.',
                'Identificar las necesidades tecnológicas de Operaciones y Excelencia Operativa.',
                'Ejecutar el seguimiento del ciclo de vida de las soluciones operativas.',
                'Desarrollar y mantener las apps móviles de gestión de tareas y auditorías.',
            ]},
            {'title':'Resp. Solucs.\nPrecios y Promociones','bullets':[
                'Documentar y mantener el backlog del dominio Precios y Promociones.',
                'Identificar las necesidades tecnológicas de Pricing y Comercial.',
                'Ejecutar el seguimiento del ciclo de vida de las soluciones de precios y promociones.',
                'Configurar y validar la lógica de jerarquía y convivencia de promociones.',
            ]},
        ]
    },
]

TI_AGENDA = [
    '1. Macroestructura OXXO MX',
    '2. Mapa Interacciones TI',
    '3. Estructura Organizacional TI',
    '4. Matriz RASCI (Metodología)',
    '5. Detalle de la Matriz RASCI por Responsabilidad',
    '6. Estructura Organizacional Desarrollo y Soluciones',
    '7. Estructura Organizacional Planeación de Proyectos',
    '8. Estructura Organizacional Soluciones Cadena de Suministro',
    '9. Estructura Organizacional Soluciones Customer y Digital',
    '10. Estructura Organizacional Soluciones Backoffice',
    '11. Estructura Organizacional Soluciones Tienda y Comercial',
    '12. Anexos Descripciones de Puesto',
]

TI_PROPOSITO = (
    'Habilitar el crecimiento y la eficiencia del negocio a través de soluciones '
    'tecnológicas robustas, escalables y seguras, que soporten las operaciones de '
    'tiendas, cadena de suministro y plataformas digitales de OXXO México, '
    'garantizando la continuidad operativa, la gobernanza de datos y la innovación '
    'tecnológica continua.'
)

# ── XML helpers (identical to rebuild_final.py) ───────────────

def clear_slide(slide):
    spTree = slide.shapes._spTree
    keep = {f'{{{NS_P}}}nvGrpSpPr', f'{{{NS_P}}}grpSpPr'}
    for child in list(spTree):
        if child.tag not in keep:
            spTree.remove(child)


def set_fill_xml(spPr, hex_color=None, no_fill=False):
    for tag in [f'{{{NS_A}}}noFill', f'{{{NS_A}}}solidFill']:
        for el in spPr.findall(tag):
            spPr.remove(el)
    if no_fill:
        etree.SubElement(spPr, f'{{{NS_A}}}noFill')
    elif hex_color:
        sf = etree.SubElement(spPr, f'{{{NS_A}}}solidFill')
        etree.SubElement(sf, f'{{{NS_A}}}srgbClr').set('val', hex_color)


def set_border_xml(spPr, hex_color='000000', width_pt=0.75):
    ln = spPr.find(f'{{{NS_A}}}ln')
    if ln is None:
        ln = etree.SubElement(spPr, f'{{{NS_A}}}ln')
    ln.set('w', str(int(width_pt * 12700)))
    for child in list(ln):
        ln.remove(child)
    sf = etree.SubElement(ln, f'{{{NS_A}}}solidFill')
    etree.SubElement(sf, f'{{{NS_A}}}srgbClr').set('val', hex_color)


def set_no_border_xml(spPr):
    ln = spPr.find(f'{{{NS_A}}}ln')
    if ln is None:
        ln = etree.SubElement(spPr, f'{{{NS_A}}}ln')
    for child in list(ln):
        ln.remove(child)
    etree.SubElement(ln, f'{{{NS_A}}}noFill')


def make_sp(spTree, shape_id, name, left, top, width, height,
            fill_hex=None, no_fill=False, border_hex=None, no_border=False,
            border_pt=0.75, wrap='square', lIns=45720, rIns=45720, tIns=30000, bIns=30000):
    sp = etree.SubElement(spTree, f'{{{NS_P}}}sp')
    nvSpPr = etree.SubElement(sp, f'{{{NS_P}}}nvSpPr')
    cNvPr = etree.SubElement(nvSpPr, f'{{{NS_P}}}cNvPr')
    cNvPr.set('id', str(shape_id)); cNvPr.set('name', name)
    etree.SubElement(nvSpPr, f'{{{NS_P}}}cNvSpPr').set('txBox', '1')
    etree.SubElement(nvSpPr, f'{{{NS_P}}}nvPr')
    spPr = etree.SubElement(sp, f'{{{NS_P}}}spPr')
    xfrm = etree.SubElement(spPr, f'{{{NS_A}}}xfrm')
    off = etree.SubElement(xfrm, f'{{{NS_A}}}off')
    off.set('x', str(int(left))); off.set('y', str(int(top)))
    ext = etree.SubElement(xfrm, f'{{{NS_A}}}ext')
    ext.set('cx', str(int(width))); ext.set('cy', str(int(height)))
    pg = etree.SubElement(spPr, f'{{{NS_A}}}prstGeom'); pg.set('prst', 'rect')
    etree.SubElement(pg, f'{{{NS_A}}}avLst')
    set_fill_xml(spPr, hex_color=fill_hex, no_fill=no_fill)
    if no_border:
        set_no_border_xml(spPr)
    elif border_hex:
        set_border_xml(spPr, hex_color=border_hex, width_pt=border_pt)
    txBody = etree.SubElement(sp, f'{{{NS_P}}}txBody')
    bodyPr = etree.SubElement(txBody, f'{{{NS_A}}}bodyPr')
    bodyPr.set('wrap', wrap)
    bodyPr.set('lIns', str(lIns)); bodyPr.set('rIns', str(rIns))
    bodyPr.set('tIns', str(tIns)); bodyPr.set('bIns', str(bIns))
    etree.SubElement(txBody, f'{{{NS_A}}}lstStyle')
    return txBody


def add_para(txBody, text, size_pt, bold=False, italic=False, color='000000',
             align='l', spc_before_pt=0):
    para = etree.SubElement(txBody, f'{{{NS_A}}}p')
    pPr = etree.SubElement(para, f'{{{NS_A}}}pPr')
    pPr.set('algn', align)
    if spc_before_pt:
        spcBef = etree.SubElement(pPr, f'{{{NS_A}}}spcBef')
        etree.SubElement(spcBef, f'{{{NS_A}}}spcPts').set('val', str(int(spc_before_pt * 100)))
    r = etree.SubElement(para, f'{{{NS_A}}}r')
    rPr = etree.SubElement(r, f'{{{NS_A}}}rPr')
    rPr.set('sz', str(int(size_pt * 100)))
    rPr.set('b', '1' if bold else '0')
    rPr.set('i', '1' if italic else '0')
    sf = etree.SubElement(rPr, f'{{{NS_A}}}solidFill')
    etree.SubElement(sf, f'{{{NS_A}}}srgbClr').set('val', color)
    etree.SubElement(r, f'{{{NS_A}}}t').text = text
    return para


def add_title_ph(spTree, shape_id, text, font_pt=22):
    sp = etree.SubElement(spTree, f'{{{NS_P}}}sp')
    nvSpPr = etree.SubElement(sp, f'{{{NS_P}}}nvSpPr')
    cNvPr = etree.SubElement(nvSpPr, f'{{{NS_P}}}cNvPr')
    cNvPr.set('id', str(shape_id)); cNvPr.set('name', 'Title 1')
    cNvSpPr = etree.SubElement(nvSpPr, f'{{{NS_P}}}cNvSpPr')
    etree.SubElement(cNvSpPr, f'{{{NS_A}}}spLocks').set('noGrp', '1')
    nvPr = etree.SubElement(nvSpPr, f'{{{NS_P}}}nvPr')
    etree.SubElement(nvPr, f'{{{NS_P}}}ph').set('type', 'title')
    etree.SubElement(sp, f'{{{NS_P}}}spPr')
    txBody = etree.SubElement(sp, f'{{{NS_P}}}txBody')
    etree.SubElement(txBody, f'{{{NS_A}}}bodyPr')
    etree.SubElement(txBody, f'{{{NS_A}}}lstStyle')
    para = etree.SubElement(txBody, f'{{{NS_A}}}p')
    r = etree.SubElement(para, f'{{{NS_A}}}r')
    rPr = etree.SubElement(r, f'{{{NS_A}}}rPr')
    rPr.set('b', '1'); rPr.set('sz', str(int(font_pt * 100)))
    sf = etree.SubElement(rPr, f'{{{NS_A}}}solidFill')
    etree.SubElement(sf, f'{{{NS_A}}}srgbClr').set('val', '000000')
    etree.SubElement(r, f'{{{NS_A}}}t').text = text


def add_line(spTree, shape_id, x, y, length, vertical=False):
    if vertical:
        make_sp(spTree, shape_id, f'Line{shape_id}', x-6350, y, 12700, length,
                fill_hex='000000', no_border=True, wrap='none')
    else:
        make_sp(spTree, shape_id, f'Line{shape_id}', x, y-6350, length, 12700,
                fill_hex='000000', no_border=True, wrap='none')


def hide_slide(slide):
    slide._element.set('show', '0')


# ── Org chart slide builder ───────────────────────────────────

def build_org_chart_slide(slide, l2_info, start_id=200):
    spTree = slide.shapes._spTree
    clear_slide(slide)
    sid = start_id
    N = len(l2_info['l3s'])

    add_title_ph(spTree, sid, 'Estructura Organizacional', font_pt=22); sid += 1

    txb = make_sp(spTree, sid, 'Subtitle', 330200, SUBTITLE_TOP, 10000000, SUBTITLE_H,
                  no_fill=True, no_border=True, tIns=0, bIns=0); sid += 1
    add_para(txb, l2_info['subtitle'], 18, italic=True, color='CC0000', align='l')

    L2_L = (SLIDE_W - L2_W) // 2
    txb = make_sp(spTree, sid, 'L2Box', L2_L, L2_TOP, L2_W, L2_H,
                  fill_hex='FFFFFF', border_hex='000000', border_pt=1,
                  tIns=35000, bIns=35000); sid += 1
    first = True
    for line in l2_info['l2_title'].split('\n'):
        add_para(txb, line, 10, bold=True, color='000000', align='ctr',
                 spc_before_pt=0 if first else 2)
        first = False

    L2_CX = L2_L + L2_W // 2
    VERT_LEN = HBAR_Y - (L2_TOP + L2_H)
    add_line(spTree, sid, L2_CX, L2_TOP + L2_H, VERT_LEN, vertical=True); sid += 1

    GAP = 50000 if N <= 4 else 35000
    L3_W = (SLIDE_W - 2 * SIDE_M - (N - 1) * GAP) // N
    L3_LEFTS = [SIDE_M + i * (L3_W + GAP) for i in range(N)]
    L3_CXS = [L + L3_W // 2 for L in L3_LEFTS]

    if N > 1:
        add_line(spTree, sid, L3_CXS[0], HBAR_Y, L3_CXS[-1] - L3_CXS[0], vertical=False); sid += 1
    for cx in L3_CXS:
        add_line(spTree, sid, cx, HBAR_Y, L3_TOP - HBAR_Y, vertical=True); sid += 1

    title_pt = 10 if N <= 4 else 9
    bullet_pt = 9 if N <= 4 else 8

    for i, l3 in enumerate(l2_info['l3s']):
        L3_L = L3_LEFTS[i]
        txb = make_sp(spTree, sid, f'L3_{i}', L3_L, L3_TOP, L3_W, L3_H,
                      fill_hex='FFFFFF', border_hex='1F3864', border_pt=1,
                      tIns=25000, bIns=25000); sid += 1
        first = True
        for line in l3['title'].split('\n'):
            add_para(txb, line, title_pt, bold=True, color='1F3864', align='ctr',
                     spc_before_pt=0 if first else 1)
            first = False
        txb = make_sp(spTree, sid, f'Desc_{i}', L3_L, DESC_TOP, L3_W, DESC_H_STD,
                      fill_hex='F2F6FC', border_hex='BDD0E9', border_pt=0.75); sid += 1
        first = True
        for bullet in l3['bullets']:
            add_para(txb, '• ' + bullet, bullet_pt, color='1F3864',
                     align='l', spc_before_pt=0 if first else 5)
            first = False

    print(f'  Built org chart: {l2_info["subtitle"]} ({N} L3s)')


# ── L1→L2 overview slide builder ─────────────────────────────

def build_l1_l2_overview(slide, overview, dept_name):
    spTree = slide.shapes._spTree
    clear_slide(slide)
    sid = 100
    N = len(overview['l2s'])

    add_title_ph(spTree, sid, 'Estructura Organizacional', font_pt=22); sid += 1

    txb = make_sp(spTree, sid, 'Subtitle', 330200, SUBTITLE_TOP, 10000000, SUBTITLE_H,
                  no_fill=True, no_border=True, tIns=0, bIns=0); sid += 1
    add_para(txb, dept_name, 18, italic=True, color='CC0000', align='l')

    L1_W = 3400000; L1_H = 500000
    L1_L = (SLIDE_W - L1_W) // 2
    txb = make_sp(spTree, sid, 'L1Box', L1_L, L2_TOP, L1_W, L1_H,
                  fill_hex='1F3864', border_hex='1F3864', border_pt=1,
                  tIns=35000, bIns=35000); sid += 1
    first = True
    for line in overview['l1_title'].split('\n'):
        add_para(txb, line, 11, bold=True, color='FFFFFF', align='ctr',
                 spc_before_pt=0 if first else 2)
        first = False

    L1_CX = L1_L + L1_W // 2
    HBAR_Y_ = L3_TOP - 100000
    add_line(spTree, sid, L1_CX, L2_TOP + L1_H, HBAR_Y_ - (L2_TOP + L1_H), vertical=True); sid += 1

    GAP = 35000
    L2_W_ = (SLIDE_W - 2 * SIDE_M - (N - 1) * GAP) // N
    L2_LEFTS = [SIDE_M + i * (L2_W_ + GAP) for i in range(N)]
    L2_CXS = [L + L2_W_ // 2 for L in L2_LEFTS]

    if N > 1:
        add_line(spTree, sid, L2_CXS[0], HBAR_Y_, L2_CXS[-1] - L2_CXS[0], vertical=False); sid += 1
    for cx in L2_CXS:
        add_line(spTree, sid, cx, HBAR_Y_, L3_TOP - HBAR_Y_, vertical=True); sid += 1

    for i, (l2_title, l2_desc) in enumerate(overview['l2s']):
        L2_L_ = L2_LEFTS[i]
        txb = make_sp(spTree, sid, f'L2_{i}', L2_L_, L3_TOP, L2_W_, L3_H,
                      fill_hex='1F3864', border_hex='1F3864', border_pt=1,
                      tIns=20000, bIns=20000); sid += 1
        first = True
        for line in l2_title.split('\n'):
            add_para(txb, line, 9, bold=True, color='FFFFFF', align='ctr',
                     spc_before_pt=0 if first else 1)
            first = False
        txb = make_sp(spTree, sid, f'L2Desc_{i}', L2_L_, DESC_TOP, L2_W_, DESC_H_STD,
                      fill_hex='EEF3FB', border_hex='BDD0E9', border_pt=0.75); sid += 1
        add_para(txb, l2_desc, 8, color='1F3864', align='l')

    print(f'  Built L1→L2 overview ({dept_name})')


# ── RASCI detail slide builder ────────────────────────────────

def add_rasci_table(spTree, shape_id, role_shorts, resp):
    """Build 2-row RASCI table matching Finanzas format."""
    N = len(role_shorts)
    total_w = 11460480
    col_w = total_w // (N + 1)

    gf = etree.SubElement(spTree, f'{{{NS_P}}}graphicFrame')
    nvGFPr = etree.SubElement(gf, f'{{{NS_P}}}nvGraphicFramePr')
    cNvPr = etree.SubElement(nvGFPr, f'{{{NS_P}}}cNvPr')
    cNvPr.set('id', str(shape_id)); cNvPr.set('name', f'Table {shape_id}')
    cNvGFPr = etree.SubElement(nvGFPr, f'{{{NS_P}}}cNvGraphicFramePr')
    etree.SubElement(cNvGFPr, f'{{{NS_A}}}graphicFrameLocks').set('noGrp', '1')
    etree.SubElement(nvGFPr, f'{{{NS_P}}}nvPr')
    xfrm = etree.SubElement(gf, f'{{{NS_P}}}xfrm')
    off = etree.SubElement(xfrm, f'{{{NS_A}}}off'); off.set('x', '365760'); off.set('y', '1508760')
    ext = etree.SubElement(xfrm, f'{{{NS_A}}}ext'); ext.set('cx', str(total_w)); ext.set('cy', '457200')
    graphic = etree.SubElement(gf, f'{{{NS_A}}}graphic')
    graphicData = etree.SubElement(graphic, f'{{{NS_A}}}graphicData')
    graphicData.set('uri', 'http://schemas.openxmlformats.org/drawingml/2006/table')
    tbl = etree.SubElement(graphicData, f'{{{NS_A}}}tbl')
    tblPr = etree.SubElement(tbl, f'{{{NS_A}}}tblPr')
    tblPr.set('firstRow', '1'); tblPr.set('bandRow', '1')
    etree.SubElement(tblPr, f'{{{NS_A}}}tableStyleId').text = '{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}'
    tblGrid = etree.SubElement(tbl, f'{{{NS_A}}}tblGrid')
    for _ in range(N + 1):
        etree.SubElement(tblGrid, f'{{{NS_A}}}gridCol').set('w', str(col_w))

    def make_tc(row_el, text, bg_hex, text_hex='FFFFFF', bold=False):
        tc = etree.SubElement(row_el, f'{{{NS_A}}}tc')
        txBody = etree.SubElement(tc, f'{{{NS_A}}}txBody')
        etree.SubElement(txBody, f'{{{NS_A}}}bodyPr')
        etree.SubElement(txBody, f'{{{NS_A}}}lstStyle')
        p = etree.SubElement(txBody, f'{{{NS_A}}}p')
        pPr = etree.SubElement(p, f'{{{NS_A}}}pPr'); pPr.set('algn', 'ctr')
        r = etree.SubElement(p, f'{{{NS_A}}}r')
        rPr = etree.SubElement(r, f'{{{NS_A}}}rPr')
        rPr.set('sz', '700'); rPr.set('b', '1' if bold else '0')
        sf = etree.SubElement(rPr, f'{{{NS_A}}}solidFill')
        etree.SubElement(sf, f'{{{NS_A}}}srgbClr').set('val', text_hex)
        etree.SubElement(r, f'{{{NS_A}}}t').text = text
        tcPr = etree.SubElement(tc, f'{{{NS_A}}}tcPr')
        sf2 = etree.SubElement(tcPr, f'{{{NS_A}}}solidFill')
        etree.SubElement(sf2, f'{{{NS_A}}}srgbClr').set('val', bg_hex)

    tr1 = etree.SubElement(tbl, f'{{{NS_A}}}tr'); tr1.set('h', '228600')
    make_tc(tr1, 'Responsabilidad', '1F3964', 'FFFFFF', bold=True)
    for short in role_shorts:
        make_tc(tr1, short, '1F3964', 'FFFFFF', bold=True)

    tr2 = etree.SubElement(tbl, f'{{{NS_A}}}tr'); tr2.set('h', '228600')
    abbr = resp[1][:38] + '...' if len(resp[1]) > 38 else resp[1]
    make_tc(tr2, abbr, 'FFFFFF', '000000', bold=False)
    for i in range(N):
        letter = resp[3].get(i, 'I')
        make_tc(tr2, letter, rasci_color(letter), 'FFFFFF', bold=False)


def build_rasci_slide(slide, role_shorts, role_fulls, resp, dept_name, start_id=50):
    """resp = (section, name, desc, roles_dict, explanation)"""
    spTree = slide.shapes._spTree
    clear_slide(slide)
    sid = start_id
    N = len(role_shorts)
    BADGE_W = 640080
    DESC_X = 1097280
    DESC_W = 10728960
    ROW_H = 582250
    START_Y = 2100000

    add_title_ph(spTree, sid, f'Matriz RASCI - Detalle | {dept_name}', font_pt=22); sid += 1

    txb = make_sp(spTree, sid, 'RespName', 365760, 777240, 11460480, 640080,
                  no_fill=True, no_border=True, lIns=0, rIns=0, tIns=0, bIns=0); sid += 1
    add_para(txb, resp[1], 14, bold=True, color='1F3964')
    add_para(txb, resp[2], 10, color='404040', spc_before_pt=4)

    add_rasci_table(spTree, sid, role_shorts, resp); sid += 1

    current_y = START_Y
    for i, full in enumerate(role_fulls):
        if current_y + ROW_H > SLIDE_H:
            break
        letter = resp[3].get(i, 'I')
        bg = rasci_color(letter)
        txb = make_sp(spTree, sid, f'RBadge_{i}', MARGIN, current_y, BADGE_W, ROW_H,
                      fill_hex=bg, no_border=True,
                      lIns=10000, rIns=10000, tIns=20000, bIns=20000); sid += 1
        add_para(txb, letter, 10, bold=True, color='FFFFFF', align='ctr')
        desc_text = ROLE_DESCS.get(letter, ROLE_DESCS['I'])
        txb = make_sp(spTree, sid, f'RDesc_{i}', DESC_X, current_y, DESC_W, ROW_H,
                      no_fill=True, border_hex=bg, border_pt=3,
                      lIns=55000, rIns=55000, tIns=20000, bIns=20000); sid += 1
        add_para(txb, f'{full}: {desc_text}', 9, color='000000')
        current_y += ROW_H


# ── Intro slide text helpers ──────────────────────────────────

def replace_text_runs(slide, old, new):
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)


def update_agenda_slide(slide, new_items):
    best = None; best_n = 0
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            txt = shape.text_frame.text
            n = len([p for p in shape.text_frame.paragraphs if p.text.strip()])
            if n > best_n and ('1.' in txt or '1 ' in txt):
                best_n = n; best = shape
    if not best:
        return
    txBody = best.text_frame._txBody
    # preserve font size from first rPr
    size_val = None
    for rPr in txBody.iter(f'{{{NS_A}}}rPr'):
        size_val = rPr.get('sz')
        break
    for p in list(txBody.findall(f'{{{NS_A}}}p')):
        txBody.remove(p)
    for item in new_items:
        p = etree.SubElement(txBody, f'{{{NS_A}}}p')
        r = etree.SubElement(p, f'{{{NS_A}}}r')
        rPr = etree.SubElement(r, f'{{{NS_A}}}rPr')
        if size_val:
            rPr.set('sz', size_val)
        etree.SubElement(r, f'{{{NS_A}}}t').text = item


def update_proposito_slide(slide, dept_name, proposito_text):
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and 'Propósito' in shape.text_frame.text:
            txBody = shape.text_frame._txBody
            size_val = None
            for rPr in txBody.iter(f'{{{NS_A}}}rPr'):
                size_val = rPr.get('sz')
                break
            for p in list(txBody.findall(f'{{{NS_A}}}p')):
                txBody.remove(p)
            # Title line
            p = etree.SubElement(txBody, f'{{{NS_A}}}p')
            r = etree.SubElement(p, f'{{{NS_A}}}r')
            rPr = etree.SubElement(r, f'{{{NS_A}}}rPr')
            rPr.set('b', '1')
            if size_val: rPr.set('sz', size_val)
            etree.SubElement(r, f'{{{NS_A}}}t').text = f'Propósito {dept_name}'
            etree.SubElement(txBody, f'{{{NS_A}}}p')  # blank line
            # Body
            p = etree.SubElement(txBody, f'{{{NS_A}}}p')
            r = etree.SubElement(p, f'{{{NS_A}}}r')
            rPr = etree.SubElement(r, f'{{{NS_A}}}rPr')
            rPr.set('b', '0')
            if size_val: rPr.set('sz', size_val)
            etree.SubElement(r, f'{{{NS_A}}}t').text = proposito_text
            break


# ── Anexos slide builder ──────────────────────────────────────

def build_anexos_slide(slide, dept_name, l2_l3_data, start_id=50):
    spTree = slide.shapes._spTree
    clear_slide(slide)
    sid = start_id

    add_title_ph(spTree, sid, 'Anexos\nDescripciones de Puesto', font_pt=22); sid += 1

    txb = make_sp(spTree, sid, 'Anexos', MARGIN, 750000, CONTENT_W, 5800000,
                  no_fill=True, no_border=True, lIns=0, rIns=0, tIns=0, bIns=0); sid += 1

    add_para(txb, f'• Director {dept_name} OXXO Mx:', 9, bold=True, color='1F3864')
    for l2_name, l3_names in l2_l3_data:
        add_para(txb, f'  ○ {l2_name} ({dept_name} L2)', 9, color='1F3864', spc_before_pt=2)
        for l3 in l3_names:
            add_para(txb, f'      ▪ {l3} ({dept_name} L3)', 8, color='262626', spc_before_pt=1)


# ── Main ──────────────────────────────────────────────────────

def main():
    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    shutil.copy2(SRC, OUT)
    print(f'Copied to: {OUT}')

    prs = Presentation(OUT)
    total = len(prs.slides)
    print(f'Slides: {total}')

    # 1. Update intro slides (0-4)
    print('Updating intro slides...')
    for i in range(5):
        replace_text_runs(prs.slides[i], 'Finanzas', TI_DEPT)
        replace_text_runs(prs.slides[i], 'finanzas', TI_DEPT.lower())
    update_agenda_slide(prs.slides[2], TI_AGENDA)
    update_proposito_slide(prs.slides[3], TI_DEPT, TI_PROPOSITO)

    # 2. Rebuild slide 5 (index 5) as L1→L2 TI overview
    print('Building L1→L2 overview...')
    build_l1_l2_overview(prs.slides[5], TI_L1_L2_OVERVIEW, TI_DEPT)

    # 3. Update RASCI methodology slide (index 6)
    replace_text_runs(prs.slides[6], 'Finanzas', TI_DEPT)

    # 4. Build 43 RASCI detail slides (indices 7-49)
    print('Building RASCI slides...')
    for i, resp in enumerate(TI_RASCI):
        idx = 7 + i
        print(f'  RASCI {i+1}/43: {resp[1][:50]}')
        build_rasci_slide(prs.slides[idx], TI_ROLE_SHORT, TI_ROLE_FULL,
                          resp, TI_DEPT, start_id=50 + i * 30)

    # 5. Clear + hide extra RASCI slides (indices 50-52, were Finanzas extras)
    print('Hiding extra RASCI slots...')
    for i in range(50, 53):
        clear_slide(prs.slides[i])
        hide_slide(prs.slides[i])

    # 6. Build 6 org chart slides (indices 53-58)
    print('Building org chart slides...')
    for i, l2_info in enumerate(TI_L2_DATA):
        build_org_chart_slide(prs.slides[53 + i], l2_info, start_id=200 + i * 50)

    # 7. Clear + hide extra org chart slot (index 59)
    clear_slide(prs.slides[59])
    hide_slide(prs.slides[59])

    # 8. Rebuild Anexos (index 66)
    print('Building Anexos...')
    ti_l2_l3 = [(l2['subtitle'], [l3['title'].replace('\n', ' ') for l3 in l2['l3s']])
                for l2 in TI_L2_DATA]
    build_anexos_slide(prs.slides[66], TI_DEPT, ti_l2_l3)

    # 9. Save
    prs.save(OUT)
    print(f'Saved: {OUT}')

    # 10. Verify ZIP
    with zipfile.ZipFile(OUT, 'r') as z:
        names = z.namelist()
        dupes = [n for n in set(names) if names.count(n) > 1]
        ct = z.read('[Content_Types].xml').decode()
        ct_slides = set(re.findall(r'PartName="/ppt/slides/(slide\d+\.xml)"', ct))
        zip_slides = set(n.replace('ppt/slides/', '') for n in names
                         if n.startswith('ppt/slides/slide') and '.rels' not in n)
        print(f'ZIP dupes: {dupes}')
        print(f'Orphans: {zip_slides - ct_slides}')

    prs2 = Presentation(OUT)
    print(f'Final slides: {len(prs2.slides)}')


if __name__ == '__main__':
    main()
